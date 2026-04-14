"""
pptx_exporter.py  —  Future Partners template-based PowerPoint exporter
Uses Template_ReportSlides.pptx layouts and populates placeholders.
Branding sourced entirely from brand_config.py.
"""

import json
import os
import zipfile
from copy import deepcopy
from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData, CategoryChartData
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from brand_config import (
    BRAND,
    SLIDE_BG_COLOR,
    FP_BLACK, FP_CREAM, FP_BLUE,
    TEXT_COLOR_PRIMARY, TEXT_COLOR_SECONDARY, GRIDLINE_COLOR,
    FONT_HEADLINE, FONT_BODY, FONT_NARRATIVE,
    FONT_SIZE,
    CHART_DEFAULTS,
    EXEC_SUMMARY, AI_INSIGHT,
    get_chart_colors, get_palette_for_table_index,
    get_data_label_format,
    SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN,
    TEMPLATE_PATH, LAYOUT, PH, CHART_TEMPLATES_DIR,
)


# ---------------------------------------------------------------------------
# Chart type mapping
# ---------------------------------------------------------------------------

CHART_TYPE_MAP = {
    "bar_h":            {"xl": XL_CHART_TYPE.BAR_CLUSTERED,    "multi": False, "layout": "primary_chart"},
    "bar horizontal":   {"xl": XL_CHART_TYPE.BAR_CLUSTERED,    "multi": False, "layout": "primary_chart"},
    "bar_v":            {"xl": XL_CHART_TYPE.COLUMN_CLUSTERED, "multi": False, "layout": "primary_chart"},
    "column":           {"xl": XL_CHART_TYPE.COLUMN_CLUSTERED, "multi": False, "layout": "primary_chart"},
    "donut":            {"xl": XL_CHART_TYPE.DOUGHNUT,         "multi": False, "layout": "primary_chart"},
    "line":             {"xl": XL_CHART_TYPE.LINE_MARKERS,     "multi": False, "layout": "primary_chart"},
    "grouped_bar_2":    {"xl": XL_CHART_TYPE.COLUMN_CLUSTERED, "multi": True,  "layout": "primary_chart"},
    "grouped_bar_3":    {"xl": XL_CHART_TYPE.COLUMN_CLUSTERED, "multi": True,  "layout": "primary_chart"},
    "multi_line":       {"xl": XL_CHART_TYPE.LINE_MARKERS,     "multi": True,  "layout": "primary_chart"},
    "table_only":       {"xl": None,                           "multi": False, "layout": "primary_chart"},
    "chart_table":      {"xl": XL_CHART_TYPE.BAR_CLUSTERED,    "multi": False, "layout": "one_two_third_alt"},
}

# User-facing label → internal key
CHART_LABEL_MAP = {
    "Horizontal Bar":       "bar_h",
    "Vertical Bar":         "bar_v",
    "Donut":                "donut",
    "Line":                 "line",
    "Grouped Bar (2)":      "grouped_bar_2",
    "Grouped Bar (3)":      "grouped_bar_3",
    "Multi-Line":           "multi_line",
    "Table Only":           "table_only",
    "Chart + Table":        "chart_table",
}

CHART_LABELS = list(CHART_LABEL_MAP.keys())


CRTX_MAP = {
    "bar_h":         "FP Blue - Horizontal Bar.crtx",
    "bar_v":         "FP Blue - Vertical Bar.crtx",
    "donut":         "FP Blue - Donut.crtx",
    "line":          "FP Blue - Line.crtx",
    "multi_line":    "FP Blue - Multi-Line.crtx",
    "grouped_bar_2": "FP Blue - Grouped Bar (2).crtx",
    "grouped_bar_3": "FP Blue - Grouped Bar (3).crtx",
}


def _resolve_chart_type(kind: str):
    kind_lower = (kind or "bar_h").lower().strip()
    if kind_lower in CHART_TYPE_MAP:
        return CHART_TYPE_MAP[kind_lower]
    for label, key in CHART_LABEL_MAP.items():
        if kind_lower == label.lower():
            return CHART_TYPE_MAP[key]
    return CHART_TYPE_MAP["bar_h"]


_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_C    = "{%s}" % _C_NS
_A    = "{%s}" % _A_NS

PLOT_TAG_MAP = {
    "bar_h":         f"{_C}barChart",
    "bar_v":         f"{_C}barChart",
    "column":        f"{_C}barChart",
    "donut":         f"{_C}doughnutChart",
    "line":          f"{_C}lineChart",
    "multi_line":    f"{_C}lineChart",
    "grouped_bar_2": f"{_C}barChart",
    "grouped_bar_3": f"{_C}barChart",
}


def _apply_crtx_template(chart_part, chart_kind: str):
    """Replace the chart's formatting with the .crtx template while keeping our data.

    Strategy: take the ENTIRE template chart XML and graft our data references
    (c:tx, c:cat, c:val, c:extLst with uniqueIds) into the template's series
    structure. This preserves 100% of the template's formatting.
    """
    kind_key = chart_kind.lower().strip()
    crtx_filename = CRTX_MAP.get(kind_key)
    if not crtx_filename:
        return False

    crtx_path = os.path.join(CHART_TEMPLATES_DIR, crtx_filename)
    if not os.path.exists(crtx_path):
        return False

    try:
        with zipfile.ZipFile(crtx_path, "r") as zf:
            tmpl_xml = zf.read("chart/chart.xml")
            style_xml = zf.read("chart/charts/style1.xml")
            colors_xml = zf.read("chart/charts/colors1.xml")

        # --- Inject style/colors OPC parts ---
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI
        import re as _re

        chart_partname = str(chart_part.partname)
        base = chart_partname.rsplit("/", 1)[0]
        chart_num = _re.search(r"(\d+)", chart_partname.rsplit("/", 1)[1])
        suffix = chart_num.group(1) if chart_num else "1"

        style_rt = "http://schemas.microsoft.com/office/2011/relationships/chartStyle"
        colors_rt = "http://schemas.microsoft.com/office/2011/relationships/chartColorStyle"

        style_part = Part(PackURI(f"{base}/style{suffix}.xml"),
                          "application/vnd.ms-office.chartstyle+xml",
                          package=chart_part.package, blob=style_xml)
        chart_part.relate_to(style_part, style_rt)

        colors_part = Part(PackURI(f"{base}/colors{suffix}.xml"),
                           "application/vnd.ms-office.chartcolorstyle+xml",
                           package=chart_part.package, blob=colors_xml)
        chart_part.relate_to(colors_part, colors_rt)

        # --- Full chart XML replacement ---
        tmpl_tree = etree.fromstring(tmpl_xml)
        chart_tree = chart_part._element

        plot_tag = PLOT_TAG_MAP.get(kind_key, f"{_C}barChart")

        # Get the template's <c:chart> and our <c:chart>
        tmpl_chart = tmpl_tree.find(f"{_C}chart")
        our_chart  = chart_tree.find(f"{_C}chart")
        if tmpl_chart is None or our_chart is None:
            return True

        # Extract data references from our series before we replace them
        our_plot = our_chart.find(f".//{plot_tag}")
        if our_plot is None:
            return True
        our_series_list = our_plot.findall(f"{_C}ser")

        data_refs = []
        for ser in our_series_list:
            refs = {
                "tx":  deepcopy(ser.find(f"{_C}tx")),
                "cat": deepcopy(ser.find(f"{_C}cat")),
                "val": deepcopy(ser.find(f"{_C}val")),
                "idx": ser.find(f"{_C}idx").get("val") if ser.find(f"{_C}idx") is not None else "0",
                "order": ser.find(f"{_C}order").get("val") if ser.find(f"{_C}order") is not None else "0",
            }
            data_refs.append(refs)

        # Replace our <c:chart> with the template's <c:chart>
        tmpl_chart_copy = deepcopy(tmpl_chart)
        parent = our_chart.getparent()
        idx = list(parent).index(our_chart)
        parent.remove(our_chart)
        parent.insert(idx, tmpl_chart_copy)

        # Now graft our data into the template's series
        new_plot = tmpl_chart_copy.find(f".//{plot_tag}")
        if new_plot is None:
            return True

        tmpl_series = new_plot.findall(f"{_C}ser")
        tmpl_ser_template = deepcopy(tmpl_series[0]) if tmpl_series else None

        # Remove template's placeholder series
        for ts in tmpl_series:
            new_plot.remove(ts)

        # Build new series from template formatting + our data
        if tmpl_ser_template is not None:
            for i, refs in enumerate(data_refs):
                new_ser = deepcopy(tmpl_ser_template)

                # Set idx and order
                idx_el = new_ser.find(f"{_C}idx")
                if idx_el is not None:
                    idx_el.set("val", str(i))
                order_el = new_ser.find(f"{_C}order")
                if order_el is not None:
                    order_el.set("val", str(i))

                # Replace data references with ours
                for tag_name in ["tx", "cat", "val"]:
                    old = new_ser.find(f"{_C}{tag_name}")
                    if old is not None:
                        new_ser.remove(old)
                    if refs[tag_name] is not None:
                        new_ser.append(refs[tag_name])

                # For multi-point charts (donut/pie), generate dPt entries
                # matching the number of categories if the template has dPt
                if i == 0:
                    tmpl_dpts = new_ser.findall(f"{_C}dPt")
                    n_cats = 0
                    cat_ref = refs.get("cat")
                    if cat_ref is not None:
                        str_cache = cat_ref.find(f".//{_C}strCache")
                        if str_cache is not None:
                            n_cats = len(str_cache.findall(f"{_C}pt"))
                        num_cache = cat_ref.find(f".//{_C}numCache")
                        if num_cache is not None and n_cats == 0:
                            n_cats = len(num_cache.findall(f"{_C}pt"))

                    if tmpl_dpts and n_cats > len(tmpl_dpts):
                        palette = get_chart_colors(n_cats, "blue")
                        for dp_idx in range(len(tmpl_dpts), n_cats):
                            new_dpt = deepcopy(tmpl_dpts[dp_idx % len(tmpl_dpts)])
                            dp_idx_el = new_dpt.find(f"{_C}idx")
                            if dp_idx_el is not None:
                                dp_idx_el.set("val", str(dp_idx))
                            fill_el = new_dpt.find(f".//{_A}srgbClr")
                            if fill_el is not None and dp_idx < len(palette):
                                fill_el.set("val", str(palette[dp_idx]))
                            new_ser.append(new_dpt)

                new_plot.append(new_ser)

        # Preserve relationship references from our original chart
        # (the embedded workbook relationship is in chartSpace, copy it back)
        for rel_tag in [f"{_C}externalData"]:
            our_rel = chart_tree.find(f".//{rel_tag}")
            if our_rel is None:
                our_rel_orig = our_chart.find(f".//{rel_tag}")

        # Remove autoTitleDeleted=0 if template has it as 1
        auto_title = tmpl_chart_copy.find(f"{_C}autoTitleDeleted")
        if auto_title is not None:
            auto_title.set("val", "1")

        # Remove legend if template doesn't have one
        tmpl_legend = tmpl_chart.find(f"{_C}legend")
        new_legend = tmpl_chart_copy.find(f"{_C}legend")
        if tmpl_legend is None and new_legend is not None:
            tmpl_chart_copy.remove(new_legend)

        # Copy chartSpace-level spPr and txPr from template
        for tag in [f"{_C}spPr", f"{_C}txPr"]:
            tmpl_elem = tmpl_tree.find(tag)
            if tmpl_elem is not None:
                old_elem = chart_tree.find(tag)
                if old_elem is not None:
                    chart_tree.remove(old_elem)
                chart_tree.append(deepcopy(tmpl_elem))

        return True
    except Exception as e:
        print(f"Warning: Could not apply .crtx template for '{chart_kind}': {e}")
        import traceback; traceback.print_exc()
        return False


# ---------------------------------------------------------------------------
# TextCallout dataclass
# ---------------------------------------------------------------------------

@dataclass
class TextCallout:
    """Represents a text callout that can be associated with table data."""

    table_title: str
    column_key:  str
    row_label:   str

    text:         Optional[str]    = None
    position:     tuple            = (0.5, 7.0, 9.0, 0.4)
    font_size:    int              = 12
    font_bold:    bool             = True
    font_color:   Optional[RGBColor] = None
    metric_type:  str              = "percentage"

    bind_question: str  = "TEXT_QUESTION"
    bind_base:     str  = "TEXT_BASE"
    auto_update:   bool = True

    def __post_init__(self):
        if self.text is None:
            self.text = f"{self.row_label}: [Value]"

    def get_display_text(self, table_data: Optional[Dict[str, Any]] = None) -> str:
        if table_data:
            row_idx = self._find_row_index(table_data)
            col_idx = self._find_column_index(table_data)
            if row_idx is not None and col_idx is not None:
                try:
                    value = table_data["values"][row_idx][col_idx]
                    if value is not None:
                        formatted_value = ""
                        if isinstance(value, (int, float)):
                            mt = (self.metric_type or "").lower()
                            if mt == "percentage":
                                formatted_value = f"{float(value) * 100:.1f}%"
                            elif mt == "currency":
                                formatted_value = f"${float(value):,.0f}"
                            else:
                                formatted_value = f"{float(value):,.1f}"
                        else:
                            formatted_value = str(value)
                        if self.text and "[Value]" in self.text:
                            return self.text.replace("[Value]", formatted_value)
                        if self.text:
                            return self.text
                        return f"{self.row_label}: {formatted_value}"
                except (IndexError, TypeError):
                    pass
        return self.text if self.text else f"{self.row_label}: [Value]"

    def _find_row_index(self, table_data):
        for i, label in enumerate(table_data.get("row_labels", [])):
            if isinstance(label, str) and self.row_label.lower() in label.lower():
                return i
        return None

    def _find_column_index(self, table_data):
        col_labels = table_data.get("col_labels", [])
        if self.column_key in col_labels:
            return col_labels.index(self.column_key)
        for fallback in ["Total", "Overall", "All", "Base"]:
            if fallback in col_labels:
                return col_labels.index(fallback)
        return 0 if col_labels else None

    def to_mapping_dict(self):
        return {
            "type": "text_callout",
            "table_title": self.table_title,
            "column": self.column_key,
            "row": self.row_label,
            "metric_type": self.metric_type,
            "auto_update": "yes" if self.auto_update else "no",
        }


# ---------------------------------------------------------------------------
# Alt-text / mapping helpers
# ---------------------------------------------------------------------------

def _set_alt_text(shape, mapping: dict):
    try:
        lines = [f"{k}: {v}" for k, v in mapping.items() if v is not None and v != ""]
        alt_text_content = "\n".join(lines)
        try:
            if hasattr(shape, "alternative_text"):
                shape.alternative_text = alt_text_content
                return
        except Exception:
            pass
        try:
            if hasattr(shape, "element"):
                ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                c_nv_pr = None
                if "graphicFrame" in shape.element.tag:
                    c_nv_pr = shape.element.find(".//p:cNvPr", namespaces=ns)
                elif "sp" in shape.element.tag:
                    c_nv_pr = shape.element.find(".//p:cNvPr", namespaces=ns)
                if c_nv_pr is not None:
                    c_nv_pr.set("descr", alt_text_content)
                    return
        except Exception:
            pass
    except Exception:
        pass


def _tag_shape(shape, obj_type: str, table_title: str, col_key=None,
               bind_question="TEXT_QUESTION", bind_base="TEXT_BASE",
               row_label=None, metric_type=None):
    base = {"table_title": table_title, "column": col_key or "Total", "auto_update": "yes"}
    type_map = {
        "chart":           {"type": "chart", "exclude_rows": "base, mean, average, avg"},
        "table":           {"type": "table", "columns": "*", "exclude_rows": "base, mean, average, avg"},
        "text_question":   {"type": "text_question"},
        "question_text":   {"type": "text_question"},
        "text_base":       {"type": "text_base"},
        "text_title":      {"type": "text_title"},
        "text_callout":    {"type": "text_callout", "row": row_label or "",
                            "metric_type": metric_type or "percentage"},
        "ai_insight":      {"type": "ai_insight"},
        "text_takeaway":   {"type": "text_takeaway"},
        "text_analysis":   {"type": "text_analysis"},
        "text_chart_title":{"type": "text_chart_title"},
    }
    extra = type_map.get(obj_type, {"type": obj_type})
    _set_alt_text(shape, {**extra, **base})


# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------

def _open_template() -> Presentation:
    """Open the branded template and strip all example slides."""
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(
            f"Template not found at {TEMPLATE_PATH}. "
            "Place Template_ReportSlides.pptx in the templates/ folder."
        )
    prs = Presentation(TEMPLATE_PATH)
    _remove_all_slides(prs)
    return prs


def _remove_all_slides(prs: Presentation):
    """Delete every slide from the presentation while keeping layouts."""
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    while len(prs.slides._sldIdLst) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(f"{ns}id")
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


def _remove_placeholder(slide, ph_idx: int):
    """Remove a placeholder from the slide DOM to prevent template defaults from showing."""
    try:
        ph = slide.placeholders[ph_idx]
        sp_elem = ph._element
        sp_elem.getparent().remove(sp_elem)
    except (KeyError, AttributeError):
        pass


def _fill_placeholder_text(slide, ph_idx: int, text: str, tag_type: str = None,
                           table_title: str = None):
    """Safely set text in a placeholder and optionally tag it for mapping."""
    try:
        ph = slide.placeholders[ph_idx]
        ph.text = text or ""
        if tag_type and table_title:
            _tag_shape(ph, tag_type, table_title)
        return ph
    except KeyError:
        return None


def _detect_value_type(values, col_indices):
    """Sample values to determine if they're proportions (0-1) or raw numbers."""
    sample = []
    for row in values:
        for ci in col_indices:
            if ci < len(row) and row[ci] is not None:
                try:
                    sample.append(float(row[ci]))
                except (TypeError, ValueError):
                    pass
    return sample


# ---------------------------------------------------------------------------
# Chart formatting
# ---------------------------------------------------------------------------

def _apply_chart_formatting(chart, palette_name: str, n_series: int,
                            value_sample: list, chart_kind: str = ""):
    """Apply brand-compliant formatting to a chart object."""
    num_fmt = get_data_label_format(value_sample)
    is_pie_type = chart_kind.lower() in ("donut", "doughnut", "pie")

    chart.has_legend = n_series > 1 or is_pie_type
    for s in chart.series:
        s.data_labels.show_value = True
        s.data_labels.number_format = num_fmt
        if not is_pie_type:
            try:
                s.data_labels.position = 2
            except Exception:
                pass
        try:
            s.data_labels.font.name = FONT_BODY
            s.data_labels.font.size = FONT_SIZE["data_label"]
            s.data_labels.font.bold = True
        except Exception:
            pass

    if not is_pie_type:
        try:
            chart.plots[0].gap_width = CHART_DEFAULTS["gap_width"]
            if n_series > 1:
                chart.plots[0].overlap = CHART_DEFAULTS["overlap"]
        except Exception:
            pass

        try:
            gl = chart.value_axis.major_gridlines
            gl.format.line.width = Pt(0.5)
            gl.format.line.fore_color.rgb = GRIDLINE_COLOR
        except Exception:
            pass

        try:
            chart.category_axis.has_title = False
            chart.value_axis.has_title = False
            chart.category_axis.tick_labels.font.size = FONT_SIZE["axis"]
            chart.category_axis.tick_labels.font.name = FONT_BODY
            chart.value_axis.tick_labels.font.size = FONT_SIZE["axis"]
            chart.value_axis.tick_labels.font.name = FONT_BODY
        except Exception:
            pass

    colors = get_chart_colors(n_series, palette_name)
    for i, series in enumerate(chart.series):
        if i < len(colors):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[i]


def _add_data_table(slide, col_labels, row_labels, values, position=None):
    """Add a formatted data table to the slide."""
    rows = 1 + len(row_labels)
    cols = 1 + len(col_labels)
    pos = position or (5.60, 1.85, 6.82, 3.24)
    left, top, width, height = [Inches(p) for p in pos]
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.cell(0, 0).text = ""
    for j, c in enumerate(col_labels, start=1):
        cell = table.cell(0, j)
        cell.text = str(c)

    for i, rlab in enumerate(row_labels, start=1):
        table.cell(i, 0).text = str(rlab)
        for j, v in enumerate(values[i - 1][:len(col_labels)], start=1):
            table.cell(i, j).text = "" if v is None else f"{v:.1f}"

    for r in range(rows):
        for c in range(cols):
            tf = table.cell(r, c).text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_BODY
                    run.font.size = FONT_SIZE["footnote"]
                    if r == 0:
                        run.font.bold = True

    return table_shape


# ---------------------------------------------------------------------------
# Row filtering / sorting
# ---------------------------------------------------------------------------

def _extract_base_text(table: Dict[str, Any], column_key: str = None) -> str:
    """Extract base N from the table data for auto-populating the base text."""
    row_labels = table.get("row_labels", [])
    col_labels = table.get("col_labels", [])
    values = table.get("values", [])
    col_idx = 0
    if column_key and column_key in col_labels:
        col_idx = col_labels.index(column_key)
    elif "Total" in col_labels:
        col_idx = col_labels.index("Total")
    for i, label in enumerate(row_labels):
        if isinstance(label, str) and label.strip().lower().startswith("base"):
            if i < len(values) and col_idx < len(values[i]):
                v = values[i][col_idx]
                if v is not None:
                    try:
                        return f"n = {int(round(float(v))):,}"
                    except (TypeError, ValueError):
                        pass
    return ""


def _get_exclude_indices(row_labels):
    """Return set of row indices that are metadata (base, mean, avg)."""
    exclude = set()
    for i, rlab in enumerate(row_labels):
        if isinstance(rlab, str):
            lab = rlab.strip().lower()
            if lab.startswith(("base", "mean", "average", "avg")):
                exclude.add(i)
    return exclude


def sort_table_rows(table: Dict[str, Any], column_key: str = "Total",
                    excluded_rows: List[str] = None) -> Dict[str, Any]:
    if excluded_rows is None:
        excluded_rows = []
    sorted_table = table.copy()
    row_labels = table["row_labels"].copy()
    values = [row.copy() for row in table["values"]]
    col_labels = table["col_labels"]
    if column_key in col_labels:
        sort_col_idx = col_labels.index(column_key)
    elif "Total" in col_labels:
        sort_col_idx = col_labels.index("Total")
    else:
        sort_col_idx = 0 if col_labels else None
    if sort_col_idx is None:
        return sorted_table
    sortable, excluded = [], []
    for label, row in zip(row_labels, values):
        if label in excluded_rows:
            excluded.append((label, row))
        else:
            sv = row[sort_col_idx] if sort_col_idx < len(row) and row[sort_col_idx] is not None else 0
            sortable.append((label, row, sv))
    sortable.sort(key=lambda x: x[2], reverse=True)
    sorted_table["row_labels"] = [x[0] for x in sortable] + [x[0] for x in excluded]
    sorted_table["values"]     = [x[1] for x in sortable] + [x[1] for x in excluded]
    return sorted_table


# ---------------------------------------------------------------------------
# Main slide builder — template-based, placeholder-driven
# ---------------------------------------------------------------------------

def add_chart_slide(prs, table: Dict[str, Any], chart_kind="bar_h",
                    chart_title=None, base_text=None, question_text=None,
                    callouts=None, enable_sorting=False, excluded_rows=None,
                    column_key=None, column_keys=None,
                    insights: Optional[Dict[str, str]] = None,
                    palette_name: str = "blue"):
    """
    Add a single data slide to the presentation using the template layout.

    Args:
        prs:          Presentation object (from _open_template).
        table:        Table dict from parse_workbook().
        chart_kind:   Chart type key or label.
        chart_title:  Override title for the chart area.
        base_text:    Base description text.
        question_text: Survey question text.
        callouts:     List of TextCallout objects.
        enable_sorting: Sort rows by value.
        excluded_rows: Rows to exclude from sorting.
        column_key:   Single column key for chart data.
        column_keys:  Multiple column keys for grouped/multi-series charts.
        insights:     {"takeaway", "analysis"} from AI.
        palette_name: Which brand palette to use.
    """
    spec = _resolve_chart_type(chart_kind)
    layout_key = spec["layout"]
    layout_idx = LAYOUT.get(layout_key, LAYOUT["primary_chart"])

    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    _remove_placeholder(slide, PH["note"])

    working_table = table
    if enable_sorting:
        sort_col = column_key or "Total"
        working_table = sort_table_rows(table, sort_col, excluded_rows or [])

    insights = insights or {}
    table_title = working_table.get("title", "")
    has_insights = bool(insights.get("takeaway") or insights.get("analysis"))

    # --- Fill text placeholders ---
    # PH 0 = Title / Long Action Takeaway
    takeaway = insights.get("takeaway") or chart_title or table_title
    _fill_placeholder_text(slide, PH["title"], takeaway, "text_takeaway", table_title)

    _remove_placeholder(slide, PH["punch"])

    # PH 1 = Supporting Analysis (left body)
    analysis = insights.get("analysis", "")
    if analysis:
        _fill_placeholder_text(slide, PH["analysis"], analysis, "text_analysis", table_title)
    else:
        _remove_placeholder(slide, PH["analysis"])

    # PH 14 = Chart Title (above chart on right side)
    ct = chart_title or table_title
    _fill_placeholder_text(slide, PH["chart_title"], ct, "text_chart_title", table_title)

    # Question / Base (combined into one text box, separated by newline)
    auto_question = question_text or table_title
    auto_base = base_text or _extract_base_text(working_table, column_key)
    qbase_parts = []
    if auto_question:
        qbase_parts.append(f"Question: {auto_question}")
    if auto_base:
        qbase_parts.append(f"Base: {auto_base}")
    qbase_text = "\n".join(qbase_parts)
    if qbase_text:
        _fill_placeholder_text(slide, PH["qbase"], qbase_text, "text_question", table_title)

    # --- Build chart data ---
    row_labels = working_table["row_labels"]
    col_labels = working_table["col_labels"]
    values     = working_table["values"]
    exclude_indices = _get_exclude_indices(row_labels)

    categories = [lab for i, lab in enumerate(row_labels) if i not in exclude_indices]

    is_multi = spec["multi"] and column_keys and len(column_keys) > 1
    if is_multi:
        series_cols = column_keys
    else:
        ck = column_key or "Total"
        series_cols = [ck] if ck in col_labels else (["Total"] if "Total" in col_labels else col_labels[:1])

    col_indices = []
    for sc in series_cols:
        if sc in col_labels:
            col_indices.append(col_labels.index(sc))
        else:
            col_indices.append(0)

    xl_type = spec["xl"]

    if xl_type is not None:
        chart_data = CategoryChartData()
        chart_data.categories = categories

        value_sample = []
        for ci_idx, (sc, ci) in enumerate(zip(series_cols, col_indices)):
            sv = []
            for i, row in enumerate(values):
                if i in exclude_indices:
                    continue
                v = row[ci] if ci < len(row) else None
                sv.append(v)
                if v is not None:
                    try:
                        value_sample.append(float(v))
                    except (TypeError, ValueError):
                        pass
            chart_data.add_series(sc, sv)

        n_series = len(series_cols)

        # Read chart placeholder bounds, then remove the placeholder to avoid overlap
        try:
            chart_ph = slide.placeholders[PH["chart"]]
            cx, cy, cw, ch = chart_ph.left, chart_ph.top, chart_ph.width, chart_ph.height
            sp_elem = chart_ph._element
            sp_elem.getparent().remove(sp_elem)
        except KeyError:
            cx, cy, cw, ch = Inches(5.51), Inches(1.64), Inches(6.92), Inches(4.43)

        chart_shape = slide.shapes.add_chart(xl_type, cx, cy, cw, ch, chart_data)
        chart = chart_shape.chart
        chart_frame = chart_shape

        _apply_crtx_template(chart_shape.chart_part, chart_kind)
        _apply_chart_formatting(chart, palette_name, n_series, value_sample, chart_kind)
        _tag_shape(chart_frame, "chart", table_title,
                   series_cols[0] if len(series_cols) == 1 else ",".join(series_cols))

    elif chart_kind.lower() in ("table_only", "table only"):
        # Table-only slide: put table in the chart placeholder area
        tbl = _add_data_table(slide, col_labels, row_labels, values,
                              position=(5.51, 1.64, 6.92, 4.43))
        _tag_shape(tbl, "table", table_title)

    # For chart+table layout, add a table as well
    if layout_key == "one_two_third_alt" and xl_type is not None:
        tbl = _add_data_table(slide, col_labels, row_labels, values,
                              position=(5.60, 1.85, 6.82, 3.24))
        _tag_shape(tbl, "table", table_title)

    # Callouts
    if callouts:
        callout_start_y = 6.2
        for i, callout in enumerate(callouts):
            pos = list(callout.position)
            pos[1] = callout_start_y + (i * 0.45)
            adj = TextCallout(
                table_title=callout.table_title, column_key=callout.column_key,
                row_label=callout.row_label, text=callout.text,
                position=tuple(pos), font_size=callout.font_size,
                font_bold=callout.font_bold, font_color=callout.font_color,
                bind_question=callout.bind_question, bind_base=callout.bind_base,
                auto_update=callout.auto_update,
            )
            _add_text_callout(slide, adj, working_table)

    # Metadata box
    meta = json.dumps({
        "table_key": table_title,
        "row_count": len(row_labels),
        "col_labels": col_labels,
        "chart_kind": chart_kind,
        "callout_count": len(callouts) if callouts else 0,
        "enable_sorting": enable_sorting,
        "excluded_rows": excluded_rows or [],
    })
    meta_box = slide.shapes.add_textbox(Inches(0.0), Inches(7.4), Inches(0.1), Inches(0.2))
    meta_box.name = "DATA_META"
    tfm = meta_box.text_frame
    tfm.clear()
    r = tfm.paragraphs[0].add_run()
    r.text = meta
    r.font.size = Pt(1)


def _add_text_callout(slide, callout: TextCallout, table_data=None):
    x, y, w, h = [Inches(pos) for pos in callout.position]
    callout_box = slide.shapes.add_textbox(x, y, w, h)
    _tag_shape(callout_box, "text_callout", callout.table_title, callout.column_key,
               callout.bind_question, callout.bind_base, callout.row_label, callout.metric_type)
    tf = callout_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = callout.get_display_text(table_data)
    run.font.size = Pt(callout.font_size)
    run.font.bold = callout.font_bold
    run.font.name = FONT_BODY
    if callout.font_color:
        run.font.color.rgb = callout.font_color
    return callout_box


# Keep old name available for backward compat
add_text_callout = _add_text_callout


# ---------------------------------------------------------------------------
# Executive Summary slide
# ---------------------------------------------------------------------------

def add_executive_summary_slide(prs, insights: Dict[str, Dict[str, str]]):
    """
    Build an Executive Summary slide with one bullet per table.
    Accepts either new-style three-tier dict or legacy string dict.
    """
    layout_idx = LAYOUT.get("section_header", 4)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title_box = slide.shapes.add_textbox(Inches(0.36), Inches(0.2), Inches(11.5), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = EXEC_SUMMARY["title_text"]
    run.font.size = EXEC_SUMMARY["title_size"]
    run.font.bold = True
    run.font.name = EXEC_SUMMARY["title_font"]
    run.font.color.rgb = EXEC_SUMMARY["title_color"]

    rule = slide.shapes.add_shape(
        1, Inches(0.36), Inches(1.0), Inches(11.5), Emu(36000)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = FP_BLUE
    rule.line.fill.background()

    content_box = slide.shapes.add_textbox(Inches(0.36), Inches(1.15), Inches(11.5), Inches(5.6))
    tf = content_box.text_frame
    tf.word_wrap = True

    items = list(insights.items())[:EXEC_SUMMARY["max_bullets"]]

    for idx, (title, tiers) in enumerate(items):
        if isinstance(tiers, str):
            analysis = tiers
        else:
            analysis = tiers.get("analysis", "")

        p_title = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_title.space_before = Pt(6)
        run_title = p_title.add_run()
        run_title.text = f"▸  {title}"
        run_title.font.name = FONT_BODY
        run_title.font.size = Pt(10)
        run_title.font.bold = True
        run_title.font.color.rgb = EXEC_SUMMARY["accent_color"]

        if analysis:
            p_analysis = tf.add_paragraph()
            p_analysis.space_before = Pt(1)
            run_analysis = p_analysis.add_run()
            run_analysis.text = f"    {analysis}"
            run_analysis.font.name = FONT_BODY
            run_analysis.font.size = Pt(9)
            run_analysis.font.bold = False
            run_analysis.font.color.rgb = TEXT_COLOR_PRIMARY

    if len(insights) > EXEC_SUMMARY["max_bullets"]:
        p_more = tf.add_paragraph()
        p_more.space_before = Pt(4)
        run_more = p_more.add_run()
        run_more.text = f"... and {len(insights) - EXEC_SUMMARY['max_bullets']} additional tables"
        run_more.font.name = FONT_BODY
        run_more.font.size = Pt(9)
        run_more.font.color.rgb = TEXT_COLOR_SECONDARY


# ---------------------------------------------------------------------------
# Main export function
# ---------------------------------------------------------------------------

def export_pptx(tables: List[Dict[str, Any]], selections: Dict[str, Dict[str, Any]],
                out_path: str, ai_insights: Dict[str, Dict[str, str]] = None,
                report_palette: str = "blue"):
    """
    Build a branded PowerPoint from parsed crosstab tables using the template.

    Args:
        tables:       List of table dicts from parse_workbook().
        selections:   Per-table configuration dict (keyed by table ID).
        out_path:     Output .pptx file path.
        ai_insights:  Optional two-tier insights {title: {takeaway, analysis}}.
        report_palette: Color palette for the entire report (default "blue").
    """
    prs = _open_template()

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[LAYOUT["title_slide"]])

    # Executive Summary (if AI insights provided)
    if ai_insights:
        add_executive_summary_slide(prs, ai_insights)

    # One slide per table
    for table_idx, t in enumerate(tables):
        sel = selections.get(t["id"], {})
        ctype = sel.get("chart_type", "bar_h")
        title = sel.get("title") or t["title"]

        # Process callouts
        callouts = None
        if "callouts" in sel:
            callouts = []
            for cd in sel["callouts"]:
                callouts.append(TextCallout(
                    table_title   = t.get("title", ""),
                    column_key    = cd.get("column_key", "Total"),
                    row_label     = cd.get("row_label", ""),
                    text          = cd.get("text"),
                    position      = cd.get("position", (0.5, 7.0, 9.0, 0.4)),
                    font_size     = cd.get("font_size", 12),
                    font_bold     = cd.get("font_bold", True),
                    font_color    = cd.get("font_color"),
                    bind_question = cd.get("bind_question", "TEXT_QUESTION"),
                    bind_base     = cd.get("bind_base", "TEXT_BASE"),
                    metric_type   = cd.get("metric_type", "percentage"),
                    auto_update   = cd.get("auto_update", True),
                ))

        palette_name = report_palette

        table_insights = {}
        if ai_insights:
            table_insights = ai_insights.get(t["title"], {})
            if isinstance(table_insights, str):
                table_insights = {"takeaway": "", "analysis": table_insights}

        add_chart_slide(
            prs,
            t,
            chart_kind     = ctype,
            chart_title    = title,
            base_text      = sel.get("base_text"),
            question_text  = sel.get("question_text"),
            callouts       = callouts,
            enable_sorting = sel.get("enable_sorting", False),
            excluded_rows  = sel.get("excluded_rows", []),
            column_key     = sel.get("column_key"),
            column_keys    = sel.get("column_keys"),
            insights       = table_insights,
            palette_name   = palette_name,
        )

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Callout factory helpers
# ---------------------------------------------------------------------------

def create_row_callout(table_title, row_label, column_key="Total",
                       custom_text=None, position=(0.5, 7.0, 9.0, 0.4)):
    return TextCallout(table_title=table_title, column_key=column_key,
                       row_label=row_label, text=custom_text, position=position)

def create_common_callouts(table_title, column_key="Total"):
    return [
        TextCallout(table_title, column_key, "Top 3 Box", "Top 3 Box", (0.5, 7.0, 4.0, 0.4)),
        TextCallout(table_title, column_key, "Top 2 Box", "Top 2 Box", (5.0, 7.0, 4.0, 0.4)),
        TextCallout(table_title, column_key, "Bottom 2 Box", "Bottom 2 Box", (0.5, 7.5, 4.0, 0.4)),
        TextCallout(table_title, column_key, "Bottom 3 Box", "Bottom 3 Box", (5.0, 7.5, 4.0, 0.4)),
    ]

def create_statistical_callouts(table_title, column_key="Total"):
    return [
        TextCallout(table_title, column_key, "Average",     "Average",     (0.5, 7.0, 4.0, 0.4)),
        TextCallout(table_title, column_key, "Mean",        "Mean",        (5.0, 7.0, 4.0, 0.4)),
        TextCallout(table_title, column_key, "Total Spend", "Total Spend", (0.5, 7.5, 4.0, 0.4)),
    ]

def create_custom_callout(table_title, column_key="Total", row_label=None,
                          custom_text=None, position=(0.5, 7.0, 9.0, 0.4)):
    return TextCallout(table_title=table_title, column_key=column_key,
                       row_label=row_label or "Custom", text=custom_text, position=position)
