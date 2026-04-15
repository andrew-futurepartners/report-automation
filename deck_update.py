"""
deck_update.py — Update an existing PowerPoint presentation with new crosstab data.

Matches shapes to tables via alt-text metadata, updates chart data (preserving
formatting), refreshes table cells, question/base/title text, and callouts.
"""

import logging
import math
import re
from typing import Dict, Any, List, Optional, Tuple

from pptx import Presentation

from chart_data_patcher import detect_value_format, patch_chart_data, patch_chart_series
from crosstab_parser import parse_workbook
from smart_match import SmartMatcher
from text_utils import (
    parse_base_text,
    format_base_text,
    format_number_with_commas,
    safe_update_text,
)

logger = logging.getLogger("report_relay.deck_update")

EXCLUDE_PREFIXES = ("base", "mean", "average", "avg")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "")).strip().lower()


def _parse_alt_text(shape) -> Dict[str, str]:
    """Parse alt text with enhanced flexibility for manual editing."""
    out: Dict[str, str] = {}

    try:
        if hasattr(shape, "element"):
            c_nv_pr = None
            if "graphicFrame" in shape.element.tag:
                c_nv_pr = shape.element.find(
                    ".//p:cNvPr",
                    namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"},
                )
            elif "sp" in shape.element.tag:
                c_nv_pr = shape.element.find(
                    ".//p:cNvPr",
                    namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"},
                )
            alt = c_nv_pr.get("descr") if c_nv_pr is not None and c_nv_pr.get("descr") else ""
        else:
            alt = ""
    except (AttributeError, TypeError):
        alt = ""

    if not alt:
        try:
            alt = shape.alternative_text or ""
        except (AttributeError, ValueError):
            alt = ""

    for line in alt.splitlines():
        line = line.strip()
        if ":" in line:
            if " : " in line:
                k, v = line.split(" : ", 1)
            else:
                k, v = line.split(":", 1)
            out[_norm(k)] = v.strip()
    return out


def _exclude_indices(labels: List[str], extra_excludes: Optional[List[str]] = None) -> set:
    """Return indices to exclude based on default prefixes and optional explicit names."""
    ex = set()
    extra_norm = []
    if extra_excludes:
        for it in extra_excludes:
            if it is None:
                continue
            s = _norm(str(it))
            if s:
                extra_norm.append(s)
    for i, lab in enumerate(labels):
        if not isinstance(lab, str):
            continue
        nlab = _norm(lab)
        if nlab.startswith(EXCLUDE_PREFIXES):
            ex.add(i)
            continue
        for pat in extra_norm:
            if not pat:
                continue
            if nlab == pat or nlab.startswith(pat):
                ex.add(i)
                break
    return ex


def _row_index_map(labels: List[str]) -> Dict[str, int]:
    return {_norm(l): i for i, l in enumerate(labels)}


def _choose_col_idx(col_labels: List[str], col_key: Optional[str]) -> Optional[int]:
    if not col_labels:
        return None
    if col_key and col_key in col_labels:
        return col_labels.index(col_key)
    for cand in ["Total", "Overall", "All", "Base"]:
        if cand in col_labels:
            return col_labels.index(cand)
    return 0


def _series_from_table(table: Dict[str, Any], col_idx: Optional[int], exclude_rows: set):
    cats, vals = [], []
    row_labels = table["row_labels"]
    values = table["values"]
    for i, lab in enumerate(row_labels):
        if i in exclude_rows:
            continue
        cats.append(lab)
        if col_idx is None:
            vals.append(None)
        else:
            row = values[i]
            val = row[col_idx] if col_idx < len(row) else None
            if val is not None:
                try:
                    float_val = float(val)
                    if math.isnan(float_val) or math.isinf(float_val):
                        val = None
                    else:
                        val = float_val
                except (ValueError, TypeError):
                    val = None
            vals.append(val)
    return cats, vals


def _find_base_row_idx(row_labels: List[str]) -> Optional[int]:
    """Find the index of the 'base' row in row_labels."""
    for i, lab in enumerate(row_labels):
        if isinstance(lab, str) and _norm(lab).startswith("base"):
            return i
    return None


def _get_base_n(table: Dict[str, Any], col_key: Optional[str] = None) -> Optional[int]:
    """Extract base N from the table for a given column."""
    base_idx = _find_base_row_idx(table.get("row_labels", []))
    if base_idx is None:
        return None
    ci = _choose_col_idx(table.get("col_labels", []), col_key or "Total")
    if ci is None:
        return None
    values = table.get("values", [])
    if base_idx < len(values) and ci < len(values[base_idx]):
        try:
            return int(round(float(values[base_idx][ci])))
        except (ValueError, TypeError):
            return None
    return None


# ---------------------------------------------------------------------------
# Chart update
# ---------------------------------------------------------------------------

def _update_chart(shape, table: Dict[str, Any], col_key: Optional[str],
                  explicit_rows: Optional[List[str]],
                  exclude_terms: Optional[List[str]] = None,
                  column_keys: Optional[List[str]] = None):
    """Update chart data in-place via XML patching (preserves all formatting)."""
    chart = shape.chart
    alt = _parse_alt_text(shape)
    ex = _exclude_indices(table["row_labels"], exclude_terms)

    # --- Extract data series ---
    multi_series: Optional[List[tuple]] = None
    if column_keys and len(column_keys) >= 2:
        cats = None
        multi_series = []
        for ck in column_keys:
            ci = _choose_col_idx(table["col_labels"], ck)
            c, v = _series_from_table(table, ci, ex)
            if cats is None:
                cats = c
            multi_series.append((ck, v))
        if cats is None:
            cats = []
        vals = multi_series[0][1] if multi_series else []
    else:
        col_idx = _choose_col_idx(table["col_labels"], col_key)
        if explicit_rows:
            idx_map = _row_index_map(table["row_labels"])
            cats, vals = [], []
            for lab in explicit_rows:
                j = idx_map.get(_norm(lab))
                if j is None or j in ex:
                    continue
                cats.append(lab)
                if col_idx is None:
                    vals.append(None)
                else:
                    row = table["values"][j]
                    vals.append(row[col_idx] if col_idx < len(row) else None)
        else:
            cats, vals = _series_from_table(table, col_idx, ex)

    # --- Patch chart data at XML level (formatting stays intact) ---
    value_fmt = detect_value_format(vals, alt)

    if multi_series:
        patch_chart_series(
            chart,
            [(name, cats, v) for name, v in multi_series],
            value_format=value_fmt,
        )
    else:
        patch_chart_data(chart, cats, vals, value_format=value_fmt)

    logger.info("Updated chart data for table: %s", table.get("title"))


# ---------------------------------------------------------------------------
# Table update
# ---------------------------------------------------------------------------

def _update_table(shape, table: Dict[str, Any]):
    if not shape.has_table:
        return
    tbl = shape.table
    hdrs = []
    for c in range(1, len(tbl.columns)):
        txt = tbl.cell(0, c).text_frame.text.strip()
        hdrs.append(txt)
    col_labels = table["col_labels"]
    col_map = [col_labels.index(h) if h in col_labels else None for h in hdrs]
    idx_map = _row_index_map(table["row_labels"])

    for r in range(1, len(tbl.rows)):
        rlab = tbl.cell(r, 0).text_frame.text.strip()
        j = idx_map.get(_norm(rlab))
        for c in range(1, len(tbl.columns)):
            txt = ""
            ci = col_map[c - 1]
            if j is not None and ci is not None:
                try:
                    val = table["values"][j][ci]
                    if val is not None:
                        float_val = float(val)
                        if math.isnan(float_val) or math.isinf(float_val):
                            txt = ""
                        else:
                            txt = f"{float_val:.1f}"
                    else:
                        txt = ""
                except (ValueError, TypeError, IndexError):
                    txt = ""

            safe_update_text(tbl.cell(r, c), txt)

    logger.info("Updated table data (preserving formatting) for table: %s", table.get("title"))


# ---------------------------------------------------------------------------
# Shape finders
# ---------------------------------------------------------------------------

def _find_shape(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def _find_shapes_by_pattern(slide, pattern: str):
    matches = []
    for shp in slide.shapes:
        if shp.name and pattern.lower() in shp.name.lower():
            matches.append(shp)
    return matches


# ---------------------------------------------------------------------------
# Shape → table mapping
# ---------------------------------------------------------------------------

def _match_shape(
    matcher: SmartMatcher, shape, selections: Optional[Dict[str, Any]] = None,
) -> Optional[Tuple[Dict[str, Any], Optional[str], Optional[List[str]]]]:
    """Use SmartMatcher to resolve a shape to (table, col_key, exclude_terms).

    Applies a selections override for ``column_key`` when present.
    """
    alt = _parse_alt_text(shape)
    name = shape.name or ""
    result = matcher.match({"name": name, "alt": alt})
    if result is None or result.table is None:
        return None

    table = result.table
    col_key = result.col_key
    exclude_terms = result.exclude_terms

    if table and selections and table.get("title") in selections:
        sel_col = selections[table["title"]].get("column_key")
        if sel_col:
            col_key = sel_col

    return (table, col_key, exclude_terms)


# ---------------------------------------------------------------------------
# Unified question / base / title update
# ---------------------------------------------------------------------------

def _update_question_and_base(slide, table: Dict[str, Any],
                              selections: Optional[dict] = None,
                              table_title: Optional[str] = None):
    """Update question, base, and title text shapes on a slide.

    When *selections* contains an entry for *table_title*, values from that
    selection dict drive the update (question_text, base_text, title).
    Otherwise the function falls back to crosstab-derived defaults.
    """
    title_key = table_title or table.get("title", "")
    table_selection = None
    if selections and title_key in selections:
        table_selection = selections[title_key]
    elif selections:
        logger.debug("No selection found for table: %s", title_key)

    for shape in slide.shapes:
        alt = _parse_alt_text(shape)

        # --- Question text ---
        if alt.get("type") in ("question_text", "text_question") and alt.get("table_title") == title_key:
            if not hasattr(shape, "text_frame"):
                continue
            if table_selection and "question_text" in table_selection:
                new_text = f"Question: {table_selection['question_text']}"
                safe_update_text(shape, new_text, preserve_font=True)
                logger.info("Updated question text for table: %s", title_key)
            else:
                current_text = shape.text_frame.text
                if current_text.startswith("Question: "):
                    existing_question = current_text[10:]
                    if existing_question == table.get("title", ""):
                        safe_update_text(shape, f"Question: {table.get('title', '')}")
                        logger.info("Updated question text for table: %s", table.get("title"))
                    else:
                        logger.info("Preserved custom question text: %s", existing_question)
                else:
                    safe_update_text(shape, f"Question: {table.get('title', '')}")
                    logger.info("Added question text for table: %s", table.get("title"))

        # --- Base text ---
        elif alt.get("type") == "text_base" and alt.get("table_title") == title_key:
            if not hasattr(shape, "text_frame"):
                continue
            if table_selection and "base_text" in table_selection:
                base_text_template = table_selection["base_text"]
                col_key_sel = table_selection.get("column_key")

                if col_key_sel:
                    new_n = _get_base_n(table, col_key_sel)
                else:
                    new_n = None

                if new_n is not None:
                    parsed = parse_base_text(base_text_template)
                    desc = parsed["description"] or "Total respondents"
                    new_text = format_base_text(desc, new_n)
                else:
                    new_text = base_text_template

                safe_update_text(shape, new_text, preserve_font=True)
                logger.info("Updated base text for table: %s", title_key)
            else:
                current_base_text = shape.text_frame.text
                parsed = parse_base_text(current_base_text)
                base_n = _get_base_n(table, "Total")

                desc = parsed["description"] or "Total respondents"
                new_text = format_base_text(desc, base_n)
                safe_update_text(shape, new_text)
                logger.info("Updated base text for table: %s (N=%s)", table.get("title"), format_number_with_commas(base_n))

        # --- Chart title ---
        elif alt.get("type") == "text_title" and alt.get("table_title") == title_key:
            if not hasattr(shape, "text_frame"):
                continue
            if table_selection and "title" in table_selection:
                safe_update_text(shape, table_selection["title"], preserve_font=True)
                logger.info("Updated chart title for table: %s", title_key)
            else:
                current_text = shape.text_frame.text
                if current_text == table.get("title", ""):
                    logger.info("Chart title already current for table: %s", table.get("title"))
                else:
                    logger.info("Preserved custom chart title: %s", current_text)


# ---------------------------------------------------------------------------
# Callout update
# ---------------------------------------------------------------------------

def _update_new_text_callout_system(slide, table: Dict[str, Any], col_key: Optional[str],
                                    selections: Optional[Dict[str, Any]] = None):
    """Update TextCallout shapes based on alt text mapping.

    When *selections* contains an entry for this table's title, the
    ``column_key`` from selections overrides the alt-text column, and
    callout-specific overrides (metric_type, text template) from the
    selections ``callouts`` list are applied when matched by row_label.
    """
    table_title = table.get("title", "")
    sel = selections.get(table_title, {}) if selections else {}
    sel_callouts = sel.get("callouts", [])

    for shape in slide.shapes:
        alt = _parse_alt_text(shape)

        if alt.get("type") != "text_callout" or alt.get("table_title") != table_title:
            continue
        if not hasattr(shape, "text_frame"):
            continue

        row_label = alt.get("row", alt.get("row_label", ""))
        column = alt.get("column", "Total")
        metric_type = alt.get("metric_type", "percentage")
        current_shape_text = shape.text_frame.text if hasattr(shape, "text_frame") else ""

        # Apply selection-level overrides for this callout
        if col_key:
            column = col_key
        for sc in sel_callouts:
            if sc.get("row_label") == row_label:
                if sc.get("column_key"):
                    column = sc["column_key"]
                if sc.get("metric_type"):
                    metric_type = sc["metric_type"]
                if sc.get("text"):
                    current_shape_text = sc["text"]
                break

        row_idx = None
        col_idx = None
        new_text = ""

        if row_label:
            row_labels = table.get("row_labels", [])
            for i, label in enumerate(row_labels):
                if isinstance(label, str) and row_label.lower() in label.lower():
                    row_idx = i
                    break

            col_labels = table.get("col_labels", [])
            if column in col_labels:
                col_idx = col_labels.index(column)
            else:
                for fallback in ["Total", "Overall", "All", "Base"]:
                    if fallback in col_labels:
                        col_idx = col_labels.index(fallback)
                        break
                if col_idx is None:
                    col_idx = 0 if col_labels else None

            if row_idx is not None and col_idx is not None:
                try:
                    values = table.get("values", [])
                    if row_idx < len(values) and col_idx < len(values[row_idx]):
                        value = values[row_idx][col_idx]
                        if value is not None:
                            formatted_value = ""
                            if isinstance(value, (int, float)):
                                mt = (metric_type or "").lower()
                                if mt == "percentage":
                                    formatted_value = f"{float(value) * 100:.1f}%"
                                elif mt == "currency":
                                    formatted_value = f"${float(value):,.0f}"
                                else:
                                    formatted_value = f"{float(value):,.1f}"
                            else:
                                formatted_value = str(value)

                            if current_shape_text and "[Value]" in current_shape_text:
                                new_text = current_shape_text.replace("[Value]", formatted_value)
                            elif current_shape_text:
                                pattern = re.compile(r"[-+]?\d{1,3}(?:,\d{3})*(?:\.\d+)?%?")
                                if pattern.search(current_shape_text):
                                    new_text = pattern.sub(formatted_value, current_shape_text, count=1)
                                else:
                                    new_text = f"{row_label}: {formatted_value}"
                            else:
                                new_text = f"{row_label}: {formatted_value}"
                except (IndexError, TypeError, AttributeError):
                    pass

        if not new_text:
            new_text = current_shape_text if current_shape_text else f"{row_label}: [Value]"

        current_text = shape.text_frame.text
        if current_text != new_text:
            safe_update_text(shape, new_text)
            logger.info("Updated text callout '%s' for table: %s", row_label, table.get("title"))


# ---------------------------------------------------------------------------
# Core slide-processing loop (shared by both entry points)
# ---------------------------------------------------------------------------

def _process_slides(prs, data: Dict[str, Any], selections: Optional[dict] = None,
                    matcher: Optional[SmartMatcher] = None,
                    progress_callback=None) -> dict:
    """Walk all slides/shapes, update charts/tables/text, return update counts.

    *progress_callback*, when provided, is called with a float 0.0–1.0
    after each slide is processed.
    """
    if matcher is None:
        matcher = SmartMatcher(data["tables"])

    update_log = {
        "charts_updated": 0,
        "tables_updated": 0,
        "text_updated": 0,
        "shapes_skipped": 0,
        "matched_titles": set(),
        "match_report": [],
    }

    slides = list(prs.slides)
    total_slides = len(slides) or 1

    for slide_idx, slide in enumerate(slides):
        for shp in slide.shapes:
            name = shp.name or ""
            alt = _parse_alt_text(shp)

            if alt.get("auto_update", "yes").lower() == "no":
                update_log["shapes_skipped"] += 1
                continue

            # --- Charts ---
            try:
                _chart_obj = shp.chart  # noqa: F841 — access triggers AttributeError for non-charts
                mapping = _match_shape(matcher, shp, selections)
                if mapping:
                    table, col_key, exclude_terms = mapping
                    sel_col_keys = None
                    table_title = table.get("title")
                    if selections and table_title and table_title in selections:
                        sel_col_keys = selections[table_title].get("column_keys")
                    _update_chart(shp, table, col_key, explicit_rows=None,
                                  exclude_terms=exclude_terms, column_keys=sel_col_keys)

                    sel_for_qb = None
                    table_title = table.get("title")
                    if selections and table_title and table_title in selections:
                        sel_for_qb = selections

                    _update_question_and_base(slide, table, sel_for_qb, table_title)
                    _update_new_text_callout_system(slide, table, col_key, selections)

                    update_log["charts_updated"] += 1
                    update_log["matched_titles"].add(table_title)
                    logger.info("Updated chart with mapping for table: %s", table_title)
                else:
                    logger.debug("Chart '%s' has no table mapping - preserving as-is", name)
                    update_log["shapes_skipped"] += 1
            except (ValueError, AttributeError):
                pass

            # --- Tables ---
            if shp.has_table:
                mapping = _match_shape(matcher, shp, selections)
                if mapping:
                    table, col_key, exclude_terms = mapping
                    _update_table(shp, table)

                    sel_for_qb = None
                    table_title = table.get("title")
                    if selections and table_title and table_title in selections:
                        sel_for_qb = selections

                    _update_question_and_base(slide, table, sel_for_qb, table_title)

                    callout_col = col_key
                    if selections and table_title and table_title in selections:
                        callout_col = selections[table_title].get("column_key") or col_key
                    _update_new_text_callout_system(slide, table, callout_col, selections)

                    update_log["tables_updated"] += 1
                    update_log["matched_titles"].add(table_title)
                    logger.info("Updated table with mapping for table: %s", table_title)
                else:
                    logger.debug("Table '%s' has no table mapping - preserving as-is", name)
                    update_log["shapes_skipped"] += 1

            # Legacy named text objects
            if name in ("TEXT_QUESTION", "OBJ_QUESTION", "TEXT_BASE", "OBJ_BASE"):
                update_log["text_updated"] += 1

        if progress_callback:
            progress_callback((slide_idx + 1) / total_slides)

    update_log["match_report"] = matcher.get_report()
    return update_log


def _log_update_summary(update_log: dict, unmapped_tables: Optional[list] = None):
    """Emit a structured update summary to the logger."""
    logger.info("=" * 50)
    logger.info("UPDATE SUMMARY")
    logger.info("=" * 50)
    logger.info("Charts updated: %d", update_log["charts_updated"])
    logger.info("Tables updated: %d", update_log["tables_updated"])
    logger.info("Text objects updated: %d", update_log["text_updated"])
    logger.info("Shapes preserved (no mapping): %d", update_log["shapes_skipped"])

    if unmapped_tables:
        logger.info("Unmapped tables added to summary page: %d", len(unmapped_tables))
        for t in unmapped_tables[:5]:
            logger.info("  - %s", t["title"])
        if len(unmapped_tables) > 5:
            logger.info("  ... and %d more", len(unmapped_tables) - 5)

    logger.info("=" * 50)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def update_presentation(pptx_in: str, crosstab_xlsx: str, pptx_out: str,
                        selections: dict = None,
                        all_tables: list = None,
                        existing_content: dict = None,
                        include_unmapped_summary: bool = False,
                        matcher: Optional[SmartMatcher] = None,
                        progress_callback=None) -> str:
    """Update an existing PowerPoint with new crosstab data.

    Args:
        selections: Optional mapping keyed by **table title** (not id)::

            {table_title: {"column_key": str, "title": str,
                           "base_text": str, "question_text": str,
                           "callouts": list, "chart_type": str,
                           "enable_sorting": bool, "excluded_rows": list}}

        matcher: Optional pre-configured ``SmartMatcher``.  When *None*,
            one is created internally from the parsed workbook tables.

        progress_callback: Optional callable(float) invoked with 0.0–1.0
            progress after each slide.

        When *include_unmapped_summary* is True and *all_tables* /
        *existing_content* are provided, tables not present in the
        existing deck are listed on an appended summary slide.
    """
    prs = Presentation(pptx_in)
    data = parse_workbook(crosstab_xlsx)

    if matcher is None:
        matcher = SmartMatcher(data["tables"])

    if selections:
        logger.debug("Selections provided: %s", list(selections.keys()))

    update_log = _process_slides(prs, data, selections, matcher=matcher,
                                 progress_callback=progress_callback)

    # --- Unmapped tables summary slide ---
    unmapped_tables = []
    if include_unmapped_summary and all_tables and existing_content is not None:
        mapped_titles = set(existing_content.keys()) if existing_content else set()
        unmapped_tables = [t for t in all_tables if t["title"] not in mapped_titles]

        if unmapped_tables:
            _add_unmapped_summary_slide(prs, unmapped_tables)

    _log_update_summary(update_log, unmapped_tables if include_unmapped_summary else None)

    prs.save(pptx_out)
    return pptx_out


def update_presentation_with_unmapped(pptx_in: str, crosstab_xlsx: str, pptx_out: str,
                                      selections: dict = None,
                                      all_tables: list = None,
                                      existing_content: dict = None,
                                      matcher: Optional[SmartMatcher] = None,
                                      progress_callback=None) -> str:
    """Convenience wrapper — calls update_presentation with unmapped summary enabled."""
    return update_presentation(
        pptx_in, crosstab_xlsx, pptx_out,
        selections=selections,
        all_tables=all_tables,
        existing_content=existing_content,
        include_unmapped_summary=True,
        matcher=matcher,
        progress_callback=progress_callback,
    )


# ---------------------------------------------------------------------------
# Unmapped tables summary slide
# ---------------------------------------------------------------------------

def _add_unmapped_summary_slide(prs, unmapped_tables: list):
    """Append a slide listing tables that had no shape match in the deck."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    unmapped_slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_box = unmapped_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Unmapped Tables Summary"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Arial"

    subtitle_box = unmapped_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.6))
    stf = subtitle_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = (
        f"The following {len(unmapped_tables)} tables from your crosstab had no "
        "existing connections and are listed here for reference:"
    )
    srun.font.size = Pt(14)
    srun.font.name = "Arial"

    list_y_start = 2.0
    line_height = 0.3

    for i, table in enumerate(unmapped_tables):
        title_box = unmapped_slide.shapes.add_textbox(
            Inches(0.5), Inches(list_y_start + i * line_height * 4), Inches(8.5), Inches(0.25)
        )
        ttf = title_box.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        trun = tp.add_run()
        trun.text = f"\u2022 {table['title']}"
        trun.font.size = Pt(12)
        trun.font.bold = True
        trun.font.name = "Arial"

        stats_box = unmapped_slide.shapes.add_textbox(
            Inches(0.8), Inches(list_y_start + i * line_height * 4 + 0.25), Inches(8.2), Inches(0.2)
        )
        stats_tf = stats_box.text_frame
        stats_tf.clear()
        stats_p = stats_tf.paragraphs[0]
        stats_run = stats_p.add_run()
        row_count = len(table.get("row_labels", []))
        col_count = len(table.get("col_labels", []))
        stats_run.text = f"  Rows: {row_count}, Columns: {col_count}"
        stats_run.font.size = Pt(10)
        stats_run.font.name = "Arial"
        stats_run.font.color.rgb = RGBColor(100, 100, 100)

        if table.get("col_labels"):
            cols_text = ", ".join(table["col_labels"][:8])
            if len(table["col_labels"]) > 8:
                cols_text += "..."
            cols_box = unmapped_slide.shapes.add_textbox(
                Inches(0.8), Inches(list_y_start + i * line_height * 4 + 0.45), Inches(8.2), Inches(0.2)
            )
            cols_tf = cols_box.text_frame
            cols_tf.clear()
            cols_p = cols_tf.paragraphs[0]
            cols_run = cols_p.add_run()
            cols_run.text = f"  Columns: {cols_text}"
            cols_run.font.size = Pt(9)
            cols_run.font.name = "Arial"
            cols_run.font.color.rgb = RGBColor(120, 120, 120)

        if i >= 11:
            remaining = len(unmapped_tables) - 12
            if remaining > 0:
                more_box = unmapped_slide.shapes.add_textbox(
                    Inches(0.5), Inches(list_y_start + 12 * line_height * 4), Inches(8.5), Inches(0.25)
                )
                more_tf = more_box.text_frame
                more_tf.clear()
                more_p = more_tf.paragraphs[0]
                more_run = more_p.add_run()
                more_run.text = f"... and {remaining} more tables"
                more_run.font.size = Pt(11)
                more_run.font.italic = True
                more_run.font.name = "Arial"
                more_run.font.color.rgb = RGBColor(150, 150, 150)
            break
