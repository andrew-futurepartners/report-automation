import logging
import os
import tempfile

import streamlit as st
from crosstab_parser import parse_workbook
from pptx_exporter import export_pptx
from deck_update import update_presentation, update_presentation_with_unmapped, _parse_alt_text
from pptx import Presentation
from smart_match import SmartMatcher
from ai_insights import generate_all_insights
from text_utils import format_number_with_commas, parse_base_text, format_base_text

logging.basicConfig(
    level=getattr(logging, os.environ.get("LOG_LEVEL", "INFO").upper(), logging.INFO),
    format="%(levelname)s | %(name)s | %(message)s",
)


def _save_temp(data_buf, suffix: str, session_key: str) -> str:
    """Write *data_buf* to a NamedTemporaryFile, store the path in session_state,
    and clean up any previous temp file stored under the same key."""
    old = st.session_state.get(session_key)
    if old and os.path.exists(old):
        try:
            os.unlink(old)
        except OSError:
            pass
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp.write(data_buf if isinstance(data_buf, bytes) else data_buf.getbuffer())
    tmp.close()
    st.session_state[session_key] = tmp.name
    return tmp.name


def _render_diff_table(old_tv: dict, new_table: dict):
    """Render a side-by-side diff of old PPT table values vs new crosstab data.

    *old_tv* has ``{"headers": [...], "rows": [[...], ...]}``.
    *new_table* is a standard crosstab table dict.
    """
    import pandas as pd

    old_headers = old_tv.get("headers", [])
    old_rows = old_tv.get("rows", [])
    if not old_headers or not old_rows:
        st.write("No previous table data available.")
        return

    old_row_labels = [r[0] for r in old_rows] if old_rows else []
    old_col_headers = old_headers[1:] if len(old_headers) > 1 else old_headers

    new_row_labels = new_table.get("row_labels", [])
    new_col_labels = new_table.get("col_labels", [])
    new_values = new_table.get("values", [])

    all_rows = list(dict.fromkeys(old_row_labels + [str(r) for r in new_row_labels]))
    common_cols = [c for c in old_col_headers if c in new_col_labels]
    if not common_cols:
        common_cols = old_col_headers[:5]

    old_map = {}
    for r in old_rows:
        if r:
            old_map[r[0]] = r[1:] if len(r) > 1 else []

    new_map = {}
    for i, rl in enumerate(new_row_labels):
        row_vals = new_values[i] if i < len(new_values) else []
        cells = {}
        for j, cl in enumerate(new_col_labels):
            cells[cl] = row_vals[j] if j < len(row_vals) else None
        new_map[str(rl)] = cells

    diff_rows = []
    statuses = []
    for rl in all_rows:
        row_data = {"Row": rl}
        in_old = rl in old_map
        in_new = rl in new_map
        if in_old and not in_new:
            statuses.append("removed")
        elif in_new and not in_old:
            statuses.append("added")
        else:
            statuses.append("existing")
        for ci, col in enumerate(common_cols):
            old_val = ""
            if in_old and ci < len(old_map.get(rl, [])):
                old_val = old_map[rl][ci]
            new_val = ""
            if in_new:
                nv = new_map.get(rl, {}).get(col)
                if nv is not None:
                    try:
                        fv = float(nv)
                        new_val = f"{fv:.1f}" if fv != int(fv) or fv <= 100 else f"{int(fv)}"
                    except (ValueError, TypeError):
                        new_val = str(nv)
            changed = str(old_val).strip() != str(new_val).strip()
            marker = " *" if changed and in_old and in_new else ""
            row_data[f"{col} (old)"] = old_val
            row_data[f"{col} (new)"] = f"{new_val}{marker}"
        diff_rows.append(row_data)

    df = pd.DataFrame(diff_rows)

    def _highlight(row):
        idx = diff_rows.index(row.to_dict()) if row.to_dict() in diff_rows else -1
        status = statuses[idx] if 0 <= idx < len(statuses) else "existing"
        if status == "added":
            return ["background-color: #d4edda"] * len(row)
        if status == "removed":
            return ["background-color: #f8d7da"] * len(row)
        styles = []
        for val in row:
            if isinstance(val, str) and val.endswith(" *"):
                styles.append("background-color: #fff3cd")
            else:
                styles.append("")
        return styles

    st.dataframe(df.style.apply(_highlight, axis=1), use_container_width=True, hide_index=True)

    legend_cols = st.columns(3)
    legend_cols[0].markdown('<span style="background:#d4edda;padding:2px 8px;border-radius:3px;">Added rows</span>', unsafe_allow_html=True)
    legend_cols[1].markdown('<span style="background:#fff3cd;padding:2px 8px;border-radius:3px;">Changed values (*)</span>', unsafe_allow_html=True)
    legend_cols[2].markdown('<span style="background:#f8d7da;padding:2px 8px;border-radius:3px;">Removed rows</span>', unsafe_allow_html=True)


def parse_existing_powerpoint(pptx_file):
    """Parse existing PowerPoint to extract current content and settings."""
    try:
        prs = Presentation(pptx_file)
        existing_content = {}
        
        for slide in prs.slides:
            for shape in slide.shapes:
                alt = _parse_alt_text(shape)
                
                if alt.get("type") in ["chart", "table", "question_text", "text_question", "text_base", "text_title", "text_callout", "text_takeaway", "text_analysis", "text_chart_title", "ai_insight"]:
                    table_title = alt.get("table_title")
                    if table_title:
                        if table_title not in existing_content:
                            existing_content[table_title] = {
                                "title": table_title,
                                "question_text": "",
                                "base_text": "",
                                "chart_type": "bar_h",
                                "custom_base_description": "",  # Store custom base description
                                "custom_question": "",  # Store custom question text
                                "callouts": [],  # Initialize callouts list
                                # Presence flags for connected objects
                                "has_chart": False,
                                "has_table": False,
                                "has_title": False,
                                "has_base": False,
                                "has_question": False,
                                "has_callouts": False
                            }
                        
                        # Track object presence
                        if alt.get("type") == "chart":
                            existing_content[table_title]["has_chart"] = True
                            if "column" in alt:
                                existing_content[table_title]["chart_column"] = alt.get("column")
                        elif alt.get("type") == "table":
                            existing_content[table_title]["has_table"] = True
                            if shape.has_table:
                                tbl = shape.table
                                headers = []
                                for c in range(len(tbl.columns)):
                                    headers.append(tbl.cell(0, c).text_frame.text.strip())
                                rows_data = []
                                for r in range(1, len(tbl.rows)):
                                    row_vals = []
                                    for c in range(len(tbl.columns)):
                                        row_vals.append(tbl.cell(r, c).text_frame.text.strip())
                                    rows_data.append(row_vals)
                                existing_content[table_title]["table_values"] = {
                                    "headers": headers,
                                    "rows": rows_data,
                                }

                        # Extract question text
                        if alt.get("type") in ["question_text", "text_question"] and hasattr(shape, "text_frame"):
                            question_text = shape.text_frame.text
                            if question_text.startswith("Question: "):
                                question_text = question_text[10:]  # Remove "Question: " prefix
                            existing_content[table_title]["question_text"] = question_text
                            existing_content[table_title]["custom_question"] = question_text
                            existing_content[table_title]["has_question"] = True
                        
                        # Extract base text
                        elif alt.get("type") == "text_base" and hasattr(shape, "text_frame"):
                            base_text = shape.text_frame.text
                            existing_content[table_title]["base_text"] = base_text
                            existing_content[table_title]["has_base"] = True
                            
                            parsed_base = parse_base_text(base_text)
                            if parsed_base["description"]:
                                existing_content[table_title]["custom_base_description"] = parsed_base["description"]
                        
                        # Extract chart title
                        elif alt.get("type") == "text_title" and hasattr(shape, "text_frame"):
                            existing_content[table_title]["title"] = shape.text_frame.text
                            existing_content[table_title]["has_title"] = True
                        
                        # Extract text callouts
                        elif alt.get("type") == "text_callout" and hasattr(shape, "text_frame"):
                            if "callouts" not in existing_content[table_title]:
                                existing_content[table_title]["callouts"] = []
                            
                            callout_text = shape.text_frame.text
                            callout_info = {
                                "row_label": alt.get("row", alt.get("row_label", "")),
                                "column_key": alt.get("column", "Total"),
                                "text": callout_text,
                                "position": (0.5, 7.0, 9.0, 0.4),  # Default position
                                "font_size": 12,
                                "font_bold": True,
                                "metric_type": alt.get("metric_type", "percentage")
                            }
                            existing_content[table_title]["callouts"].append(callout_info)
                            existing_content[table_title]["has_callouts"] = True
        
        return existing_content
    except Exception as e:
        st.error(f"Error parsing PowerPoint: {e}")
        return {}

st.set_page_config(page_title="Report Relay", layout="wide")

st.title("Report Relay")
st.write("An efficient handoff from crosstab to report.")

# Step 1: Choose workflow type
st.subheader("Step 1: Choose your workflow")
workflow_type = st.radio(
    "Are you creating a new report or updating an existing one?",
    ["Create New Report", "Update Existing Report"],
    index=0,
    help="""
    **Create New Report**: Upload crosstab Excel → Choose chart types and customize → Export new PowerPoint
    
    **Update Existing Report**: Upload PowerPoint → Detect connections → Upload crosstab → Review mappings → Update PowerPoint
    """
)

# Add helpful workflow information
with st.expander("💡 How each workflow works", expanded=False):
    if workflow_type == "Create New Report":
        st.write("""
        **New Report Workflow:**
        1. Upload crosstab Excel file
        2. Choose chart types and customize titles for each table
        3. Export new PowerPoint presentation
        
        **Perfect for:** First-time report creation, one-off reports
        """)
    else:
        st.write("""
        **Update Existing Report Workflow:**
        1. Upload existing PowerPoint presentation
        2. System detects existing table connections
        3. Upload new crosstab Excel file  
        4. Review tables with/without existing connections
        5. Configure column selection (switch segments quickly)
        6. Update presentation (preserves formatting, adds unmapped tables summary)
        
        **Perfect for:** Regular report updates, preserving custom formatting and content
        
        **Key Benefits:**
        - Custom work is preserved when updating reports
        - Base N values are automatically refreshed with new data
        - Chart titles, questions, and base descriptions can be customized and preserved
        - New tables are automatically listed on a summary page
        - Quick segment switching (Total → Male → Female → Gen Z, etc.)
        """)

# Initialize session state for workflow management
if "workflow_step" not in st.session_state:
    st.session_state.workflow_step = "start"
if "existing_content" not in st.session_state:
    st.session_state.existing_content = {}
if "data" not in st.session_state:
    st.session_state.data = None
if "selections" not in st.session_state:
    st.session_state["selections"] = {}

# Workflow branching
if workflow_type == "Create New Report":
    # NEW REPORT WORKFLOW
    st.subheader("Step 2: Upload crosstab data")
    uploaded = st.file_uploader("Upload crosstab Excel", type=["xlsx", "xls"], key="new_report_excel")
    
    if uploaded:
        with st.spinner("Parsing workbook..."):
            xlsx_path = _save_temp(uploaded, ".xlsx", "_tmp_xlsx")
            data = parse_workbook(xlsx_path)
            st.session_state.data = data

        st.success(f"Found {len(data['tables'])} tables")
        st.divider()
        
        # Show default choice controls for new reports
        st.subheader("Step 3: Configure visualizations")
        
        # Build banner list preserving datafile order (from the first table that has banners)
        all_columns = []
        found_order = False
        for table in data["tables"]:
            meta = table.get("meta", {})
            banners = meta.get("col_banners") or table.get("col_labels", [])
            if banners and not found_order:
                for b in banners:
                    if b is None:
                        continue
                    norm_b = str(b).replace("\xa0", " ").strip()
                    if norm_b and norm_b not in all_columns:
                        all_columns.append(norm_b)
                found_order = len(all_columns) > 0
        # Fallback: aggregate in encountered order if no table provided banners
        if not all_columns:
            for table in data["tables"]:
                banners = table.get("col_labels", [])
                for b in banners:
                    if b is None:
                        continue
                    norm_b = str(b).replace("\xa0", " ").strip()
                    if norm_b and norm_b not in all_columns:
                        all_columns.append(norm_b)
        
        # Default selections
        col1, col2 = st.columns(2)
        with col1:
            default_choice = st.selectbox("Default visualization", ["Horizontal Bar", "Vertical Bar", "Donut", "Line", "Grouped Bar (2)", "Grouped Bar (3)", "Multi-Line", "Table Only", "Chart + Table"], index=0)
            apply_all_btn = st.button("Apply chart type to all")
        
        with col2:
            default_column = st.selectbox("Default column", all_columns, index=all_columns.index("Total") if "Total" in all_columns else 0)
            apply_column_btn = st.button("Apply column to all")

else:
    # UPDATE EXISTING REPORT WORKFLOW
    st.subheader("Step 2: Upload existing PowerPoint")
    existing_ppt = st.file_uploader("Upload the PowerPoint to update", type=["pptx"], key="existing_ppt")
    
    if existing_ppt:
        with st.spinner("Parsing existing PowerPoint..."):
            existing_content = parse_existing_powerpoint(existing_ppt)
            st.session_state.existing_content = existing_content
            _save_temp(existing_ppt, ".pptx", "_tmp_pptx")

        if existing_content:
            st.success(f"Found connections for {len(existing_content)} tables")
            
            # Show found connections
            with st.expander(f"📊 Found connections for {len(existing_content)} tables", expanded=True):
                for table_title in existing_content.keys():
                    st.write(f"• {table_title}")
        else:
            st.info("No existing table connections found in this PowerPoint.")
        
        st.divider()
        
        # Step 3: Upload crosstab for updates
        st.subheader("Step 3: Upload new crosstab data")
        uploaded = st.file_uploader("Upload crosstab Excel", type=["xlsx", "xls"], key="update_report_excel")
        
        if uploaded:
            with st.spinner("Parsing workbook..."):
                xlsx_path = _save_temp(uploaded, ".xlsx", "_tmp_xlsx")
                data = parse_workbook(xlsx_path)
                st.session_state.data = data

            st.success(f"Found {len(data['tables'])} tables in crosstab")
            
            # Categorize tables into connected vs unconnected
            connected_tables = []
            unconnected_tables = []
            
            for table in data["tables"]:
                if table["title"] in existing_content:
                    connected_tables.append(table)
                else:
                    unconnected_tables.append(table)
            
            st.divider()
            
            # Show categorized results
            st.subheader("Step 4: Review table connections")
            
            # Connected tables section (collapsible)
            if connected_tables:
                with st.expander(f"✅ Tables with existing connections ({len(connected_tables)})", expanded=True):
                    st.success("These tables will be updated with new data while preserving your custom formatting.")
                    for table in connected_tables:
                        st.write(f"• {table['title']}")
            
            # Unconnected tables section (collapsible)  
            if unconnected_tables:
                with st.expander(f"➕ Tables with NO connections ({len(unconnected_tables)})", expanded=True):
                    st.info("These tables are new or don't have existing connections. They will be added to a summary page.")
                    for table in unconnected_tables:
                        st.write(f"• {table['title']}")
            
            # Show column selection for update workflow
            if data["tables"]:
                # Build banner list preserving datafile order (from the first table that has banners)
                all_columns = []
                found_order = False
                for table in data["tables"]:
                    meta = table.get("meta", {})
                    banners = meta.get("col_banners") or table.get("col_labels", [])
                    if banners and not found_order:
                        for b in banners:
                            if b is None:
                                continue
                            norm_b = str(b).replace("\xa0", " ").strip()
                            if norm_b and norm_b not in all_columns:
                                all_columns.append(norm_b)
                        found_order = len(all_columns) > 0
                if not all_columns:
                    for table in data["tables"]:
                        banners = table.get("col_labels", [])
                        for b in banners:
                            if b is None:
                                continue
                            norm_b = str(b).replace("\xa0", " ").strip()
                            if norm_b and norm_b not in all_columns:
                                all_columns.append(norm_b)
                
                st.subheader("Step 5: Configure column selection")
                
                col1, _ = st.columns([1, 1])
                with col1:
                    default_column = st.selectbox("Default column for updates", all_columns, 
                                                index=all_columns.index("Total") if "Total" in all_columns else 0,
                                                key="update_default_column")
                    apply_column_btn = st.button("Apply column to all connected tables", key="update_apply_column")
                
            
            # Set defaults for the rest of the logic
            default_choice = "Horizontal Bar"
            apply_all_btn = False
        else:
            # If no data yet, set default values
            default_column = "Total"
            apply_column_btn = False

# Only show table configuration if we have data loaded
if st.session_state.data is not None:
    data = st.session_state.data
    existing_content = st.session_state.existing_content
    
    # Ensure default_column and apply_column_btn are defined for both workflows
    if workflow_type == "Create New Report":
        # These were defined earlier in the new report section
        pass
    else:
        # For update workflow, check if they were defined in the update section
        if 'default_column' not in locals():
            default_column = "Total"
            apply_column_btn = False
    
    # Determine which tables to show configuration for
    if workflow_type == "Create New Report":
        # Show all tables for new reports
        tables_to_configure = data["tables"]
        config_title = "Configure all tables"
    else:
        # For updates, only show connected tables that have existing settings
        # These can be modified if needed
        tables_to_configure = []
        for table in data["tables"]:
            if table["title"] in existing_content:
                tables_to_configure.append(table)
        config_title = f"Configure existing tables ({len(tables_to_configure)})"
    
    # Only show configuration section if there are tables to configure
    if tables_to_configure:
        st.divider()
        st.subheader(config_title)

    # Per-table controls
    for t in tables_to_configure:
        tid = t["id"]
        with st.expander(f"{t['title']}  ({tid})", expanded=False):
            cols = st.columns([2, 1, 2])
            # Lookup any existing mapping info for this table (from PPT)
            existing_table = existing_content.get(t["title"], {})
            with cols[0]:
                st.write("**Row labels**")
                # Show row labels as a concise bullet list; limit to first 12
                _all_rows = [rl for rl in t["row_labels"] if isinstance(rl, str) and rl.strip()]
                _preview = _all_rows[:12]
                if _preview:
                    st.markdown("\n".join([f"- {rl}" for rl in _preview]))
                _remaining = max(0, len(_all_rows) - len(_preview))
                if _remaining > 0:
                    st.markdown(f"- … and {_remaining} more")

            with cols[1]:
                options = ["Horizontal Bar", "Vertical Bar", "Donut", "Line", "Grouped Bar (2)", "Grouped Bar (3)", "Multi-Line", "Table Only", "Chart + Table"]
                if apply_all_btn:
                    choice = default_choice
                else:
                    choice = st.session_state["selections"].get(tid, {}).get("chart_type_label", default_choice)
                choice = st.selectbox("Chart type", options, key=f"ctype_{tid}", index=options.index(choice))
                
                # Column selection for this table
                # Keep Data column as banners only, add Metric selector when available
                meta = t.get("meta", {})
                banners = meta.get("col_banners") or t.get("col_labels", [])
                groups = meta.get("col_groups") or ["" for _ in banners]
                combined_labels = t.get("col_labels", [])

                # Derive metric options (unique, order-preserving, non-empty)
                metric_options = []
                seen_metrics = set()
                for g in groups:
                    if isinstance(g, str):
                        g2 = g.strip()
                        if g2 and g2 not in seen_metrics:
                            seen_metrics.add(g2)
                            metric_options.append(g2)
                has_metrics = len(metric_options) > 0
                if not has_metrics:
                    metric_options = [""]

                # Default metric: prefer session state; else metric parsed from mapped chart column; else left-most
                leftmost_metric = next((g for g in groups if isinstance(g, str) and g.strip()), "")
                current_metric = st.session_state["selections"].get(tid, {}).get("metric_key")
                if not current_metric:
                    mapped_col = existing_table.get("chart_column")
                    if isinstance(mapped_col, str) and "|" in mapped_col:
                        try:
                            banner_part, metric_part = [p.strip() for p in mapped_col.split("|", 1)]
                            current_metric = metric_part
                        except ValueError:
                            current_metric = None
                if not current_metric:
                    current_metric = leftmost_metric
                if current_metric not in metric_options:
                    current_metric = metric_options[0]

                if has_metrics:
                    current_metric = st.selectbox("Metric", metric_options, key=f"metric_{tid}", index=metric_options.index(current_metric))
                
                # Handle column application for both new reports and updates
                if workflow_type == "Create New Report" and apply_column_btn:
                    # Apply default banner to all tables for new reports
                    current_col = default_column if default_column in banners else ("Total" if "Total" in banners else banners[0] if banners else "Total")
                elif workflow_type == "Update Existing Report" and apply_column_btn:
                    # Apply default banner to all connected tables for updates
                    current_col = default_column if default_column in banners else ("Total" if "Total" in banners else banners[0] if banners else "Total")
                else:
                    # Use existing banner selection or mapped chart banner; else default
                    current_col = st.session_state["selections"].get(tid, {}).get("banner_key")
                    if not current_col:
                        mapped_col = existing_table.get("chart_column")
                        if isinstance(mapped_col, str):
                            if "|" in mapped_col:
                                try:
                                    banner_part, metric_part = [p.strip() for p in mapped_col.split("|", 1)]
                                    current_col = banner_part
                                except ValueError:
                                    current_col = mapped_col.strip()
                            else:
                                current_col = mapped_col.strip()
                    if not current_col:
                        current_col = "Total"
                    if current_col not in banners:
                        current_col = "Total" if "Total" in banners else banners[0] if banners else "Total"
                
                selected_col = st.selectbox("Data column", banners, 
                                         index=banners.index(current_col) if current_col in banners else 0,
                                         key=f"col_{tid}")

                # Resolve combined label for selected (metric, banner)
                pair_to_combined = {}
                for i, b in enumerate(banners):
                    g = groups[i] if i < len(groups) else ""
                    lab = combined_labels[i] if i < len(combined_labels) else b
                    pair_to_combined[(g or "", b)] = lab
                combined_selected = pair_to_combined.get(((current_metric or ""), selected_col))
                if not combined_selected:
                    # Fallback to the first column with this banner
                    for i, b in enumerate(banners):
                        if b == selected_col:
                            combined_selected = combined_labels[i]
                            break
                if not combined_selected and combined_labels:
                    combined_selected = combined_labels[0]

                # Multi-column selection for grouped / multi-series chart types
                multi_series_types = {"Grouped Bar (2)", "Grouped Bar (3)", "Multi-Line"}
                multi_columns_selected = None
                if choice in multi_series_types:
                    n_series = 3 if choice == "Grouped Bar (3)" else 2
                    default_multi = st.session_state["selections"].get(tid, {}).get("column_keys", [])
                    if not default_multi or len(default_multi) < 2:
                        default_multi = combined_labels[:n_series] if len(combined_labels) >= n_series else combined_labels
                    multi_columns_selected = st.multiselect(
                        f"Select {n_series} columns for series",
                        combined_labels,
                        default=default_multi[:n_series],
                        key=f"multi_col_{tid}",
                        max_selections=n_series if choice != "Multi-Line" else 5,
                    )

            with cols[2]:
                # Use existing content if available, otherwise use defaults
                has_title_obj = existing_table.get("has_title", False)
                has_base_obj = existing_table.get("has_base", False)
                has_question_obj = existing_table.get("has_question", False)
                
                # Chart title - prioritize existing custom title, then session state, then table title
                existing_title = existing_table.get("title", "")
                session_title = st.session_state["selections"].get(tid, {}).get("title", "")
                
                # Priority: existing content > session state > table title
                if existing_title and existing_title != t["title"]:
                    default_title = existing_title
                elif session_title:
                    default_title = session_title
                else:
                    default_title = t["title"]
                
                # Show indicator if using existing custom title
                title_label = "Chart title"
                if existing_title and existing_title != t["title"]:
                    title_label += " (Previously: " + existing_title + ")"
                
                title_val = None
                if has_title_obj:
                    title_val = st.text_input(title_label, value=default_title, key=f"title_{tid}")

                # Base text logic - preserve custom descriptions while updating N values
                def _find_base_idx(labels):
                    for i, lab in enumerate(labels):
                        if isinstance(lab, str) and lab.strip().lower().startswith("base"):
                            return i
                    return None
                
                base_idx = _find_base_idx(t["row_labels"])
                # Determine the Total column under the selected metric if applicable
                total_idx = None
                total_candidates = []
                if (" | " in (combined_selected or "")):
                    # Prefer banner-qualified by current metric first (Banner | Metric)
                    metric_name = (current_metric or "").strip()
                    if metric_name:
                        total_candidates.append(f"Total | {metric_name}")
                total_candidates.append("Total")
                for cand in total_candidates:
                    if cand in t["col_labels"]:
                        total_idx = t["col_labels"].index(cand)
                        break
                if total_idx is None:
                    total_idx = 0 if t["col_labels"] else None
                
                # Calculate new base N value
                new_base_n = None
                if base_idx is not None and total_idx is not None and base_idx < len(t["values"]) and total_idx < len(t["values"][base_idx]):
                    try:
                        new_base_n = int(round(float(t["values"][base_idx][total_idx])))
                    except (ValueError, TypeError, IndexError):
                        new_base_n = None
                
                # Determine base text default - preserve custom descriptions
                # Priority: existing content > session state > calculated default
                if existing_table.get("custom_base_description"):
                    default_base = format_base_text(existing_table["custom_base_description"], new_base_n)
                elif existing_table.get("base_text"):
                    default_base = existing_table["base_text"]
                else:
                    default_base = st.session_state["selections"].get(tid, {}).get("base_text")
                    if default_base is None:
                        default_base = format_base_text("Total respondents", new_base_n)
                
                # Show indicator if using existing custom base description
                base_label = "Base text"
                if existing_table.get("custom_base_description"):
                    base_label += f" (Previously: {existing_table['custom_base_description']})"
                
                base_text_val = None
                if has_base_obj:
                    base_text_val = st.text_input(base_label, value=default_base, key=f"base_{tid}")
                
                # Question text - preserve custom questions
                existing_q = existing_table.get("custom_question", "")
                session_q = st.session_state["selections"].get(tid, {}).get("question_text", "")
                
                # Priority: existing content > session state > table title
                if existing_q and existing_q != t["title"]:
                    default_q = existing_q
                elif session_q:
                    default_q = session_q
                else:
                    default_q = t["title"]
                
                # Show indicator if using existing custom question
                question_label = "Question text"
                if existing_q and existing_q != t["title"]:
                    question_label += " (Previously: " + existing_q + ")"
                
                question_text_val = None
                if has_question_obj:
                    question_text_val = st.text_input(question_label, value=default_q, key=f"qtext_{tid}")

            # Row sorting management section
            st.write("**Row Sorting**")
            
            # Toggle checkbox for row sorting
            enable_sorting = st.checkbox(f"Enable row sorting for this table", value=False, key=f"enable_sorting_{tid}")
            
            if enable_sorting:
                st.info("💡 Rows will be sorted by their values (descending). You can exclude certain rows from sorting.")
                
                # Get available rows for exclusion
                available_rows = [row for row in t["row_labels"] if isinstance(row, str) and row.strip()]
                
                # Common rows that users might want to exclude from sorting
                suggested_excludes = []
                for row in available_rows:
                    row_lower = row.lower().strip()
                    if any(pattern in row_lower for pattern in ["other", "none", "don't know", "no answer", "prefer not", "n/a"]):
                        suggested_excludes.append(row)
                
                # Row exclusion selection
                exclude_options = ["None (sort all rows)"] + available_rows
                default_excludes = ["None (sort all rows)"]
                if suggested_excludes:
                    default_excludes = suggested_excludes
                
                excluded_rows = st.multiselect(
                    "Rows to exclude from sorting (will appear at bottom)",
                    exclude_options,
                    default=default_excludes,
                    key=f"sort_exclude_{tid}",
                    help="Select rows that should remain at the bottom and not be sorted by value"
                )
                
                # Remove "None" option if other selections are made
                if "None (sort all rows)" in excluded_rows and len(excluded_rows) > 1:
                    excluded_rows = [row for row in excluded_rows if row != "None (sort all rows)"]
                elif excluded_rows == ["None (sort all rows)"]:
                    excluded_rows = []
                
                # Show preview of what will be sorted vs excluded
                if excluded_rows:
                    sortable_rows = [row for row in available_rows if row not in excluded_rows]
                    st.write(f"**Will be sorted:** {len(sortable_rows)} rows")
                    st.write(f"**Will stay at bottom:** {len(excluded_rows)} rows ({', '.join(excluded_rows[:3])}{'...' if len(excluded_rows) > 3 else ''})")
                else:
                    st.write(f"**Will be sorted:** All {len(available_rows)} rows")
            else:
                excluded_rows = []
            
            st.divider()
            
            # Callout management section (only if connected in PPT)
            # Check if there are existing callouts from PowerPoint or session state
            existing_table = existing_content.get(t["title"], {})
            existing_callouts = existing_table.get("callouts", []) if existing_table else []
            has_callouts_obj = existing_table.get("has_callouts", False)
            current_callouts = st.session_state["selections"].get(tid, {}).get("callouts", [])
            total_callout_count = len(existing_callouts) + len(current_callouts)
            
            # Determine if callouts should be enabled by default
            # Enable if there are existing callouts or if user has previously enabled them
            default_enabled = (total_callout_count > 0 or 
                             st.session_state.get(f"enable_callouts_{tid}", False))
            
            enable_callouts = False
            if has_callouts_obj:
                st.write("**Callouts**")
                # Toggle checkbox for callouts
                toggle_label = f"Enable callouts for this table"
                if total_callout_count > 0:
                    toggle_label += f" ({total_callout_count} active)"
                enable_callouts = st.checkbox(toggle_label, value=default_enabled, key=f"enable_callouts_{tid}")
            
            if not enable_callouts:
                # Clear any existing callouts when disabled
                if "callouts" in st.session_state["selections"].get(tid, {}):
                    st.session_state["selections"][tid]["callouts"] = []
            
            if has_callouts_obj and enable_callouts:
                # Only show callout controls when enabled
                if total_callout_count == 0:
                    cc1, cc2, cc3 = st.columns([2, 1, 1])
                    
                    with cc1:
                        # Row selection for callouts
                        available_rows = [row for row in t["row_labels"] if isinstance(row, str) and row.strip()]
                        selected_row = st.selectbox("Select row for callout", available_rows, key=f"callout_row_{tid}")
                    
                    with cc2:
                        # Callout selectors: Banner only to users, Metric shown only if multiple
                        callout_metric = current_metric
                        if has_metrics:
                            callout_metric = st.selectbox("Callout metric", metric_options, key=f"callout_metric_{tid}", index=metric_options.index(current_metric) if current_metric in metric_options else 0)
                        callout_banner_default = selected_col if selected_col in banners else (banners[0] if banners else "")
                        callout_banner = st.selectbox("Callout banner", banners, key=f"callout_banner_{tid}", index=banners.index(callout_banner_default) if callout_banner_default in banners else 0)
                        # Metric type selector (display capitalized, store lowercase)
                        metric_type_display = ["Percentage", "Number", "Currency"]
                        default_metric_type = st.session_state["selections"].get(tid, {}).get("callout_metric_type", "percentage")
                        def_idx = metric_type_display.index(default_metric_type.capitalize()) if default_metric_type and default_metric_type.capitalize() in metric_type_display else 0
                        selected_metric_display = st.selectbox("Callout metric type", metric_type_display, key=f"callout_metric_type_{tid}", index=def_idx)
                        callout_metric_type = selected_metric_display.lower()
                        # Resolve to combined column label internally
                        callout_col = pair_to_combined.get(((callout_metric or ""), callout_banner))
                        if not callout_col:
                            # Fallback to first occurrence of the banner
                            for i, b in enumerate(banners):
                                if b == callout_banner:
                                    callout_col = combined_labels[i]
                                    break
                        if not callout_col and combined_labels:
                            callout_col = combined_labels[0]
                    
                    with cc3:
                        if st.button("Add Callout", key=f"add_callout_{tid}"):
                            if "callouts" not in st.session_state["selections"][tid]:
                                st.session_state["selections"][tid]["callouts"] = []
                            
                            # Create new callout with default text format
                            default_text = f"{selected_row}: [Value]"
                            new_callout = {
                                "row_label": selected_row,
                                "column_key": callout_col,
                                "text": default_text,  # Start with default "Row Name: Value" format
                                "position": (0.5, 7.0, 9.0, 0.4),
                                "font_size": 12,
                                "font_bold": True,
                                "metric_type": callout_metric_type
                            }
                            
                            st.session_state["selections"][tid]["callouts"].append(new_callout)
                            st.success(f"Added callout for '{selected_row}'")
                            st.rerun()  # Refresh to update the UI
                else:
                    st.info("📝 Callouts already exist for this table")
                
                # Display existing callouts with editable text boxes
                
                # Show existing callouts from PowerPoint with "Previously:" indicator
                if has_callouts_obj and (existing_table.get("callouts", []) ):
                    st.write("**Existing Callouts from PowerPoint:**")
                    ecs = existing_table.get("callouts", [])
                    for i, existing_callout in enumerate(ecs):
                        callout_display_cols = st.columns([4, 1, 1, 1])
                        
                        with callout_display_cols[0]:
                            # Show the callout info with "Previously:" indicator
                            st.write(f"• {existing_callout['row_label']} → {existing_callout['column_key']}")
                            
                            # Editable text box showing existing callout text with "Previously:" label
                            callout_label = "Callout text (Previously: " + existing_callout['text'] + ")"
                            updated_text = st.text_input(
                                callout_label,
                                value=existing_callout['text'],
                                key=f"existing_callout_text_{tid}_{i}",
                                help="Customize your callout text. Use [Value] as a placeholder for the actual data value."
                            )
                            # Metric type dropdown for existing callout (display capitalized, store lowercase)
                            mt_display = ["Percentage", "Number", "Currency"]
                            mt_current = existing_callout.get('metric_type', 'percentage')
                            mt_idx = mt_display.index(mt_current.capitalize()) if mt_current and mt_current.capitalize() in mt_display else 0
                            mt_updated_display = st.selectbox("Metric type", mt_display, index=mt_idx, key=f"existing_callout_metric_{tid}_{i}")
                            mt_updated = mt_updated_display.lower()
                            
                            # Update the existing callout text in real-time
                            if updated_text != existing_callout['text']:
                                ecs[i]['text'] = updated_text
                            if mt_updated != existing_callout.get('metric_type'):
                                ecs[i]['metric_type'] = mt_updated
                        
                        with callout_display_cols[1]:
                            if st.button("Remove", key=f"remove_existing_callout_{tid}_{i}"):
                                ecs.pop(i)
                                st.rerun()
                        
                        with callout_display_cols[2]:
                            if st.button("Edit", key=f"edit_existing_callout_{tid}_{i}"):
                                st.session_state["editing_existing_callout"] = (tid, i)
                                st.rerun()
                        
                        # Add some spacing between callouts
                        st.divider()
                
                # Show current callouts from session state
                if current_callouts:
                    st.write("**Current Callouts:**")
                    for i, callout in enumerate(current_callouts):
                        callout_display_cols = st.columns([4, 1, 1, 1])
                        
                        with callout_display_cols[0]:
                            # Show the callout info and editable text box
                            st.write(f"• {callout['row_label']} → {callout['column_key']}")
                            
                            # Editable text box for customizing callout text
                            current_text = callout.get('text', f"{callout['row_label']}: [Value]")
                            updated_text = st.text_input(
                                "Callout text (customize as needed):",
                                value=current_text,
                                key=f"callout_text_{tid}_{i}",
                                help="Customize your callout text. Use [Value] as a placeholder for the actual data value."
                            )
                            mt_display = ["Percentage", "Number", "Currency"]
                            mt_current = callout.get('metric_type', 'percentage')
                            mt_idx = mt_display.index(mt_current.capitalize()) if mt_current and mt_current.capitalize() in mt_display else 0
                            mt_updated_display = st.selectbox("Metric type", mt_display, index=mt_idx, key=f"callout_metric_{tid}_{i}")
                            mt_updated = mt_updated_display.lower()
                            
                            # Update the callout text in real-time
                            if updated_text != current_text:
                                st.session_state["selections"][tid]["callouts"][i]["text"] = updated_text
                            if mt_updated != mt_current:
                                st.session_state["selections"][tid]["callouts"][i]["metric_type"] = mt_updated
                        
                        with callout_display_cols[1]:
                            if st.button("Remove", key=f"remove_callout_{tid}_{i}"):
                                st.session_state["selections"][tid]["callouts"].pop(i)
                                st.rerun()
                        
                        with callout_display_cols[2]:
                            if st.button("Edit", key=f"edit_callout_{tid}_{i}"):
                                st.session_state["editing_callout"] = (tid, i)
                                st.rerun()
                        
                        # Add some spacing between callouts
                        st.divider()

            # Build selections dictionary
            chart_label_to_key = {
                "Horizontal Bar":    "bar_h",
                "Vertical Bar":      "bar_v",
                "Donut":             "donut",
                "Line":              "line",
                "Grouped Bar (2)":   "grouped_bar_2",
                "Grouped Bar (3)":   "grouped_bar_3",
                "Multi-Line":        "multi_line",
                "Table Only":        "table_only",
                "Chart + Table":     "chart_table",
            }
            selection_dict = {
                "chart_type_label": choice,
                "chart_type": chart_label_to_key.get(choice, "bar_h"),
                # Persist both banner and metric plus the resolved combined label for export
                "banner_key": selected_col,
                "metric_key": current_metric,
                "column_key": combined_selected,
                "column_keys": multi_columns_selected if multi_columns_selected else None,
                "enable_sorting": enable_sorting,
                "excluded_rows": excluded_rows if enable_sorting else []
            }

            # Only include text fields if those objects are linked in PPT
            if title_val is not None:
                selection_dict["title"] = title_val
            if base_text_val is not None:
                selection_dict["base_text"] = base_text_val
            if question_text_val is not None:
                selection_dict["question_text"] = question_text_val
            
            # Include only user-created callouts in selections; do NOT merge in existing PPT callouts
            if enable_callouts:
                user_callouts = st.session_state["selections"].get(tid, {}).get("callouts", [])
                if user_callouts:
                    selection_dict["callouts"] = user_callouts
            
            st.session_state["selections"][tid] = selection_dict

    # Export/Update section
    if st.session_state.data is not None:
        st.divider()
        
        if workflow_type == "Create New Report":
            # NEW REPORT EXPORT
            st.subheader("Step 4: Export new PowerPoint")
            if st.button("Export PowerPoint", type="primary"):
                sels = {tid: {
                    "chart_type": v.get("chart_type", "bar_h"), 
                    "column_key": v.get("column_key", "Total"),
                    "column_keys": v.get("column_keys"),
                    "title": v.get("title"), 
                    "base_text": v.get("base_text"), 
                    "question_text": v.get("question_text"), 
                    "callouts": v.get("callouts", []),
                    "enable_sorting": v.get("enable_sorting", False),
                    "excluded_rows": v.get("excluded_rows", [])
                } for tid, v in st.session_state["selections"].items()}
                out = _save_temp(b"", ".pptx", "_tmp_out_pptx")
                ai_data = None
                try:
                    with st.spinner("Generating AI insights..."):
                        ai_data = generate_all_insights(data["tables"], selections=sels)
                except Exception as e:
                    st.warning(f"AI insights skipped: {e}")
                export_pptx(data["tables"], sels, out, ai_insights=ai_data)
                with open(out, "rb") as f:
                    st.download_button("Download report.pptx", f, file_name="report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        
        else:
            # UPDATE EXISTING REPORT
            st.subheader("Step 6: Update PowerPoint")
            
            # Show what will happen
            connected_count = len([t for t in data["tables"] if t["title"] in existing_content])
            unconnected_count = len([t for t in data["tables"] if t["title"] not in existing_content])
            
            st.info(f"""
            **Update Summary:**
            - {connected_count} existing tables will be updated with new data
            - {unconnected_count} new tables will be added to an "Unmapped Tables" summary page
            - Custom formatting and content will be preserved where possible
            """)
            
            # --- Match review ---
            if "match_review" not in st.session_state:
                st.session_state.match_review = None
            if "match_overrides" not in st.session_state:
                st.session_state.match_overrides = {}

            if st.button("Preview Matches", key="preview_matches"):
                prs_preview = Presentation(st.session_state["_tmp_pptx"])
                matcher = SmartMatcher(data["tables"])
                shapes_meta = []
                for slide in prs_preview.slides:
                    for shp in slide.shapes:
                        alt = _parse_alt_text(shp)
                        if alt.get("auto_update", "yes").lower() == "no":
                            continue
                        is_chart = False
                        try:
                            _ = shp.chart
                            is_chart = True
                        except (ValueError, AttributeError):
                            pass
                        if is_chart or shp.has_table:
                            stype = "Chart" if is_chart else "Table"
                            shapes_meta.append({"name": shp.name or "", "alt": alt, "shape_type": stype})
                matcher.match_all(shapes_meta)
                st.session_state.match_review = matcher.get_report()
                st.session_state.match_overrides = {}

            if st.session_state.match_review:
                report = st.session_state.match_review
                table_titles = [t["title"] for t in data["tables"]]

                # Build title → resolved column_key from current selections
                _title_to_col = {}
                for _tid, _sel in st.session_state["selections"].items():
                    for _t in data["tables"]:
                        if _t["id"] == _tid:
                            _title_to_col[_t["title"]] = (
                                _sel.get("column_key") or _sel.get("banner_key") or "Total"
                            )
                            break

                with st.expander("Match Review", expanded=True):
                    # Approve All / Skip unmatched helper
                    approve_col, skip_col, _ = st.columns([1, 1, 4])
                    with approve_col:
                        if st.button("Approve All", key="approve_all_matches"):
                            st.session_state.match_overrides = {}
                            for entry in report:
                                label = entry["shape_alt_title"] or entry["shape_name"]
                                if entry["status"] in ("failed", "duplicate", "skipped"):
                                    st.session_state.match_overrides[label] = "__skip__"
                            st.rerun()
                    with skip_col:
                        if st.button("Skip All Unmatched", key="skip_unmatched"):
                            for entry in report:
                                if entry["status"] in ("failed", "duplicate"):
                                    label = entry["shape_alt_title"] or entry["shape_name"]
                                    st.session_state.match_overrides[label] = "__skip__"
                            st.rerun()

                    # Header row
                    hdr = st.columns([1, 3, 2, 1, 1, 1, 3])
                    hdr[0].markdown("**Type**")
                    hdr[1].markdown("**Old Title**")
                    hdr[2].markdown("**Matched Title**")
                    hdr[3].markdown("**Conf.**")
                    hdr[4].markdown("**Tier**")
                    hdr[5].markdown("**Column**")
                    hdr[6].markdown("**Action**")

                    _TIER_LABELS = {0: "—", 1: "Exact", 2: "Fuzzy", 3: "LLM"}

                    for i, entry in enumerate(report):
                        tier = entry["tier"]
                        conf = entry["confidence"]
                        status = entry["status"]
                        shape_label = entry["shape_alt_title"] or entry["shape_name"] or f"Shape {i}"
                        matched = entry["matched_table"] or "—"
                        stype = entry.get("shape_type", "?")
                        resolved_col = _title_to_col.get(matched, "—") if matched != "—" else "—"

                        if status in ("failed", "duplicate", "skipped"):
                            row_color = "#FADADD"
                        elif status == "low_confidence" or 0.60 <= conf < 0.85:
                            row_color = "#FFF3CD"
                        elif conf >= 0.85:
                            row_color = "#D4EDDA"
                        else:
                            row_color = "#FADADD"

                        st.markdown(
                            f'<div style="background:{row_color};padding:2px 6px;border-radius:4px;margin-bottom:2px;">&nbsp;</div>',
                            unsafe_allow_html=True,
                        )
                        row = st.columns([1, 3, 2, 1, 1, 1, 3])
                        row[0].write(stype)
                        row[1].write(shape_label)
                        row[2].write(matched)
                        row[3].write(f"{conf:.0%}" if conf else "—")
                        row[4].write(_TIER_LABELS.get(tier, str(tier)))
                        row[5].write(resolved_col)

                        action_options = [f"Update ({matched})", "Skip"] + table_titles
                        current_override = st.session_state.match_overrides.get(shape_label)
                        if current_override == "__skip__":
                            default_idx = 1
                        elif current_override and current_override in table_titles:
                            default_idx = 2 + table_titles.index(current_override)
                        else:
                            default_idx = 0

                        action = row[6].selectbox(
                            "Action",
                            action_options,
                            index=default_idx,
                            key=f"match_action_{i}",
                            label_visibility="collapsed",
                        )
                        if action == "Skip":
                            st.session_state.match_overrides[shape_label] = "__skip__"
                        elif action.startswith("Update ("):
                            st.session_state.match_overrides.pop(shape_label, None)
                        else:
                            st.session_state.match_overrides[shape_label] = action

            # --- Data diff preview for connected tables ---
            if st.session_state.match_review:
                connected_with_data = [
                    title for title, info in existing_content.items()
                    if info.get("table_values")
                ]
                if connected_with_data:
                    with st.expander("Data Changes Preview", expanded=False):
                        for title in connected_with_data:
                            old_tv = existing_content[title]["table_values"]
                            new_table = next(
                                (t for t in data["tables"] if t["title"] == title), None
                            )
                            if not new_table:
                                continue
                            st.markdown(f"**{title}**")
                            _render_diff_table(old_tv, new_table)
                            st.divider()

            if st.button("Update PowerPoint", type="primary"):
                # Convert selections to use table titles as keys with all fields
                table_selections = {}
                for tid, v in st.session_state["selections"].items():
                    for t in data["tables"]:
                        if t["id"] == tid:
                            table_selections[t["title"]] = {
                                "chart_type": v.get("chart_type", "bar_h"),
                                "column_key": v.get("column_key", "Total"),
                                "column_keys": v.get("column_keys"),
                                "title": v.get("title"),
                                "base_text": v.get("base_text"),
                                "question_text": v.get("question_text"),
                                "callouts": v.get("callouts", []),
                                "enable_sorting": v.get("enable_sorting", False),
                                "excluded_rows": v.get("excluded_rows", []),
                            }
                            break

                # Build SmartMatcher with user overrides from match review
                overrides = st.session_state.get("match_overrides", {})
                matcher = SmartMatcher(
                    data["tables"],
                    overrides=overrides if overrides else None,
                )

                progress_bar = st.progress(0, text="Updating presentation...")

                def _on_progress(pct: float):
                    progress_bar.progress(min(pct, 1.0), text=f"Updating... {pct:.0%}")

                out_path = _save_temp(b"", ".pptx", "_tmp_updated_pptx")
                updated = update_presentation_with_unmapped(
                    st.session_state["_tmp_pptx"],
                    st.session_state["_tmp_xlsx"],
                    out_path,
                    table_selections, data["tables"], existing_content,
                    matcher=matcher,
                    progress_callback=_on_progress,
                )

                progress_bar.progress(1.0, text="Update complete!")

                with open(updated, "rb") as f:
                    st.download_button("Download updated_report.pptx", f, file_name="updated_report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# Show instructions when no data is loaded
if st.session_state.data is None:
    if workflow_type == "Create New Report":
        st.info("👆 Upload a crosstab Excel file to begin creating your report.")
