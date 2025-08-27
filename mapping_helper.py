"""
PowerPoint Mapping Helper

This script helps you manually create mappings between PowerPoint shapes and crosstab data.
It can be used to:
1. List all shapes in a PowerPoint file
2. Show current mapping status
3. Generate mapping templates for manual editing
4. Validate existing mappings
"""

from pptx import Presentation
from crosstab_parser import parse_workbook
import json
from typing import Dict, Any, List

def list_all_shapes(pptx_path: str) -> List[Dict[str, Any]]:
    """List all shapes in a PowerPoint file with their current mapping status."""
    prs = Presentation(pptx_path)
    shapes_info = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape_num, shape in enumerate(slide.shapes, 1):
            # Check if shape has chart (GraphicFrame with chart)
            has_chart = False
            try:
                if hasattr(shape, "chart") and shape.chart is not None:
                    has_chart = True
            except:
                pass
            
            # Check if shape has table
            has_table = False
            try:
                if hasattr(shape, "has_table") and shape.has_table:
                    has_table = True
            except:
                pass
            
            shape_info = {
                "slide": slide_num,
                "shape_num": shape_num,
                "name": shape.name or f"Unnamed_{slide_num}_{shape_num}",
                "type": "unknown",
                "has_chart": has_chart,
                "has_table": has_table,
                "alt_text": "",
                "mapping_status": "unmapped"
            }
            
            # Determine shape type
            if shape_info["has_chart"]:
                shape_info["type"] = "chart"
            elif shape_info["has_table"]:
                shape_info["type"] = "table"
            elif hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                shape_info["type"] = "text"
            else:
                shape_info["type"] = "other"
            
            # Get alt text
            try:
                # Method 1: Try to read from XML descr attribute (most reliable)
                alt_text = ""
                if hasattr(shape, 'element'):
                    # Look for the cNvPr element which contains the description
                    c_nv_pr = None
                    
                    # For GraphicFrame (charts/tables)
                    if 'graphicFrame' in shape.element.tag:
                        c_nv_pr = shape.element.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    # For Shape (text boxes, etc.)
                    elif 'sp' in shape.element.tag:
                        c_nv_pr = shape.element.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    
                    if c_nv_pr is not None and c_nv_pr.get('descr'):
                        alt_text = c_nv_pr.get('descr')
                
                # Method 2: Fallback to alternative_text property (if it exists)
                if not alt_text:
                    try:
                        alt_text = shape.alternative_text or ""
                    except:
                        alt_text = ""
                
                shape_info["alt_text"] = alt_text
            except:
                shape_info["alt_text"] = ""
            
            # Determine mapping status
            if shape_info["alt_text"]:
                if "table_title:" in shape_info["alt_text"] or "type:" in shape_info["alt_text"]:
                    shape_info["mapping_status"] = "mapped"
                elif shape_info["name"].startswith(("CHART_", "TABLE_", "TEXT_")):
                    shape_info["mapping_status"] = "named_but_unmapped"
            
            shapes_info.append(shape_info)
    
    return shapes_info

def generate_mapping_template(pptx_path: str, crosstab_path: str = None) -> str:
    """Generate a mapping template that can be manually edited."""
    shapes_info = list_all_shapes(pptx_path)
    
    template = "# PowerPoint Mapping Template\n\n"
    template += "# Instructions:\n"
    template += "# 1. Edit the mapping values below\n"
    template += "# 2. Save this file\n"
    template += "# 3. Use the apply_mapping function to update your PowerPoint\n\n"
    
    if crosstab_path:
        try:
            data = parse_workbook(crosstab_path)
            template += f"# Available crosstab tables:\n"
            for table in data["tables"]:
                template += f"# - {table['title']} (columns: {', '.join(table['col_labels'])})\n"
            template += "\n"
        except Exception as e:
            template += f"# Could not parse crosstab: {e}\n\n"
    
    template += "MAPPINGS = {\n"
    
    for shape in shapes_info:
        if shape["type"] in ["chart", "table"]:
            template += f"    # Slide {shape['slide']}, Shape {shape['shape_num']}: {shape['name']}\n"
            template += f"    '{shape['name']}': {{\n"
            template += f"        'type': '{shape['type']}',\n"
            template += f"        'table_title': 'REPLACE_WITH_TABLE_TITLE',\n"
            if shape["type"] == "chart":
                template += f"        'column': 'Total',  # or specific column name\n"
            template += f"        'exclude_rows': 'base, mean, average, avg',\n"
            template += f"        'auto_update': 'yes'\n"
            template += f"    }},\n\n"
    
    template += "}\n"
    
    return template

def apply_mapping_from_file(pptx_path: str, mapping_file_path: str, output_path: str = None) -> str:
    """Apply mappings from a Python file to a PowerPoint presentation."""
    if output_path is None:
        output_path = pptx_path.replace(".pptx", "_mapped.pptx")
    
    # Load the mapping file
    with open(mapping_file_path, 'r') as f:
        exec(f.read(), globals())
    
    if 'MAPPINGS' not in globals():
        raise ValueError("No MAPPINGS dictionary found in the file")
    
    # Apply mappings to PowerPoint
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name and shape.name in MAPPINGS:
                mapping = MAPPINGS[shape.name]
                
                # Update alt text
                alt_lines = []
                for key, value in mapping.items():
                    if value is not None and value != "":
                        alt_lines.append(f"{key}: {value}")
                
                if alt_lines:
                    try:
                        shape.alternative_text = "\n".join(alt_lines)
                    except Exception as e:
                        print(f"Warning: Could not set alt text for {shape.name}: {e}")
    
    prs.save(output_path)
    return output_path

def validate_mappings(pptx_path: str, crosstab_path: str) -> Dict[str, Any]:
    """Validate existing mappings against crosstab data."""
    shapes_info = list_all_shapes(pptx_path)
    data = parse_workbook(crosstab_path)
    
    validation_results = {
        "total_shapes": len(shapes_info),
        "mapped_shapes": 0,
        "valid_mappings": 0,
        "invalid_mappings": 0,
        "issues": []
    }
    
    for shape in shapes_info:
        if shape["mapping_status"] == "mapped":
            validation_results["mapped_shapes"] += 1
            
            # Parse alt text to get mapping
            alt_text = shape["alt_text"]
            mapping = {}
            for line in alt_text.splitlines():
                if ":" in line and not line.startswith("---"):
                    key, value = line.split(":", 1)
                    mapping[key.strip()] = value.strip()
            
            # Validate mapping
            if "table_title" in mapping:
                table_title = mapping["table_title"]
                table_found = False
                
                for table in data["tables"]:
                    if table.get("title") == table_title:
                        table_found = True
                        
                        # Validate column if specified
                        if "column" in mapping and shape["type"] == "chart":
                            column = mapping["column"]
                            if column not in table.get("col_labels", []):
                                validation_results["issues"].append(
                                    f"Shape '{shape['name']}': Column '{column}' not found in table '{table_title}'"
                                )
                                validation_results["invalid_mappings"] += 1
                                break
                        else:
                            validation_results["valid_mappings"] += 1
                        break
                
                if not table_found:
                    validation_results["issues"].append(
                        f"Shape '{shape['name']}': Table '{table_title}' not found in crosstab"
                    )
                    validation_results["invalid_mappings"] += 1
    
    return validation_results

def main():
    """Interactive command-line interface for the mapping helper."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python mapping_helper.py <command> [options]")
        print("\nCommands:")
        print("  list <pptx_file>                    - List all shapes in PowerPoint")
        print("  template <pptx_file> [crosstab]     - Generate mapping template")
        print("  apply <pptx_file> <mapping_file>    - Apply mappings from file")
        print("  validate <pptx_file> <crosstab>     - Validate existing mappings")
        return
    
    command = sys.argv[1]
    
    if command == "list" and len(sys.argv) >= 3:
        pptx_path = sys.argv[2]
        shapes = list_all_shapes(pptx_path)
        
        print(f"\nShapes in {pptx_path}:")
        print("-" * 80)
        for shape in shapes:
            status_icon = "✅" if shape["mapping_status"] == "mapped" else "❌"
            print(f"{status_icon} Slide {shape['slide']:2d} | {shape['type']:8s} | {shape['name']:30s} | {shape['mapping_status']}")
    
    elif command == "template" and len(sys.argv) >= 3:
        pptx_path = sys.argv[2]
        crosstab_path = sys.argv[3] if len(sys.argv) >= 4 else None
        
        template = generate_mapping_template(pptx_path, crosstab_path)
        output_file = "mapping_template.py"
        
        with open(output_file, 'w') as f:
            f.write(template)
        
        print(f"Mapping template saved to {output_file}")
        print("Edit this file with your mappings, then use 'apply' command to update PowerPoint")
    
    elif command == "apply" and len(sys.argv) >= 4:
        pptx_path = sys.argv[2]
        mapping_file = sys.argv[3]
        
        try:
            output_path = apply_mapping_from_file(pptx_path, mapping_file)
            print(f"Mappings applied successfully! Updated file saved as: {output_path}")
        except Exception as e:
            print(f"Error applying mappings: {e}")
    
    elif command == "validate" and len(sys.argv) >= 4:
        pptx_path = sys.argv[2]
        crosstab_path = sys.argv[3]
        
        try:
            results = validate_mappings(pptx_path, crosstab_path)
            print(f"\nValidation Results for {pptx_path}:")
            print("-" * 50)
            print(f"Total shapes: {results['total_shapes']}")
            print(f"Mapped shapes: {results['mapped_shapes']}")
            print(f"Valid mappings: {results['valid_mappings']}")
            print(f"Invalid mappings: {results['invalid_mappings']}")
            
            if results['issues']:
                print(f"\nIssues found:")
                for issue in results['issues']:
                    print(f"  - {issue}")
        except Exception as e:
            print(f"Error validating mappings: {e}")
    
    else:
        print("Invalid command or missing arguments. Use 'python mapping_helper.py' for help.")

if __name__ == "__main__":
    main()
