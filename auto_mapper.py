"""
auto_mapper.py — AI-powered bootstrap pass that maps un-annotated PowerPoint
shapes to crosstab tables by extracting structural fingerprints from existing
chart/table data and matching them against parsed crosstab tables.

v2: Multi-signal matching engine with value correlation, slide context,
series-to-row fuzzy matching, and time-series support.

Produces a copy of the PPTX with alt text written onto matched shapes so that
the normal deck_update pipeline can process it without manual pre-mapping.
"""

import json
import logging
import math
import re
from dataclasses import dataclass, field
from difflib import SequenceMatcher
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from crosstab_parser import parse_workbook

logger = logging.getLogger("report_relay.auto_mapper")

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# ---------------------------------------------------------------------------
# Blended scoring weights (configurable)
# ---------------------------------------------------------------------------

WEIGHT_VALUE_CORR = 0.45
WEIGHT_LABEL_OVERLAP = 0.25
WEIGHT_SLIDE_CONTEXT = 0.20
WEIGHT_SERIES_NAME = 0.10

# Confidence thresholds
HIGH_CONFIDENCE_THRESHOLD = 0.75
LLM_THRESHOLD = 0.40
ALT_TEXT_WRITE_THRESHOLD = 0.75


# ---------------------------------------------------------------------------
# Normalisation helpers
# ---------------------------------------------------------------------------

def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "")).strip().lower()


def _jaccard(a: set, b: set) -> float:
    if not a and not b:
        return 1.0
    union = a | b
    if not union:
        return 0.0
    return len(a & b) / len(union)


# ---------------------------------------------------------------------------
# Detection helpers
# ---------------------------------------------------------------------------

_QCODE_RE = re.compile(r"\(Q\d+[_\d]*\)")

_DATE_PATTERNS = [
    re.compile(
        r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d{4}",
        re.IGNORECASE,
    ),
    re.compile(
        r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d{1,2}",
        re.IGNORECASE,
    ),
    re.compile(
        r"\d{1,2}[/-]\d{1,2}[/-]\d{2,4}",
    ),
]

# Month name → month number for date normalisation
_MONTH_MAP = {
    "jan": 1, "feb": 2, "mar": 3, "apr": 4, "may": 5, "jun": 6,
    "jul": 7, "aug": 8, "sep": 9, "oct": 10, "nov": 11, "dec": 12,
    "january": 1, "february": 2, "march": 3, "april": 4, "june": 6,
    "july": 7, "august": 8, "september": 9, "october": 10,
    "november": 11, "december": 12,
}

_UNINFORMATIVE_SERIES = frozenset({"%", "series 1", "series1", ""})

_SEGMENT_KEYWORDS = frozenset({
    "total", "gen z", "millennial", "gen x", "boomer", "baby boomer",
    "west", "midwest", "northeast", "south", "male", "female",
    "18-24", "25-34", "35-44", "45-54", "55-64", "65+",
})


def _extract_qcode(text: str) -> Optional[str]:
    """Extract a question code like ``(Q1369)`` from arbitrary text."""
    m = _QCODE_RE.search(text or "")
    return m.group() if m else None


def _is_uninformative_series(names: List[str]) -> bool:
    if not names:
        return True
    return all(_norm(n) in _UNINFORMATIVE_SERIES for n in names)


def _is_trend_chart(categories: List[str]) -> bool:
    """True when most categories look like dates (historical trend data)."""
    if len(categories) < 3:
        return False
    hits = 0
    for cat in categories:
        if any(p.search(cat) for p in _DATE_PATTERNS):
            hits += 1
    return hits / len(categories) >= 0.5


def _is_segment_chart(categories: List[str], col_banners_pool: set) -> bool:
    """True when chart categories are demographic/segment columns, not response options.

    Uses a two-pronged check:
    1. High overlap with actual column banners from the crosstab
    2. Fallback to common segment keywords
    """
    if not categories:
        return False
    norm_cats = {_norm(c) for c in categories}
    # Check against actual crosstab column banners
    if col_banners_pool:
        overlap = len(norm_cats & col_banners_pool)
        if overlap / len(norm_cats) >= 0.5:
            return True
    # Fallback: keyword check
    keyword_hits = sum(1 for c in norm_cats if any(kw in c for kw in _SEGMENT_KEYWORDS))
    return keyword_hits / len(norm_cats) >= 0.5


# ---------------------------------------------------------------------------
# Value correlation helpers
# ---------------------------------------------------------------------------

def _pearson_r(x: List[float], y: List[float]) -> float:
    """Compute Pearson correlation coefficient between two equal-length vectors.

    Returns 0.0 if either vector has zero variance or lengths differ.
    """
    n = min(len(x), len(y))
    if n < 3:
        return 0.0

    x, y = x[:n], y[:n]
    mx = sum(x) / n
    my = sum(y) / n

    sx = sum((xi - mx) ** 2 for xi in x)
    sy = sum((yi - my) ** 2 for yi in y)

    if sx == 0 or sy == 0:
        return 0.0

    sxy = sum((xi - mx) * (yi - my) for xi, yi in zip(x, y))
    return sxy / (sx * sy) ** 0.5


def _fuzzy_match_label(
    label: str,
    candidates: List[str],
    threshold: float = 0.6,
) -> Optional[str]:
    """Find the best fuzzy match for *label* among *candidates*."""
    if not label or not candidates:
        return None
    norm_label = _norm(label)
    best_score = 0.0
    best_match: Optional[str] = None
    for cand in candidates:
        norm_cand = _norm(cand)
        # Exact substring match gets a boost
        if norm_label in norm_cand or norm_cand in norm_label:
            ratio = max(0.85, SequenceMatcher(None, norm_label, norm_cand).ratio())
        else:
            ratio = SequenceMatcher(None, norm_label, norm_cand).ratio()
        if ratio > best_score:
            best_score = ratio
            best_match = cand
    if best_score >= threshold:
        return best_match
    return None


def _build_segment_mapping(
    chart_categories: List[str],
    crosstab_banners: List[str],
    threshold: float = 0.55,
) -> Dict[int, int]:
    """Map chart category indices to crosstab banner indices via fuzzy matching.

    Returns {chart_idx: banner_idx} for matched pairs.
    """
    mapping: Dict[int, int] = {}
    used_banner_idxs: set = set()

    # First pass: exact (normalised) matches
    norm_banners = [_norm(b) for b in crosstab_banners]
    for ci, cat in enumerate(chart_categories):
        nc = _norm(cat)
        for bi, nb in enumerate(norm_banners):
            if bi in used_banner_idxs:
                continue
            if nc == nb:
                mapping[ci] = bi
                used_banner_idxs.add(bi)
                break

    # Second pass: fuzzy matches for unmatched categories
    _stop_words = {"the", "a", "an", "of", "in", "for", "and", "or", "to", "is", "it"}
    for ci, cat in enumerate(chart_categories):
        if ci in mapping:
            continue
        nc = _norm(cat).rstrip("+").strip()
        nc_words = set(nc.split()) - _stop_words
        best_score = 0.0
        best_bi = None
        for bi, nb in enumerate(norm_banners):
            if bi in used_banner_idxs:
                continue
            nb_clean = nb.rstrip("+").strip()
            nb_words = set(nb_clean.split()) - _stop_words
            # Substring containment
            if nc in nb_clean or nb_clean in nc:
                score = max(0.85, SequenceMatcher(None, nc, nb_clean).ratio())
            # Word overlap: if a significant keyword is shared (e.g., "boomer")
            elif nc_words and nb_words and len(nc_words & nb_words) / max(len(nc_words), len(nb_words)) >= 0.3:
                score = max(0.70, SequenceMatcher(None, nc, nb_clean).ratio())
            else:
                score = SequenceMatcher(None, nc, nb_clean).ratio()
            if score > best_score:
                best_score = score
                best_bi = bi
        if best_bi is not None and best_score >= threshold:
            mapping[ci] = best_bi
            used_banner_idxs.add(best_bi)

    return mapping


def _value_correlation_distribution(
    fp_values: List[float],
    fp_categories: List[str],
    table: Dict[str, Any],
) -> float:
    """Correlate chart values against a table's Total column (distribution chart).

    Aligns chart categories to table row labels by fuzzy matching, then
    computes Pearson r on the aligned value pairs.
    """
    row_labels = table.get("row_labels", [])
    col_labels = table.get("col_labels", [])
    values = table.get("values", [])

    if not row_labels or not values:
        return 0.0

    # Find Total column
    total_idx = None
    for preferred in ["Total", "Overall", "All"]:
        if preferred in col_labels:
            total_idx = col_labels.index(preferred)
            break
    if total_idx is None:
        total_idx = 0 if col_labels else None
    if total_idx is None:
        return 0.0

    # Align chart categories to table row labels
    chart_vals_aligned: List[float] = []
    table_vals_aligned: List[float] = []

    norm_rows = {_norm(rl): i for i, rl in enumerate(row_labels)}

    used_row_indices: set = set()
    for ci, cat in enumerate(fp_categories):
        if ci >= len(fp_values) or fp_values[ci] is None:
            continue
        nc = _norm(cat)
        # Try exact match first
        ri = norm_rows.get(nc)
        # Try prefix/startswith match (handles "UNCHANGED - Neither..." → "UNCHANGED")
        if ri is None:
            for nrl, idx in norm_rows.items():
                if idx in used_row_indices:
                    continue
                if nc.startswith(nrl) or nrl.startswith(nc):
                    ri = idx
                    break
        # Fall back to fuzzy only if no prefix match
        if ri is None:
            best_ratio = 0.0
            best_ri = None
            for nrl, idx in norm_rows.items():
                if idx in used_row_indices:
                    continue
                ratio = SequenceMatcher(None, nc, nrl).ratio()
                if ratio > best_ratio:
                    best_ratio = ratio
                    best_ri = idx
            if best_ratio >= 0.65:
                ri = best_ri
        if ri is not None and ri not in used_row_indices and ri < len(values) and total_idx < len(values[ri]):
            used_row_indices.add(ri)
            tv = values[ri][total_idx]
            if tv is not None:
                try:
                    chart_vals_aligned.append(float(fp_values[ci]))
                    table_vals_aligned.append(float(tv))
                except (ValueError, TypeError):
                    pass

    if len(chart_vals_aligned) < 3:
        return 0.0

    r = _pearson_r(chart_vals_aligned, table_vals_aligned)
    return max(0.0, r)


def _value_correlation_segment(
    fp_values: List[float],
    fp_categories: List[str],
    table: Dict[str, Any],
) -> Tuple[float, Optional[str]]:
    """Correlate chart values against each row's values across segment columns.

    For segment breakout charts where categories are segments (Total, Gen Z, etc.)
    and the chart shows one metric (row) across those segments.

    Returns (best_correlation, matched_row_label).
    """
    row_labels = table.get("row_labels", [])
    col_labels = table.get("col_labels", [])
    banners = table.get("meta", {}).get("col_banners", col_labels)
    values = table.get("values", [])

    if not row_labels or not values or not banners:
        return 0.0, None

    # Build chart-category → banner-index mapping
    seg_map = _build_segment_mapping(fp_categories, banners)
    if len(seg_map) < 3:
        return 0.0, None

    best_r = 0.0
    best_row = None

    for ri, rl in enumerate(row_labels):
        # Skip base/mean rows
        nrl = _norm(rl)
        if nrl.startswith(("base", "mean", "average", "avg", "median")):
            continue
        if ri >= len(values):
            continue

        chart_aligned: List[float] = []
        table_aligned: List[float] = []

        for ci, bi in sorted(seg_map.items()):
            if ci >= len(fp_values) or fp_values[ci] is None:
                continue
            if bi >= len(values[ri]):
                continue
            tv = values[ri][bi]
            if tv is not None:
                try:
                    chart_aligned.append(float(fp_values[ci]))
                    table_aligned.append(float(tv))
                except (ValueError, TypeError):
                    pass

        if len(chart_aligned) < 3:
            continue

        r = _pearson_r(chart_aligned, table_aligned)
        if r > best_r:
            best_r = r
            best_row = rl

    return max(0.0, best_r), best_row


def _value_correlation_timeseries(
    fp_values: List[float],
    fp_categories: List[str],
    table: Dict[str, Any],
) -> Tuple[float, Optional[str]]:
    """Correlate a time-series chart's values against time-series table rows.

    Time-series charts have date categories and the crosstab time-series tables
    have date column headers (e.g., "January 2022", "February 2022", ...).

    Returns (best_correlation, matched_row_label).
    """
    row_labels = table.get("row_labels", [])
    col_labels = table.get("col_labels", [])
    values = table.get("values", [])

    if not row_labels or not values or not col_labels:
        return 0.0, None

    # Normalise chart date categories to "mon yyyy" format
    chart_date_keys = [_normalise_date_label(c) for c in fp_categories]
    # Normalise crosstab column labels to the same format
    col_date_keys = [_normalise_date_label(c) for c in col_labels]

    # Build col_date_key → col_idx mapping
    col_key_to_idx: Dict[str, int] = {}
    for ci, dk in enumerate(col_date_keys):
        if dk and dk not in col_key_to_idx:
            col_key_to_idx[dk] = ci

    best_r = 0.0
    best_row = None

    for ri, rl in enumerate(row_labels):
        nrl = _norm(rl)
        if nrl.startswith(("base", "mean", "average", "avg", "median")):
            continue
        if ri >= len(values):
            continue

        chart_aligned: List[float] = []
        table_aligned: List[float] = []

        for ci, dk in enumerate(chart_date_keys):
            if not dk or ci >= len(fp_values) or fp_values[ci] is None:
                continue
            col_idx = col_key_to_idx.get(dk)
            if col_idx is None:
                continue
            if col_idx >= len(values[ri]):
                continue
            tv = values[ri][col_idx]
            if tv is not None:
                try:
                    chart_aligned.append(float(fp_values[ci]))
                    table_aligned.append(float(tv))
                except (ValueError, TypeError):
                    pass

        if len(chart_aligned) < 3:
            continue

        r = _pearson_r(chart_aligned, table_aligned)
        if r > best_r:
            best_r = r
            best_row = rl

    return max(0.0, best_r), best_row


def _normalise_date_label(label: str) -> Optional[str]:
    """Normalise a date string to 'mon yyyy' for comparison.

    Handles: 'Feb 2026', 'February 2026', 'Feb 17-Mar 3, 2026',
             'Mar 26-28, 2021', 'Oct 2025', etc.
    Returns e.g. 'feb 2026' or None if not parseable.
    """
    if not label:
        return None
    s = label.strip().lower()

    # Try "Month YYYY" or "Mon YYYY" pattern
    m = re.match(r"(\w+)\s+(\d{4})", s)
    if m:
        month_str = m.group(1)
        year = m.group(2)
        if month_str in _MONTH_MAP:
            return f"{month_str[:3]} {year}"

    # Try "Mon DD-DD, YYYY" or "Mon DD, YYYY" pattern (exact collection dates)
    m = re.match(r"(\w+)\s+\d{1,2}(?:[,-]\s*(?:\w+\s+)?\d{1,2})?,?\s*(\d{4})", s)
    if m:
        month_str = m.group(1)
        year = m.group(2)
        if month_str in _MONTH_MAP:
            return f"{month_str[:3]} {year}"

    return None


# ---------------------------------------------------------------------------
# Slide context
# ---------------------------------------------------------------------------

@dataclass
class SlideContext:
    """Text context extracted from a slide's text boxes."""
    slide_idx: int
    question_text: str = ""
    chart_title: str = ""
    base_text: str = ""
    q_code: Optional[str] = None
    all_text: str = ""  # concatenation of all text on the slide


def _extract_slide_contexts(prs: Presentation) -> Dict[int, SlideContext]:
    """Extract text context from every slide that has a chart or table."""
    contexts: Dict[int, SlideContext] = {}

    for slide_idx, slide in enumerate(prs.slides):
        # Check if slide has chart or table
        has_data_shape = False
        for shp in slide.shapes:
            try:
                _ = shp.chart
                has_data_shape = True
            except (ValueError, AttributeError):
                pass
            if shp.has_table:
                has_data_shape = True

        if not has_data_shape:
            continue

        ctx = SlideContext(slide_idx=slide_idx)
        all_texts: List[str] = []

        for shp in slide.shapes:
            # Skip data shapes
            is_data = False
            try:
                _ = shp.chart
                is_data = True
            except (ValueError, AttributeError):
                pass
            if shp.has_table:
                is_data = True
            if is_data:
                continue

            if not hasattr(shp, "text_frame"):
                continue
            text = shp.text_frame.text.strip()
            if not text or len(text) < 3:
                continue

            all_texts.append(text)
            lower = text.lower()

            # Classify text
            if lower.startswith("question:") or lower.startswith("q:"):
                ctx.question_text = text
            elif lower.startswith("base:") or ("respondent" in lower and ("base" in lower or "n=" in lower)):
                ctx.base_text = text
            elif not ctx.chart_title and len(text) < 120:
                # Short text that isn't question/base is likely the chart title
                # Skip if it looks like a callout (contains %)
                if not re.search(r"\d+\.?\d*%", text):
                    ctx.chart_title = text

            # Extract Q-code from any text
            qc = _extract_qcode(text)
            if qc:
                ctx.q_code = qc

        ctx.all_text = " ".join(all_texts)
        contexts[slide_idx] = ctx

    return contexts


def _slide_context_score(
    ctx: Optional[SlideContext],
    table: Dict[str, Any],
) -> float:
    """Score how well a slide's text context matches a crosstab table.

    Returns 0.0–1.0.
    """
    if ctx is None:
        return 0.0

    table_title = table.get("title", "")

    # Q-code exact match is the strongest signal
    if ctx.q_code and ctx.q_code in table_title:
        return 1.0

    # Also check if the table title contains a Q-code that appears in slide text
    table_qcode = _extract_qcode(table_title)
    if table_qcode and ctx.all_text and table_qcode in ctx.all_text:
        return 0.95

    # Fuzzy match question text against table title
    best_score = 0.0

    if ctx.question_text:
        # Strip "Question: " prefix for cleaner matching
        q = ctx.question_text
        if q.lower().startswith("question:"):
            q = q[9:].strip()
        # The table title often has the Q-code + full question text
        # Extract just the question part from the table title for comparison
        table_q = re.sub(r"\(Q\d+[_\d]*\)\s*\d*\.\s*", "", table_title).strip()
        # Remove trailing "by Demographic Banner" etc
        table_q = re.sub(r"\s+by\s+\w+\s+Banner.*$", "", table_q, flags=re.IGNORECASE).strip()

        ratio = SequenceMatcher(None, _norm(q)[:200], _norm(table_q)[:200]).ratio()
        best_score = max(best_score, ratio)

    if ctx.chart_title:
        # Chart titles are often short descriptive names like "Travel as a Budget Priority"
        # Check if keywords from the chart title appear in the table's row labels
        title_words = set(_norm(ctx.chart_title).split()) - {"by", "the", "a", "an", "of", "in", "for", "and", "or", "to", "is", "it"}
        row_text = " ".join(_norm(rl) for rl in table.get("row_labels", []))
        if title_words:
            word_hits = sum(1 for w in title_words if w in row_text or w in _norm(table_title))
            word_score = min(1.0, word_hits / max(3, len(title_words)))
            best_score = max(best_score, word_score * 0.7)

    return best_score


# ---------------------------------------------------------------------------
# Series-to-row matching
# ---------------------------------------------------------------------------

def _series_name_score(
    series_names: List[str],
    table: Dict[str, Any],
) -> Tuple[float, Optional[str]]:
    """Score how well chart series names match table row labels.

    Returns (score, matched_row_label).
    """
    if _is_uninformative_series(series_names):
        return 0.0, None

    row_labels = table.get("row_labels", [])
    if not row_labels:
        return 0.0, None

    best_score = 0.0
    best_row = None

    for sname in series_names:
        ns = _norm(sname)
        if not ns or ns in _UNINFORMATIVE_SERIES:
            continue

        # Try exact/substring match against row labels
        for rl in row_labels:
            nrl = _norm(rl)
            if nrl.startswith(("base", "mean", "average", "avg")):
                continue

            # Exact match
            if ns == nrl:
                return 1.0, rl

            # Substring containment (either direction)
            if ns in nrl or nrl in ns:
                score = max(0.80, SequenceMatcher(None, ns, nrl).ratio())
                if score > best_score:
                    best_score = score
                    best_row = rl
                continue

            # Fuzzy match
            ratio = SequenceMatcher(None, ns, nrl).ratio()
            if ratio > best_score and ratio >= 0.55:
                best_score = ratio
                best_row = rl

        # Check for "Top N Box" pattern:
        # Series name like "Better or Much Better Off" maps to "Top 2 Box"
        if " or " in ns or " and " in ns:
            for rl in row_labels:
                nrl = _norm(rl)
                if nrl.startswith("top") and "box" in nrl:
                    # This is likely the aggregate row
                    score = 0.70  # Good confidence for Top N Box match
                    if score > best_score:
                        best_score = score
                        best_row = rl

        # Also try matching series name against keywords in the table title
        table_title = _norm(table.get("title", ""))
        if ns and table_title:
            # If series name words appear in the table title, that's a weak but useful signal
            s_words = set(ns.split()) - {"the", "a", "an", "of", "in", "for", "and", "or", "to"}
            title_words = set(table_title.split())
            if s_words and len(s_words & title_words) / len(s_words) >= 0.3:
                score = 0.40
                if score > best_score:
                    best_score = score
                    # Don't set best_row here — the series name didn't match a specific row

    return best_score, best_row


# ---------------------------------------------------------------------------
# Fingerprinting
# ---------------------------------------------------------------------------

@dataclass
class ShapeFingerprint:
    """Structural data extracted from a PowerPoint shape (no alt text needed)."""
    slide_idx: int
    shape_idx: int
    shape_name: str
    shape_type: str                        # "chart" | "table" | "text"
    categories: List[str] = field(default_factory=list)
    series_names: List[str] = field(default_factory=list)
    col_headers: List[str] = field(default_factory=list)
    row_labels: List[str] = field(default_factory=list)
    text_content: str = ""
    values_sample: List[float] = field(default_factory=list)
    existing_alt: str = ""
    # Full value data per series (for multi-series correlation)
    series_values: List[List[float]] = field(default_factory=list)


def _extract_chart_fingerprint(shape, slide_idx: int, shape_idx: int) -> Optional[ShapeFingerprint]:
    """Pull categories, series names, and sample values from a chart shape."""
    try:
        chart = shape.chart
    except (ValueError, AttributeError):
        return None

    categories = []
    try:
        for pt in chart.plots[0].categories:
            if pt is not None:
                categories.append(str(pt))
    except Exception:
        pass

    series_names = []
    values_sample = []
    series_values = []
    try:
        for s in chart.series:
            try:
                series_names.append(str(s.name) if s.name else "")
            except Exception:
                series_names.append("")
            s_vals = []
            try:
                for v in s.values:
                    fv = float(v) if v is not None else None
                    if fv is not None:
                        values_sample.append(fv)
                    s_vals.append(fv)
            except Exception:
                pass
            series_values.append(s_vals)
    except Exception:
        pass

    if not categories and not series_names:
        return None

    return ShapeFingerprint(
        slide_idx=slide_idx,
        shape_idx=shape_idx,
        shape_name=shape.name or "",
        shape_type="chart",
        categories=categories,
        series_names=series_names,
        values_sample=values_sample[:50],
        series_values=series_values,
        existing_alt=_read_existing_alt(shape),
    )


def _extract_table_fingerprint(shape, slide_idx: int, shape_idx: int) -> Optional[ShapeFingerprint]:
    """Pull header row and first-column labels from a table shape."""
    if not shape.has_table:
        return None
    tbl = shape.table
    if len(tbl.rows) < 2 or len(tbl.columns) < 2:
        return None

    col_headers = []
    for c in range(1, len(tbl.columns)):
        col_headers.append(tbl.cell(0, c).text_frame.text.strip())

    row_labels = []
    for r in range(1, len(tbl.rows)):
        row_labels.append(tbl.cell(r, 0).text_frame.text.strip())

    if not col_headers and not row_labels:
        return None

    return ShapeFingerprint(
        slide_idx=slide_idx,
        shape_idx=shape_idx,
        shape_name=shape.name or "",
        shape_type="table",
        col_headers=col_headers,
        row_labels=row_labels,
        existing_alt=_read_existing_alt(shape),
    )


def _extract_text_fingerprint(shape, slide_idx: int, shape_idx: int) -> Optional[ShapeFingerprint]:
    """Pull text content from a text box shape."""
    if not hasattr(shape, "text_frame"):
        return None
    text = shape.text_frame.text.strip()
    if not text or len(text) < 5:
        return None
    return ShapeFingerprint(
        slide_idx=slide_idx,
        shape_idx=shape_idx,
        shape_name=shape.name or "",
        shape_type="text",
        text_content=text,
    )


def extract_all_fingerprints(prs: Presentation) -> List[ShapeFingerprint]:
    """Walk the entire presentation and extract fingerprints from all shapes."""
    fingerprints = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            fp = _extract_chart_fingerprint(shape, slide_idx, shape_idx)
            if fp:
                fingerprints.append(fp)
                continue
            fp = _extract_table_fingerprint(shape, slide_idx, shape_idx)
            if fp:
                fingerprints.append(fp)
                continue
            fp = _extract_text_fingerprint(shape, slide_idx, shape_idx)
            if fp:
                fingerprints.append(fp)
    return fingerprints


# ---------------------------------------------------------------------------
# Label overlap scoring (preserved from v1)
# ---------------------------------------------------------------------------

_LABEL_ROW_WEIGHT = 0.65
_LABEL_COL_WEIGHT = 0.35


def _containment(subset: set, superset: set) -> float:
    if not subset:
        return 0.0
    return len(subset & superset) / len(subset)


def _fuzzy_containment(subset_labels: set, superset_labels: set, threshold: float = 0.7) -> float:
    if not subset_labels:
        return 0.0
    hits = 0
    for label in subset_labels:
        if label in superset_labels:
            hits += 1
            continue
        for cand in superset_labels:
            if SequenceMatcher(None, label, cand).ratio() >= threshold:
                hits += 1
                break
    return hits / len(subset_labels)


def _sim_blend(fp_set: set, t_set: set, use_fuzzy: bool = False) -> float:
    if use_fuzzy:
        contain = _fuzzy_containment(fp_set, t_set, threshold=0.7)
    else:
        contain = _containment(fp_set, t_set)
    jacc = _jaccard(fp_set, t_set)
    return 0.6 * contain + 0.4 * jacc


@dataclass
class MapCandidate:
    """A candidate table match for a fingerprint."""
    table: Dict[str, Any]
    score: float
    # Sub-scores for debugging/UI
    value_corr_score: float = 0.0
    label_score: float = 0.0
    context_score: float = 0.0
    series_score: float = 0.0
    row_score: float = 0.0
    col_score: float = 0.0
    orientation: str = "normal"
    matched_row: Optional[str] = None


def _label_overlap_score(
    fp: ShapeFingerprint,
    table: Dict[str, Any],
) -> Tuple[float, float, str]:
    """Compute label overlap score between a fingerprint and a table.

    Returns (row_score, col_score, orientation).
    """
    t_rows = {_norm(r) for r in table.get("row_labels", []) if isinstance(r, str)}
    t_cols = {_norm(c) for c in table.get("col_labels", []) if isinstance(c, str)}

    if fp.shape_type == "chart":
        fp_cats = {_norm(c) for c in fp.categories if c}
        fp_ser = {_norm(s) for s in fp.series_names if s}
        uninformative = _is_uninformative_series(fp.series_names)

        # Normal orientation
        row_sim_n = _sim_blend(fp_cats, t_rows)
        col_sim_n = _sim_blend(fp_ser, t_cols) if not uninformative else 0.0
        rw = 1.0 if uninformative else _LABEL_ROW_WEIGHT
        cw = 0.0 if uninformative else _LABEL_COL_WEIGHT
        score_n = rw * row_sim_n + cw * col_sim_n

        best_score = score_n
        best_row = row_sim_n
        best_col = col_sim_n
        best_orient = "normal"

        # Flipped orientation
        if not uninformative:
            row_sim_f = _sim_blend(fp_cats, t_cols)
            col_sim_f = _sim_blend(fp_ser, t_rows, use_fuzzy=True)
            score_f = _LABEL_ROW_WEIGHT * row_sim_f + _LABEL_COL_WEIGHT * col_sim_f
            if score_f > best_score:
                best_row = row_sim_f
                best_col = col_sim_f
                best_orient = "flipped"

        return best_row, best_col, best_orient

    elif fp.shape_type == "table":
        fp_rows = {_norm(r) for r in fp.row_labels if r}
        fp_cols = {_norm(c) for c in fp.col_headers if c}
        row_sim = _sim_blend(fp_rows, t_rows)
        col_sim = _sim_blend(fp_cols, t_cols)
        return row_sim, col_sim, "normal"

    return 0.0, 0.0, "normal"


# ---------------------------------------------------------------------------
# Blended scoring
# ---------------------------------------------------------------------------

def score_fingerprint_blended(
    fp: ShapeFingerprint,
    tables: List[Dict[str, Any]],
    slide_ctx: Optional[SlideContext] = None,
    col_banners_pool: Optional[set] = None,
) -> List[MapCandidate]:
    """Score a fingerprint against all tables using all available signals.

    Returns candidates sorted by final blended score (descending).
    """
    if fp.shape_type not in ("chart", "table"):
        return []

    # Determine chart orientation type
    is_trend = fp.shape_type == "chart" and _is_trend_chart(fp.categories)
    is_segment = (
        fp.shape_type == "chart"
        and not is_trend
        and _is_segment_chart(fp.categories, col_banners_pool or set())
    )

    # Get first series values for correlation
    first_series_vals = []
    if fp.series_values:
        first_series_vals = [v for v in fp.series_values[0] if v is not None]
    elif fp.values_sample:
        first_series_vals = fp.values_sample

    candidates: List[MapCandidate] = []

    for t in tables:
        # --- Value correlation ---
        val_score = 0.0
        matched_row = None

        if first_series_vals and len(first_series_vals) >= 3:
            if is_trend:
                val_score, matched_row = _value_correlation_timeseries(
                    first_series_vals, fp.categories, t,
                )
            elif is_segment:
                val_score, matched_row = _value_correlation_segment(
                    first_series_vals, fp.categories, t,
                )
            else:
                val_score = _value_correlation_distribution(
                    first_series_vals, fp.categories, t,
                )

        # --- Label overlap ---
        row_sim, col_sim, orient = _label_overlap_score(fp, t)
        uninformative = _is_uninformative_series(fp.series_names) if fp.shape_type == "chart" else False
        rw = 1.0 if uninformative else _LABEL_ROW_WEIGHT
        cw = 0.0 if uninformative else _LABEL_COL_WEIGHT
        label_score = rw * row_sim + cw * col_sim

        # --- Slide context ---
        ctx_score = _slide_context_score(slide_ctx, t)

        # --- Series-to-row matching ---
        series_score = 0.0
        series_matched_row = None
        if fp.shape_type == "chart":
            series_score, series_matched_row = _series_name_score(fp.series_names, t)

        # If segment/trend correlation identified a row, prefer that; else use series match
        if matched_row is None and series_matched_row:
            matched_row = series_matched_row

        # --- Blend ---
        # For time series and segment charts, value correlation is more important
        if is_trend or is_segment:
            w_val = 0.55
            w_label = 0.15
            w_ctx = 0.20
            w_series = 0.10
        else:
            w_val = WEIGHT_VALUE_CORR
            w_label = WEIGHT_LABEL_OVERLAP
            w_ctx = WEIGHT_SLIDE_CONTEXT
            w_series = WEIGHT_SERIES_NAME

        # Cascade-redistribute weight from absent signals to active ones.
        # A signal is "absent" when its input data is missing entirely
        # (not just a low score — a low score is still informative).
        has_val = bool(first_series_vals) and len(first_series_vals) >= 3
        has_ctx = (
            slide_ctx is not None
            and bool(slide_ctx.q_code or slide_ctx.question_text or slide_ctx.chart_title)
        )
        has_series = (
            fp.shape_type == "chart"
            and not _is_uninformative_series(fp.series_names)
        )

        absent_weight = 0.0
        if not has_val:
            absent_weight += w_val
            w_val = 0.0
        if not has_ctx:
            absent_weight += w_ctx
            w_ctx = 0.0
        if not has_series:
            absent_weight += w_series
            w_series = 0.0

        # Distribute the collected absent weight proportionally among
        # the signals that are still active.
        active_total = w_val + w_label + w_ctx + w_series
        if active_total > 0 and absent_weight > 0:
            scale = (active_total + absent_weight) / active_total
            w_val *= scale
            w_label *= scale
            w_ctx *= scale
            w_series *= scale

        final_score = (
            w_val * val_score
            + w_label * label_score
            + w_ctx * ctx_score
            + w_series * series_score
        )

        # Override orientation for segment/trend charts
        if is_segment:
            orient = "flipped"
        elif is_trend:
            orient = "timeseries"

        candidates.append(MapCandidate(
            table=t,
            score=final_score,
            value_corr_score=val_score,
            label_score=label_score,
            context_score=ctx_score,
            series_score=series_score,
            row_score=row_sim,
            col_score=col_sim,
            orientation=orient,
            matched_row=matched_row,
        ))

    candidates.sort(key=lambda c: c.score, reverse=True)
    return candidates


# Also keep the old function name for backward compat
def score_fingerprint_against_tables(
    fp: ShapeFingerprint,
    tables: List[Dict[str, Any]],
) -> List[MapCandidate]:
    """Legacy wrapper — calls score_fingerprint_blended without context signals."""
    return score_fingerprint_blended(fp, tables)


# ---------------------------------------------------------------------------
# LLM disambiguation (enhanced with richer context)
# ---------------------------------------------------------------------------

def _llm_disambiguate(
    fp: ShapeFingerprint,
    candidates: List[MapCandidate],
    slide_ctx: Optional[SlideContext] = None,
    max_candidates: int = 5,
) -> Optional[Tuple[MapCandidate, str]]:
    """Use an LLM to pick the best candidate when scoring is ambiguous."""
    try:
        from openai import OpenAI
    except ImportError:
        logger.error("openai package not installed — skipping LLM disambiguation")
        return None

    top = candidates[:max_candidates]
    if not top:
        return None

    if fp.shape_type == "chart":
        shape_info = {
            "shape_name": fp.shape_name,
            "type": "chart",
            "categories": fp.categories[:15],
            "series_names": fp.series_names,
        }
    elif fp.shape_type == "table":
        shape_info = {
            "shape_name": fp.shape_name,
            "type": "table",
            "column_headers": fp.col_headers,
            "row_labels": fp.row_labels[:15],
        }
    else:
        return None

    # Add slide context to help the LLM
    if slide_ctx:
        shape_info["slide_question_text"] = slide_ctx.question_text[:200] if slide_ctx.question_text else ""
        shape_info["slide_chart_title"] = slide_ctx.chart_title[:100] if slide_ctx.chart_title else ""
        shape_info["slide_q_code"] = slide_ctx.q_code or ""

    prompt_candidates = []
    for i, c in enumerate(top):
        prompt_candidates.append({
            "index": i,
            "title": c.table.get("title", ""),
            "row_labels": c.table.get("row_labels", [])[:15],
            "col_labels": c.table.get("col_labels", [])[:10],
            "value_correlation": round(c.value_corr_score, 3),
            "label_overlap": round(c.label_score, 3),
            "context_score": round(c.context_score, 3),
        })

    user_msg = (
        "I have a PowerPoint chart/table shape that needs to be matched to one of "
        "several crosstab data tables. The shape's existing data was extracted as follows:\n\n"
        f"Shape: {json.dumps(shape_info)}\n\n"
        "Candidate tables from the crosstab data (with algorithmic scores):\n"
        f"{json.dumps(prompt_candidates)}\n\n"
        "Which candidate table is the best match? Consider: value correlation scores "
        "(higher = chart data pattern matches table data pattern), label overlap, "
        "and whether the slide's question text matches the table title.\n\n"
        'Return JSON: {"best_index": <int or null>, "reason": "<one sentence>"}'
    )

    try:
        client = OpenAI()
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            max_tokens=150,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You are a data-matching assistant that maps PowerPoint chart/table shapes to their corresponding crosstab data tables."},
                {"role": "user", "content": user_msg},
            ],
        )
        body = json.loads(response.choices[0].message.content)
        idx = body.get("best_index")
        reason = body.get("reason", "")
        if idx is not None and 0 <= idx < len(top):
            return top[idx], reason
    except Exception as e:
        logger.error("LLM disambiguation failed: %s", e)

    return None


# ---------------------------------------------------------------------------
# Mapping result
# ---------------------------------------------------------------------------

@dataclass
class AutoMapResult:
    """Result of auto-mapping one shape."""
    fingerprint: ShapeFingerprint
    table_title: Optional[str]
    col_key: Optional[str]
    confidence: float
    method: str  # "qcode" | "structural" | "blended" | "llm" | "slide_context" | "trend_match" | "unmatched"
    reason: str = ""
    candidates: List[MapCandidate] = field(default_factory=list)
    orientation: str = "normal"
    row_key: Optional[str] = None
    column_keys: Optional[List[str]] = None
    # Sub-scores for UI display
    value_corr_score: float = 0.0
    label_score: float = 0.0
    context_score: float = 0.0
    series_score: float = 0.0
    matched_row: Optional[str] = None


# ---------------------------------------------------------------------------
# Alt text writer
# ---------------------------------------------------------------------------

def _build_alt_text(
    table_title: str,
    col_key: Optional[str],
    shape_type: str,
    *,
    orientation: str = "normal",
    row_key: Optional[str] = None,
    column_keys: Optional[List[str]] = None,
) -> str:
    """Build the standard alt text format that deck_update expects."""
    lines = [f"table_title: {table_title}"]
    if orientation == "flipped" and row_key:
        lines.append(f"row_key: {row_key}")
    if orientation == "timeseries":
        lines.append("chart_mode: timeseries")
        if row_key:
            lines.append(f"row_key: {row_key}")
    if column_keys:
        lines.append(f"column: {','.join(column_keys)}")
    elif col_key:
        lines.append(f"column: {col_key}")
    if shape_type == "chart":
        lines.append("type: chart")
    elif shape_type == "table":
        lines.append("type: table")
    if orientation == "flipped":
        lines.append("orientation: flipped")
    elif orientation == "timeseries":
        lines.append("orientation: timeseries")
    lines.append("auto_update: yes")
    return "\n".join(lines)


def _build_text_alt_text(table_title: str, text_type: str) -> str:
    return f"type: {text_type}\ntable_title: {table_title}"


def _write_alt_text(shape, alt_text: str):
    """Write alt text onto a shape's cNvPr descr attribute."""
    el = shape.element
    c_nv_pr = el.find(
        ".//p:cNvPr",
        namespaces={"p": P_NS},
    )
    if c_nv_pr is not None:
        c_nv_pr.set("descr", alt_text)
        return
    try:
        shape.alternative_text = alt_text
    except (AttributeError, ValueError):
        logger.warning("Could not write alt text to shape '%s'", shape.name)


def _read_existing_alt(shape) -> str:
    try:
        el = shape.element
        c_nv_pr = el.find(
            ".//p:cNvPr",
            namespaces={"p": P_NS},
        )
        if c_nv_pr is not None:
            return c_nv_pr.get("descr", "")
    except Exception:
        pass
    try:
        return shape.alternative_text or ""
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Text shape association (slide context)
# ---------------------------------------------------------------------------

def _classify_text_shape(text: str) -> Optional[str]:
    lower = text.strip().lower()
    if lower.startswith("question:") or lower.startswith("q:"):
        return "text_question"
    if lower.startswith("base:") or lower.startswith("base n"):
        return "text_base"
    if "respondent" in lower and ("base" in lower or "n=" in lower or "n =" in lower):
        return "text_base"
    return None


def associate_text_shapes(
    prs: Presentation,
    map_results: List[AutoMapResult],
) -> List[AutoMapResult]:
    slide_to_table: Dict[int, str] = {}
    for r in map_results:
        if r.table_title and r.fingerprint.shape_type in ("chart", "table"):
            slide_idx = r.fingerprint.slide_idx
            if slide_idx not in slide_to_table:
                slide_to_table[slide_idx] = r.table_title

    text_results: List[AutoMapResult] = []
    slides = list(prs.slides)

    for slide_idx, table_title in slide_to_table.items():
        if slide_idx >= len(slides):
            continue
        slide = slides[slide_idx]
        for shape_idx, shape in enumerate(slide.shapes):
            existing_alt = _read_existing_alt(shape)
            if existing_alt and "table_title" in existing_alt.lower():
                continue
            is_data_shape = False
            try:
                _ = shape.chart
                is_data_shape = True
            except (ValueError, AttributeError):
                pass
            if shape.has_table:
                is_data_shape = True
            if is_data_shape:
                continue
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text_frame.text.strip()
            if not text or len(text) < 3:
                continue
            text_type = _classify_text_shape(text)
            if text_type is None:
                continue
            fp = ShapeFingerprint(
                slide_idx=slide_idx,
                shape_idx=shape_idx,
                shape_name=shape.name or "",
                shape_type="text",
                text_content=text,
            )
            text_results.append(AutoMapResult(
                fingerprint=fp,
                table_title=table_title,
                col_key=None,
                confidence=0.80,
                method="slide_context",
                reason=f"Text shape classified as '{text_type}' on same slide as '{table_title}'",
            ))

    return text_results


# ---------------------------------------------------------------------------
# Key inference helpers
# ---------------------------------------------------------------------------

def _infer_col_key(fp: ShapeFingerprint, table: Dict[str, Any]) -> Optional[str]:
    col_labels = table.get("col_labels", [])
    if not col_labels:
        return None
    if fp.shape_type == "chart" and len(fp.series_names) == 1:
        sn = fp.series_names[0]
        if sn in col_labels:
            return sn
    for preferred in ["Total", "Overall", "All", "Base"]:
        if preferred in col_labels:
            return preferred
    return col_labels[0] if col_labels else None


def _infer_row_key(fp: ShapeFingerprint, table: Dict[str, Any]) -> Optional[str]:
    row_labels = table.get("row_labels", [])
    if not row_labels or not fp.series_names:
        return None
    for sn in fp.series_names:
        if not sn or _norm(sn) in _UNINFORMATIVE_SERIES:
            continue
        if sn in row_labels:
            return sn
        fuzzy = _fuzzy_match_label(sn, row_labels, threshold=0.55)
        if fuzzy:
            return fuzzy
    return None


def _infer_column_keys(fp: ShapeFingerprint, table: Dict[str, Any]) -> Optional[List[str]]:
    col_labels = table.get("col_labels", [])
    if not col_labels or not fp.categories:
        return None
    keys: List[str] = []
    for cat in fp.categories:
        if cat in col_labels:
            keys.append(cat)
        else:
            fuzzy = _fuzzy_match_label(cat, col_labels, threshold=0.6)
            if fuzzy:
                keys.append(fuzzy)
    return keys if keys else None


def _infer_column_keys_from_banners(fp: ShapeFingerprint, table: Dict[str, Any]) -> Optional[List[str]]:
    """For segment breakout charts: map chart categories to crosstab column banners."""
    banners = table.get("meta", {}).get("col_banners", [])
    col_labels = table.get("col_labels", [])
    if not banners or not fp.categories:
        return None

    seg_map = _build_segment_mapping(fp.categories, banners)
    if not seg_map:
        return None

    # Return the actual col_labels (which may include metric qualifiers like "Total | %")
    keys = []
    for ci in sorted(seg_map.keys()):
        bi = seg_map[ci]
        if bi < len(col_labels):
            keys.append(col_labels[bi])
        elif bi < len(banners):
            keys.append(banners[bi])
    return keys if keys else None


def _infer_keys(
    fp: ShapeFingerprint,
    table: Dict[str, Any],
    orientation: str,
    matched_row: Optional[str] = None,
) -> Tuple[Optional[str], Optional[str], Optional[List[str]]]:
    """Return (col_key, row_key, column_keys) appropriate for orientation."""
    if orientation == "flipped":
        row_key = matched_row or _infer_row_key(fp, table)
        column_keys = _infer_column_keys_from_banners(fp, table) or _infer_column_keys(fp, table)
        return None, row_key, column_keys
    elif orientation == "timeseries":
        row_key = matched_row or _infer_row_key(fp, table)
        return None, row_key, None
    return _infer_col_key(fp, table), None, None


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------

def auto_map_presentation(
    pptx_in: str,
    crosstab_xlsx: str,
    pptx_out: str,
    *,
    use_llm: bool = True,
    write_alt: bool = True,
    progress_callback=None,
) -> List[AutoMapResult]:
    prs = Presentation(pptx_in)
    data = parse_workbook(crosstab_xlsx)
    tables = data["tables"]
    return auto_map_presentation_obj(
        prs, tables,
        pptx_out=pptx_out,
        use_llm=use_llm,
        write_alt=write_alt,
        progress_callback=progress_callback,
    )


def auto_map_presentation_obj(
    prs: Presentation,
    tables: List[Dict[str, Any]],
    *,
    pptx_out: Optional[str] = None,
    use_llm: bool = True,
    write_alt: bool = True,
    progress_callback=None,
) -> List[AutoMapResult]:
    """Core auto-mapping logic with multi-signal blended scoring.

    v2 changes:
    - Value correlation as primary matching signal
    - Slide context (question text, Q-codes) as secondary signal
    - Series-to-row fuzzy matching
    - Time series charts are matched (not skipped)
    - Enhanced LLM fallback with richer context
    """
    fingerprints = extract_all_fingerprints(prs)
    data_fps = [fp for fp in fingerprints if fp.shape_type in ("chart", "table")]

    # Extract slide contexts
    slide_contexts = _extract_slide_contexts(prs)

    # Build pool of all column banners for segment detection
    col_banners_pool: set = set()
    for t in tables:
        banners = t.get("meta", {}).get("col_banners", t.get("col_labels", []))
        for b in banners:
            if b and isinstance(b, str):
                col_banners_pool.add(_norm(b))

    # Build Q-code → table index for Tier 0 matching
    qcode_index: Dict[str, List[Dict[str, Any]]] = {}
    for t in tables:
        qc = _extract_qcode(t.get("title", ""))
        if qc:
            qcode_index.setdefault(qc, []).append(t)

    results: List[AutoMapResult] = []
    total = len(data_fps) or 1

    for i, fp in enumerate(data_fps):
        slide_ctx = slide_contexts.get(fp.slide_idx)

        # --- Tier 0: Q-code match from existing alt text ---
        fp_qcode = _extract_qcode(fp.existing_alt)
        if fp_qcode and fp_qcode in qcode_index:
            matched_tables = qcode_index[fp_qcode]
            qc_table = matched_tables[0]
            col_key = _infer_col_key(fp, qc_table)
            results.append(AutoMapResult(
                fingerprint=fp, table_title=qc_table["title"], col_key=col_key,
                confidence=0.95, method="qcode",
                reason=f"Q-code {fp_qcode} matched table title",
            ))
            if progress_callback:
                progress_callback((i + 1) / total)
            continue

        # --- Blended multi-signal scoring ---
        candidates = score_fingerprint_blended(
            fp, tables, slide_ctx=slide_ctx, col_banners_pool=col_banners_pool,
        )

        if not candidates:
            results.append(AutoMapResult(
                fingerprint=fp, table_title=None, col_key=None,
                confidence=0.0, method="unmatched",
                reason="No candidate tables available",
            ))
            if progress_callback:
                progress_callback((i + 1) / total)
            continue

        best = candidates[0]

        if best.score >= HIGH_CONFIDENCE_THRESHOLD:
            orient = best.orientation
            col_key, row_key, column_keys = _infer_keys(
                fp, best.table, orient, matched_row=best.matched_row,
            )
            title = best.table["title"]
            results.append(AutoMapResult(
                fingerprint=fp, table_title=title, col_key=col_key,
                confidence=best.score, method="blended",
                reason=(
                    f"Val r={best.value_corr_score:.2f}, "
                    f"Labels={best.label_score:.2f}, "
                    f"Ctx={best.context_score:.2f}, "
                    f"Series={best.series_score:.2f}"
                ),
                candidates=candidates[:5],
                orientation=orient,
                row_key=row_key,
                column_keys=column_keys,
                value_corr_score=best.value_corr_score,
                label_score=best.label_score,
                context_score=best.context_score,
                series_score=best.series_score,
                matched_row=best.matched_row,
            ))
        elif best.score >= LLM_THRESHOLD and use_llm:
            llm_result = _llm_disambiguate(fp, candidates, slide_ctx=slide_ctx)
            if llm_result:
                chosen, reason = llm_result
                orient = chosen.orientation
                col_key, row_key, column_keys = _infer_keys(
                    fp, chosen.table, orient, matched_row=chosen.matched_row,
                )
                title = chosen.table["title"]
                results.append(AutoMapResult(
                    fingerprint=fp, table_title=title, col_key=col_key,
                    confidence=max(chosen.score, 0.75), method="llm",
                    reason=reason,
                    candidates=candidates[:5],
                    orientation=orient,
                    row_key=row_key,
                    column_keys=column_keys,
                    value_corr_score=chosen.value_corr_score,
                    label_score=chosen.label_score,
                    context_score=chosen.context_score,
                    series_score=chosen.series_score,
                    matched_row=chosen.matched_row,
                ))
            else:
                results.append(AutoMapResult(
                    fingerprint=fp, table_title=None, col_key=None,
                    confidence=best.score, method="unmatched",
                    reason="LLM could not resolve ambiguity",
                    candidates=candidates[:5],
                ))
        else:
            results.append(AutoMapResult(
                fingerprint=fp, table_title=None, col_key=None,
                confidence=best.score, method="unmatched",
                reason=f"Best score {best.score:.2f} below threshold",
                candidates=candidates[:5],
            ))

        if progress_callback:
            progress_callback((i + 1) / total)

    # --- Text shape association ---
    text_results = associate_text_shapes(prs, results)
    results.extend(text_results)

    # --- Write alt text onto shapes ---
    if write_alt:
        slides = list(prs.slides)
        for r in results:
            if r.table_title is None:
                continue
            if r.confidence < ALT_TEXT_WRITE_THRESHOLD:
                continue
            fp = r.fingerprint
            if fp.slide_idx >= len(slides):
                continue
            slide = slides[fp.slide_idx]
            shapes = list(slide.shapes)
            if fp.shape_idx >= len(shapes):
                continue
            shape = shapes[fp.shape_idx]

            if fp.shape_type in ("chart", "table"):
                alt = _build_alt_text(
                    r.table_title, r.col_key, fp.shape_type,
                    orientation=r.orientation,
                    row_key=r.row_key,
                    column_keys=r.column_keys,
                )
            elif fp.shape_type == "text":
                text_type = _classify_text_shape(fp.text_content)
                if text_type:
                    alt = _build_text_alt_text(r.table_title, text_type)
                else:
                    continue
            else:
                continue

            _write_alt_text(shape, alt)
            logger.info(
                "Wrote alt text for shape '%s' → table '%s' (method=%s, conf=%.2f)",
                fp.shape_name, r.table_title, r.method, r.confidence,
            )

    if pptx_out:
        prs.save(pptx_out)
        logger.info("Saved auto-mapped presentation to %s", pptx_out)

    return results


# ---------------------------------------------------------------------------
# Report for UI
# ---------------------------------------------------------------------------

def results_to_report(results: List[AutoMapResult]) -> List[dict]:
    """Convert AutoMapResult list into a serializable report for the Streamlit UI."""
    report = []
    for r in results:
        entry = {
            "shape_name": r.fingerprint.shape_name,
            "shape_type": r.fingerprint.shape_type,
            "slide_idx": r.fingerprint.slide_idx,
            "table_title": r.table_title,
            "col_key": r.col_key,
            "confidence": round(r.confidence, 3),
            "method": r.method,
            "reason": r.reason,
            "orientation": r.orientation,
            "row_key": r.row_key,
            "matched_row": r.matched_row,
            "value_corr_score": round(r.value_corr_score, 3),
            "label_score": round(r.label_score, 3),
            "context_score": round(r.context_score, 3),
            "series_score": round(r.series_score, 3),
            "candidates": [
                {
                    "title": c.table.get("title", ""),
                    "score": round(c.score, 3),
                    "value_corr": round(c.value_corr_score, 3),
                    "label": round(c.label_score, 3),
                    "context": round(c.context_score, 3),
                    "series": round(c.series_score, 3),
                    "row_score": round(c.row_score, 3),
                    "col_score": round(c.col_score, 3),
                }
                for c in r.candidates[:5]
            ],
        }
        report.append(entry)
    return report
