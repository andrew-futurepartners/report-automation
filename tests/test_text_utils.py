"""Tests for text_utils.py — base-text parsing, formatting, and safe_update_text."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from text_utils import (
    format_base_text,
    format_number_with_commas,
    parse_base_text,
    safe_update_text,
)


# ---------------------------------------------------------------------------
# format_number_with_commas
# ---------------------------------------------------------------------------

class TestFormatNumberWithCommas:
    def test_integer(self):
        assert format_number_with_commas(1448) == "1,448"

    def test_large_integer(self):
        assert format_number_with_commas(1000000) == "1,000,000"

    def test_small_integer(self):
        assert format_number_with_commas(42) == "42"

    def test_float(self):
        result = format_number_with_commas(1234.5)
        assert result == "1,234.5"

    def test_none_returns_none(self):
        assert format_number_with_commas(None) is None

    def test_zero(self):
        assert format_number_with_commas(0) == "0"


# ---------------------------------------------------------------------------
# parse_base_text
# ---------------------------------------------------------------------------

class TestParseBaseText:
    def test_standard_pattern(self):
        text = "Base: Total respondents. 1,448 complete surveys."
        result = parse_base_text(text)
        assert result["description"] == "Total respondents"
        assert result["n_value"] == 1448

    def test_no_period(self):
        text = "Base: Adults 500"
        result = parse_base_text(text)
        assert result["description"] == "Adults"
        assert result["n_value"] == 500

    def test_no_base_prefix(self):
        result = parse_base_text("something else entirely")
        assert result["description"] == ""
        assert result["n_value"] is None

    def test_empty_string(self):
        result = parse_base_text("")
        assert result["description"] == ""
        assert result["n_value"] is None

    def test_base_with_equals(self):
        text = "Base: All adults =. 2,000 complete surveys."
        result = parse_base_text(text)
        assert result["description"] == "All adults"
        assert result["n_value"] == 2000

    def test_number_without_commas(self):
        text = "Base: Teens. 350 complete surveys."
        result = parse_base_text(text)
        assert result["description"] == "Teens"
        assert result["n_value"] == 350

    def test_base_only_no_n(self):
        text = "Base: Total respondents."
        result = parse_base_text(text)
        assert result["description"] == "Total respondents"
        assert result["n_value"] is None


# ---------------------------------------------------------------------------
# format_base_text
# ---------------------------------------------------------------------------

class TestFormatBaseText:
    def test_with_n(self):
        result = format_base_text("Total respondents", 1448)
        assert result == "Base: Total respondents. 1,448 complete surveys."

    def test_without_n(self):
        result = format_base_text("Adults")
        assert result == "Base: Adults."

    def test_empty_description_defaults(self):
        result = format_base_text("")
        assert result == "Base: Total respondents."

    def test_none_description_defaults(self):
        result = format_base_text(None)
        assert result == "Base: Total respondents."

    def test_n_value_zero(self):
        result = format_base_text("Sample", 0)
        assert result == "Base: Sample. 0 complete surveys."


# ---------------------------------------------------------------------------
# Round-trip: parse_base_text ↔ format_base_text
# ---------------------------------------------------------------------------

class TestRoundTrip:
    @pytest.mark.parametrize("text", [
        "Base: Total respondents. 1,448 complete surveys.",
        "Base: Adults. 500 complete surveys.",
        "Base: Online panel. 2,000 complete surveys.",
    ])
    def test_round_trip_preserves_semantics(self, text):
        parsed = parse_base_text(text)
        rebuilt = format_base_text(parsed["description"], parsed["n_value"])
        assert rebuilt == text

    def test_round_trip_description_only(self):
        text = "Base: Total respondents."
        parsed = parse_base_text(text)
        rebuilt = format_base_text(parsed["description"], parsed["n_value"])
        assert rebuilt == text


# ---------------------------------------------------------------------------
# safe_update_text
# ---------------------------------------------------------------------------

class TestSafeUpdateText:
    def _make_textbox(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        return tb

    def test_basic_update(self):
        tb = self._make_textbox()
        tb.text_frame.text = "old text"
        result = safe_update_text(tb, "new text")
        assert result is True
        assert tb.text_frame.text == "new text"

    def test_preserve_font_false(self):
        tb = self._make_textbox()
        tb.text_frame.text = "original"
        result = safe_update_text(tb, "updated", preserve_font=False)
        assert result is True
        assert tb.text_frame.text == "updated"

    def test_preserve_font_true_with_single_run(self):
        tb = self._make_textbox()
        tf = tb.text_frame
        tf.text = "hello"
        p = tf.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(18)

        result = safe_update_text(tb, "world", preserve_font=True)
        assert result is True
        assert tb.text_frame.text == "world"

    def test_preserve_font_true_with_multiple_runs(self):
        tb = self._make_textbox()
        tf = tb.text_frame
        tf.text = "bold"
        p = tf.paragraphs[0]
        run1 = p.runs[0]
        run1.font.bold = True
        run1.font.size = Pt(14)

        run2 = p.add_run()
        run2.text = " not bold"

        result = safe_update_text(tb, "replaced", preserve_font=True)
        assert result is True
        assert tb.text_frame.text == "replaced"
        assert len(p.runs) == 1
        assert p.runs[0].font.bold is True
        assert p.runs[0].font.size == Pt(14)

    def test_returns_false_for_object_without_text(self):
        class Dummy:
            pass
        assert safe_update_text(Dummy(), "text") is False

    def test_object_with_text_attribute(self):
        class FakeCell:
            text = "old"
        cell = FakeCell()
        result = safe_update_text(cell, "new")
        assert result is True
        assert cell.text == "new"
