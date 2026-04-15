"""Tests for chart_data_patcher.py — value format detection and XML-level patching."""

import pytest
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from chart_data_patcher import (
    _C,
    _C_NS,
    detect_value_format,
    patch_chart_data,
    patch_chart_series,
)


# ---------------------------------------------------------------------------
# detect_value_format
# ---------------------------------------------------------------------------

class TestDetectValueFormat:
    def test_proportions(self):
        assert detect_value_format([0.25, 0.50, 0.75]) == "percentage"

    def test_integers(self):
        assert detect_value_format([10, 20, 30]) == "number"

    def test_large_numbers(self):
        assert detect_value_format([1500, 700, 800]) == "number"

    def test_explicit_override_percentage(self):
        result = detect_value_format([10, 20], alt_metadata={"value_format": "percentage"})
        assert result == "percentage"

    def test_explicit_override_number(self):
        result = detect_value_format([0.5, 0.8], alt_metadata={"value_format": "number"})
        assert result == "number"

    def test_explicit_valueformat_key(self):
        result = detect_value_format([10], alt_metadata={"valueformat": "percentage"})
        assert result == "percentage"

    def test_empty_list(self):
        assert detect_value_format([]) == "number"

    def test_mixed_with_none(self):
        assert detect_value_format([0.5, None, 0.3]) == "percentage"

    def test_all_none(self):
        assert detect_value_format([None, None]) == "number"

    def test_boundary_values(self):
        assert detect_value_format([0.0, 1.0]) == "percentage"

    def test_just_above_one(self):
        assert detect_value_format([0.5, 1.05]) == "percentage"

    def test_above_range(self):
        assert detect_value_format([0.5, 1.2]) == "number"

    def test_negative_values(self):
        assert detect_value_format([-0.1, 0.5]) == "number"

    def test_very_small_max(self):
        assert detect_value_format([0.01, 0.02]) == "number"


# ---------------------------------------------------------------------------
# Helpers for building charts
# ---------------------------------------------------------------------------

def _make_chart_shape(categories, series_values, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """Create a Presentation with one chart and return the chart shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Series 1", series_values)
    shape = slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(5), Inches(3.5), cd,
    )
    return shape, prs


def _make_multi_series_chart(categories, series_dict):
    """Create a chart with multiple series; series_dict maps name→values."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series_dict.items():
        cd.add_series(name, vals)
    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(5), Inches(3.5), cd,
    )
    return shape, prs


def _read_chart_values(chart):
    """Extract (categories, values) from the first series' XML caches."""
    tree = chart.part._element
    ns = {"c": _C_NS}
    ser = tree.find(".//c:ser", ns)
    cats = [v.text for v in ser.findall(".//c:cat//c:strCache//c:v", ns)]
    vals = [v.text for v in ser.findall(".//c:val//c:numCache//c:v", ns)]
    return cats, vals


def _count_series(chart):
    tree = chart.part._element
    ns = {"c": _C_NS}
    plot_tags = ["barChart", "lineChart", "pieChart", "areaChart"]
    for tag in plot_tags:
        plot = tree.find(f".//c:{tag}", ns)
        if plot is not None:
            return len(plot.findall(f"c:ser", ns))
    return 0


# ---------------------------------------------------------------------------
# patch_chart_data
# ---------------------------------------------------------------------------

class TestPatchChartData:
    def test_patch_replaces_data(self):
        shape, _ = _make_chart_shape(["A", "B", "C"], (10, 20, 30))
        chart = shape.chart

        patch_chart_data(chart, ["X", "Y"], [100, 200], value_format="number")

        cats, vals = _read_chart_values(chart)
        assert cats == ["X", "Y"]
        assert vals == ["100", "200"]

    def test_patch_with_percentage_format(self):
        shape, _ = _make_chart_shape(["A", "B"], (10, 20))
        chart = shape.chart

        patch_chart_data(chart, ["P", "Q"], [0.5, 0.8], value_format="percentage")

        ns = {"c": _C_NS}
        fmt = chart.part._element.find(".//c:val//c:numCache//c:formatCode", ns)
        assert fmt is not None
        assert fmt.text == "0.0%"

    def test_patch_with_auto_format(self):
        shape, _ = _make_chart_shape(["A"], (10,))
        chart = shape.chart
        patch_chart_data(chart, ["Cat1"], [0.45], value_format="auto")

        ns = {"c": _C_NS}
        fmt = chart.part._element.find(".//c:val//c:numCache//c:formatCode", ns)
        assert fmt is not None
        assert fmt.text == "0.0%"

    def test_patch_preserves_formatting_nodes(self):
        shape, _ = _make_chart_shape(["A", "B"], (10, 20))
        chart = shape.chart

        tree = chart.part._element
        ns = {"c": _C_NS}
        sp_pr_before = tree.findall(".//c:spPr", ns)

        patch_chart_data(chart, ["X", "Y"], [100, 200], value_format="number")

        sp_pr_after = tree.findall(".//c:spPr", ns)
        assert len(sp_pr_after) >= len(sp_pr_before)

    def test_patch_with_none_values(self):
        shape, _ = _make_chart_shape(["A", "B", "C"], (10, 20, 30))
        chart = shape.chart

        patch_chart_data(chart, ["X", "Y", "Z"], [100, None, 300], value_format="number")

        cats, vals = _read_chart_values(chart)
        assert cats == ["X", "Y", "Z"]
        assert "100" in vals
        assert "300" in vals
        assert len(vals) == 2  # None is skipped


# ---------------------------------------------------------------------------
# patch_chart_series (multi-series)
# ---------------------------------------------------------------------------

class TestPatchChartSeries:
    def test_patch_two_series(self):
        shape, _ = _make_multi_series_chart(
            ["A", "B"], {"S1": (1, 2), "S2": (3, 4)},
        )
        chart = shape.chart

        series_data = [
            ("NewS1", ["X", "Y"], [10, 20]),
            ("NewS2", ["X", "Y"], [30, 40]),
        ]
        patch_chart_series(chart, series_data, value_format="number")

        ns = {"c": _C_NS}
        tree = chart.part._element
        all_ser = tree.findall(f".//{_C}ser")
        assert len(all_ser) == 2

        for ser in all_ser:
            cats = [v.text for v in ser.findall(f".//{_C}strCache//{_C}v")]
            assert "X" in cats or len(cats) > 0

    def test_expand_series(self):
        shape, _ = _make_chart_shape(["A", "B"], (1, 2))
        chart = shape.chart

        series_data = [
            ("S1", ["A", "B"], [10, 20]),
            ("S2", ["A", "B"], [30, 40]),
            ("S3", ["A", "B"], [50, 60]),
        ]
        patch_chart_series(chart, series_data, value_format="number")

        all_ser = chart.part._element.findall(f".//{_C}ser")
        assert len(all_ser) == 3

    def test_shrink_series(self):
        shape, _ = _make_multi_series_chart(
            ["A", "B"], {"S1": (1, 2), "S2": (3, 4), "S3": (5, 6)},
        )
        chart = shape.chart

        series_data = [("Only", ["X"], [99])]
        patch_chart_series(chart, series_data, value_format="number")

        all_ser = chart.part._element.findall(f".//{_C}ser")
        assert len(all_ser) == 1
