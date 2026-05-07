"""Integration tests for deck_update.update_presentation — end-to-end pipeline."""

import os

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from chart_data_patcher import _C_NS
from deck_update import update_presentation

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _set_descr(shape, text: str):
    el = shape.element
    c_nv_pr = el.find(".//p:cNvPr", namespaces={"p": P_NS})
    if c_nv_pr is not None:
        c_nv_pr.set("descr", text)


def _read_descr(shape) -> str:
    el = shape.element
    c_nv_pr = el.find(".//p:cNvPr", namespaces={"p": P_NS})
    if c_nv_pr is not None:
        return c_nv_pr.get("descr", "")
    return ""


def _get_chart_categories(chart):
    ns = {"c": _C_NS}
    tree = chart.part._element
    return [v.text for v in tree.findall(".//c:ser[1]//c:cat//c:strCache//c:v", ns)]


def _get_chart_values(chart):
    ns = {"c": _C_NS}
    tree = chart.part._element
    return [v.text for v in tree.findall(".//c:ser[1]//c:val//c:numCache//c:v", ns)]


# ---------------------------------------------------------------------------
# test_update_presentation_chart
# ---------------------------------------------------------------------------

class TestUpdatePresentationChart:
    def test_chart_and_text_updated(self, mini_pptx, sample_crosstab_xlsx, tmp_path):
        out_path = str(tmp_path / "output.pptx")

        update_presentation(mini_pptx, sample_crosstab_xlsx, out_path)

        prs = Presentation(out_path)
        slide = prs.slides[0]

        chart_shape = None
        q_shape = None
        b_shape = None

        for shp in slide.shapes:
            descr = _read_descr(shp)
            if "type: chart" in descr:
                chart_shape = shp
            elif "type: text_question" in descr:
                q_shape = shp
            elif "type: text_base" in descr:
                b_shape = shp

        # Chart should have been updated with crosstab data
        assert chart_shape is not None
        chart = chart_shape.chart
        cats = _get_chart_categories(chart)
        # The crosstab has rows Aided, Unaided, Top of Mind (Base is excluded by default)
        assert len(cats) >= 1

        # Question text: since original was "Question: Old question text"
        # and it differs from the table title, the pipeline preserves the custom text
        assert q_shape is not None
        q_text = q_shape.text_frame.text
        assert q_text.startswith("Question:")

        # Base text should be formatted
        assert b_shape is not None
        b_text = b_shape.text_frame.text
        assert "Base:" in b_text


# ---------------------------------------------------------------------------
# test_update_presentation_with_selections
# ---------------------------------------------------------------------------

class TestUpdateWithSelections:
    def test_column_key_override(self, mini_pptx, sample_crosstab_xlsx, tmp_path):
        out_path = str(tmp_path / "output_sel.pptx")

        selections = {
            "Brand Awareness": {
                "column_key": "Female",
            }
        }

        update_presentation(
            mini_pptx, sample_crosstab_xlsx, out_path,
            selections=selections,
        )

        prs = Presentation(out_path)
        slide = prs.slides[0]

        for shp in slide.shapes:
            descr = _read_descr(shp)
            if "type: chart" in descr:
                vals = _get_chart_values(shp.chart)
                assert len(vals) >= 1
                break


# ---------------------------------------------------------------------------
# test_update_presentation_skip_auto_update_no
# ---------------------------------------------------------------------------

class TestAutoUpdateNo:
    def test_shape_with_auto_update_no_preserved(self, sample_crosstab_xlsx, tmp_path):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        cd = CategoryChartData()
        cd.categories = ["Old A", "Old B"]
        cd.add_series("Old", (1, 2))
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(4), Inches(3), cd,
        )
        _set_descr(
            chart_frame,
            "table_title: Brand Awareness\nauto_update: no\ntype: chart",
        )

        pptx_in = str(tmp_path / "skip_test.pptx")
        prs.save(pptx_in)

        out_path = str(tmp_path / "skip_out.pptx")
        update_presentation(pptx_in, sample_crosstab_xlsx, out_path)

        prs_out = Presentation(out_path)
        slide_out = prs_out.slides[0]

        for shp in slide_out.shapes:
            try:
                chart = shp.chart
                cats = _get_chart_categories(chart)
                assert "Old A" in cats
                assert "Old B" in cats
                break
            except (ValueError, AttributeError):
                continue


# ---------------------------------------------------------------------------
# test progress_callback fires
# ---------------------------------------------------------------------------

class TestProgressCallback:
    def test_callback_invoked(self, mini_pptx, sample_crosstab_xlsx, tmp_path):
        out_path = str(tmp_path / "progress_out.pptx")
        progress_values = []

        update_presentation(
            mini_pptx, sample_crosstab_xlsx, out_path,
            progress_callback=lambda p: progress_values.append(p),
        )

        assert len(progress_values) >= 1
        assert progress_values[-1] == pytest.approx(1.0)


# ---------------------------------------------------------------------------
# Flipped-orientation chart update
# ---------------------------------------------------------------------------

class TestFlippedChartUpdate:
    """Test _update_chart with flipped-orientation alt text (categories = columns)."""

    def _table(self):
        return {
            "title": "Brand Awareness",
            "row_labels": ["Aided", "Unaided", "Top of Mind", "Base"],
            "col_labels": ["Total", "Male", "Female"],
            "values": [
                [0.85, 0.80, 0.90],
                [0.45, 0.42, 0.48],
                [0.22, 0.20, 0.24],
                [1500, 700, 800],
            ],
        }

    def test_flipped_chart_updated_from_row(self, tmp_path):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from deck_update import _update_chart

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        cd = CategoryChartData()
        cd.categories = ["Total", "Male", "Female"]
        cd.add_series("Old Label", (0, 0, 0))
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        _set_descr(
            chart_frame,
            (
                "table_title: Brand Awareness\n"
                "row_key: Aided\n"
                "column: Total,Male,Female\n"
                "type: chart\n"
                "orientation: flipped\n"
                "auto_update: yes"
            ),
        )

        _update_chart(chart_frame, self._table(), "Total", explicit_rows=None)

        chart = chart_frame.chart
        cats = _get_chart_categories(chart)
        vals = _get_chart_values(chart)
        assert cats == ["Total", "Male", "Female"]
        assert len(vals) == 3
        assert vals[0] == "0.85"
        assert vals[1] == "0.8"
        assert vals[2] == "0.9"

    def test_flipped_chart_fuzzy_row_key(self, tmp_path):
        """row_key with different casing should still resolve via fuzzy match."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from deck_update import _update_chart

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        cd = CategoryChartData()
        cd.categories = ["Total", "Male"]
        cd.add_series("Label", (0, 0))
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        _set_descr(
            chart_frame,
            (
                "table_title: Brand Awareness\n"
                "row_key: top of mind\n"
                "column: Total,Male\n"
                "type: chart\n"
                "orientation: flipped\n"
                "auto_update: yes"
            ),
        )

        _update_chart(chart_frame, self._table(), "Total", explicit_rows=None)

        vals = _get_chart_values(chart_frame.chart)
        assert len(vals) == 2
        assert vals[0] == "0.22"
        assert vals[1] == "0.2"

    def test_normal_chart_unaffected_by_flipped_code(self, tmp_path):
        """A chart without orientation: flipped should still update normally."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from deck_update import _update_chart

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        cd = CategoryChartData()
        cd.categories = ["Aided", "Unaided", "Top of Mind"]
        cd.add_series("Total", (0, 0, 0))
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        _set_descr(
            chart_frame,
            "table_title: Brand Awareness\ncolumn: Total\ntype: chart\nauto_update: yes",
        )

        _update_chart(chart_frame, self._table(), "Total", explicit_rows=None)

        vals = _get_chart_values(chart_frame.chart)
        assert len(vals) == 3
        assert vals[0] == "0.85"
        assert vals[1] == "0.45"
        assert vals[2] == "0.22"
