"""Tests for auto_mapper.py — AI-powered shape-to-table mapping."""

import sys
import os

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from openpyxl import Workbook

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from auto_mapper import (
    ShapeFingerprint,
    MapCandidate,
    AutoMapResult,
    _extract_chart_fingerprint,
    _extract_table_fingerprint,
    _extract_text_fingerprint,
    extract_all_fingerprints,
    score_fingerprint_against_tables,
    _classify_text_shape,
    associate_text_shapes,
    auto_map_presentation_obj,
    results_to_report,
    _norm,
    _jaccard,
    _build_alt_text,
    _write_alt_text,
    _read_existing_alt,
    _extract_qcode,
    _is_uninformative_series,
    _is_trend_chart,
    _fuzzy_match_label,
    _fuzzy_containment,
    _containment,
    _infer_row_key,
    _infer_column_keys,
)

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _set_descr(shape, text: str):
    el = shape.element
    c_nv_pr = el.find(".//p:cNvPr", namespaces={"p": P_NS})
    if c_nv_pr is not None:
        c_nv_pr.set("descr", text)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _sample_tables():
    """Crosstab tables mimicking parse_workbook output."""
    return [
        {
            "id": "Sheet1#1",
            "title": "Brand Awareness",
            "row_labels": ["Aided", "Unaided", "Top of Mind", "Base"],
            "col_labels": ["Total", "Male", "Female"],
            "values": [
                [85.0, 80.0, 90.0],
                [45.0, 42.0, 48.0],
                [22.0, 20.0, 24.0],
                [1500, 700, 800],
            ],
        },
        {
            "id": "Sheet1#2",
            "title": "Purchase Intent",
            "row_labels": ["Definitely", "Probably", "Might", "Probably Not", "Definitely Not", "Base"],
            "col_labels": ["Total", "18-34", "35-54"],
            "values": [
                [30.0, 35.0, 25.0],
                [40.0, 38.0, 42.0],
                [20.0, 22.0, 18.0],
                [5.0, 3.0, 2.0],
                [5.0, 2.0, 13.0],
                [1000, 500, 500],
            ],
        },
        {
            "id": "Sheet1#3",
            "title": "Net Promoter Score",
            "row_labels": ["Promoters", "Passives", "Detractors"],
            "col_labels": ["Total", "Q1", "Q2"],
            "values": [[45, 42, 48], [35, 38, 32], [20, 20, 20]],
        },
    ]


@pytest.fixture
def pptx_with_chart(tmp_path):
    """PPTX with a chart whose categories match Brand Awareness rows."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    cd = CategoryChartData()
    cd.categories = ["Aided", "Unaided", "Top of Mind"]
    cd.add_series("Total", (85, 45, 22))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
    )

    path = str(tmp_path / "chart_deck.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_with_table(tmp_path):
    """PPTX with a table whose headers/rows match Purchase Intent."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    rows, cols = 5, 4
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(6), Inches(3))
    tbl = tbl_shape.table

    headers = ["", "Total", "18-34", "35-54"]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h

    row_labels = ["Definitely", "Probably", "Might", "Probably Not"]
    for r, label in enumerate(row_labels, start=1):
        tbl.cell(r, 0).text = label
        for c in range(1, cols):
            tbl.cell(r, c).text = "0.0"

    path = str(tmp_path / "table_deck.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_with_chart_and_text(tmp_path):
    """PPTX with a chart + question + base text shapes (no alt text)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    cd = CategoryChartData()
    cd.categories = ["Aided", "Unaided", "Top of Mind"]
    cd.add_series("Total", (85, 45, 22))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
    )

    q_box = slide.shapes.add_textbox(Inches(5), Inches(0.5), Inches(4), Inches(0.5))
    q_box.text_frame.text = "Question: Which brands are you aware of?"

    b_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(0.5))
    b_box.text_frame.text = "Base: Total respondents (n=1500)"

    path = str(tmp_path / "chart_text_deck.pptx")
    prs.save(path)
    return path


@pytest.fixture
def sample_crosstab_xlsx(tmp_path):
    """Create a small .xlsx that parse_workbook can consume."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    ws.cell(row=1, column=1, value="Brand Awareness")
    ws.cell(row=3, column=1, value="")
    ws.cell(row=3, column=2, value="Total")
    ws.cell(row=3, column=3, value="Male")
    ws.cell(row=3, column=4, value="Female")

    rows = [
        ("Aided", 85, 80, 90),
        ("Unaided", 45, 42, 48),
        ("Top of Mind", 22, 20, 24),
        ("Base", 1500, 700, 800),
    ]
    for i, (label, *vals) in enumerate(rows, start=4):
        ws.cell(row=i, column=1, value=label)
        for j, v in enumerate(vals, start=2):
            ws.cell(row=i, column=j, value=v)

    path = str(tmp_path / "crosstab.xlsx")
    wb.save(path)
    return path


# ---------------------------------------------------------------------------
# Fingerprint extraction
# ---------------------------------------------------------------------------

class TestChartFingerprint:
    def test_extracts_categories_and_series(self, pptx_with_chart):
        prs = Presentation(pptx_with_chart)
        slide = prs.slides[0]
        for idx, shp in enumerate(slide.shapes):
            fp = _extract_chart_fingerprint(shp, 0, idx)
            if fp is not None:
                assert fp.shape_type == "chart"
                assert "Aided" in fp.categories
                assert "Unaided" in fp.categories
                assert "Top of Mind" in fp.categories
                assert len(fp.series_names) >= 1
                assert len(fp.values_sample) >= 1
                return
        pytest.fail("No chart fingerprint extracted")

    def test_non_chart_returns_none(self, pptx_with_table):
        prs = Presentation(pptx_with_table)
        slide = prs.slides[0]
        for idx, shp in enumerate(slide.shapes):
            if shp.has_table:
                fp = _extract_chart_fingerprint(shp, 0, idx)
                assert fp is None
                return


class TestTableFingerprint:
    def test_extracts_headers_and_rows(self, pptx_with_table):
        prs = Presentation(pptx_with_table)
        slide = prs.slides[0]
        for idx, shp in enumerate(slide.shapes):
            fp = _extract_table_fingerprint(shp, 0, idx)
            if fp is not None:
                assert fp.shape_type == "table"
                assert "Total" in fp.col_headers
                assert "18-34" in fp.col_headers
                assert "Definitely" in fp.row_labels
                assert "Probably" in fp.row_labels
                return
        pytest.fail("No table fingerprint extracted")


class TestTextFingerprint:
    def test_extracts_text_content(self, pptx_with_chart_and_text):
        prs = Presentation(pptx_with_chart_and_text)
        slide = prs.slides[0]
        found_text = False
        for idx, shp in enumerate(slide.shapes):
            fp = _extract_text_fingerprint(shp, 0, idx)
            if fp is not None and fp.text_content.startswith("Question:"):
                found_text = True
                assert fp.shape_type == "text"
                assert "aware" in fp.text_content.lower()
        assert found_text

    def test_short_text_returns_none(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
        box.text_frame.text = "Hi"
        fp = _extract_text_fingerprint(box, 0, 0)
        assert fp is None


class TestExtractAll:
    def test_finds_all_shapes(self, pptx_with_chart_and_text):
        prs = Presentation(pptx_with_chart_and_text)
        fps = extract_all_fingerprints(prs)
        types = {fp.shape_type for fp in fps}
        assert "chart" in types
        assert "text" in types
        assert len(fps) >= 3


# ---------------------------------------------------------------------------
# Structural scoring
# ---------------------------------------------------------------------------

class TestStructuralScoring:
    def test_perfect_row_match(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=["Aided", "Unaided", "Top of Mind"],
            series_names=["Total"],
        )
        tables = _sample_tables()
        candidates = score_fingerprint_against_tables(fp, tables)
        assert len(candidates) == 3
        best = candidates[0]
        assert best.table["title"] == "Brand Awareness"
        assert best.row_score > 0.6
        assert best.score > 0.5

    def test_partial_row_match(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=["Aided", "Unaided", "Something Else"],
            series_names=["Total"],
        )
        tables = _sample_tables()
        candidates = score_fingerprint_against_tables(fp, tables)
        best = candidates[0]
        assert best.table["title"] == "Brand Awareness"
        assert best.row_score > 0

    def test_table_fingerprint_scoring(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Table1", shape_type="table",
            row_labels=["Definitely", "Probably", "Might", "Probably Not"],
            col_headers=["Total", "18-34", "35-54"],
        )
        tables = _sample_tables()
        candidates = score_fingerprint_against_tables(fp, tables)
        best = candidates[0]
        assert best.table["title"] == "Purchase Intent"
        assert best.score > 0.5

    def test_text_fingerprint_returns_empty(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="TextBox", shape_type="text",
            text_content="Question: What brands do you know?",
        )
        candidates = score_fingerprint_against_tables(fp, _sample_tables())
        assert candidates == []

    def test_no_match_scores_low(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=["Alpha", "Beta", "Gamma", "Delta"],
            series_names=["Region A"],
        )
        tables = _sample_tables()
        candidates = score_fingerprint_against_tables(fp, tables)
        assert candidates[0].score < 0.3

    def test_column_overlap_contributes(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=["Promoters", "Passives", "Detractors"],
            series_names=["Total", "Q1", "Q2"],
        )
        tables = _sample_tables()
        candidates = score_fingerprint_against_tables(fp, tables)
        best = candidates[0]
        assert best.table["title"] == "Net Promoter Score"
        assert best.col_score > 0.5


# ---------------------------------------------------------------------------
# Text classification
# ---------------------------------------------------------------------------

class TestClassifyTextShape:
    def test_question_prefix(self):
        assert _classify_text_shape("Question: Which brands?") == "text_question"

    def test_q_prefix(self):
        assert _classify_text_shape("Q: Which brands?") == "text_question"

    def test_base_prefix(self):
        assert _classify_text_shape("Base: Total respondents (n=1500)") == "text_base"

    def test_base_n_prefix(self):
        assert _classify_text_shape("Base N=1500") == "text_base"

    def test_respondent_base(self):
        assert _classify_text_shape("Total respondents, base n=500") == "text_base"

    def test_unclassifiable(self):
        assert _classify_text_shape("Some random slide content") is None


# ---------------------------------------------------------------------------
# Text shape association
# ---------------------------------------------------------------------------

class TestTextShapeAssociation:
    def test_associates_text_on_same_slide(self, pptx_with_chart_and_text):
        prs = Presentation(pptx_with_chart_and_text)
        chart_fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
        )
        chart_result = AutoMapResult(
            fingerprint=chart_fp, table_title="Brand Awareness",
            col_key="Total", confidence=0.9, method="structural",
        )
        text_results = associate_text_shapes(prs, [chart_result])
        assert len(text_results) >= 1
        types = {r.reason for r in text_results}
        assert any("text_question" in t for t in types)

    def test_skips_shapes_with_existing_alt(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        q_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(0.5))
        q_box.text_frame.text = "Question: Already mapped"
        _set_descr(q_box, "type: text_question\ntable_title: Existing Table")

        chart_fp = ShapeFingerprint(slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart")
        chart_result = AutoMapResult(
            fingerprint=chart_fp, table_title="Brand Awareness",
            col_key="Total", confidence=0.9, method="structural",
        )
        text_results = associate_text_shapes(prs, [chart_result])
        for tr in text_results:
            assert tr.fingerprint.text_content != "Question: Already mapped"


# ---------------------------------------------------------------------------
# Alt text writing
# ---------------------------------------------------------------------------

class TestAltTextWriting:
    def test_build_chart_alt_text(self):
        alt = _build_alt_text("Brand Awareness", "Total", "chart")
        assert "table_title: Brand Awareness" in alt
        assert "column: Total" in alt
        assert "type: chart" in alt
        assert "auto_update: yes" in alt

    def test_build_table_alt_text(self):
        alt = _build_alt_text("Purchase Intent", None, "table")
        assert "table_title: Purchase Intent" in alt
        assert "type: table" in alt
        assert "column:" not in alt

    def test_write_and_read_alt(self, pptx_with_chart):
        prs = Presentation(pptx_with_chart)
        slide = prs.slides[0]
        for shp in slide.shapes:
            try:
                _ = shp.chart
                _write_alt_text(shp, "table_title: Test Table\ntype: chart")
                result = _read_existing_alt(shp)
                assert "Test Table" in result
                return
            except (ValueError, AttributeError):
                continue
        pytest.fail("No chart shape found to test alt text writing")


# ---------------------------------------------------------------------------
# End-to-end orchestration (no LLM)
# ---------------------------------------------------------------------------

class TestAutoMapOrchestration:
    def test_auto_map_matches_chart(self, pptx_with_chart):
        prs = Presentation(pptx_with_chart)
        tables = _sample_tables()
        results = auto_map_presentation_obj(
            prs, tables, use_llm=False, write_alt=False,
        )
        data_results = [r for r in results if r.fingerprint.shape_type == "chart"]
        assert len(data_results) == 1
        assert data_results[0].table_title == "Brand Awareness"
        assert data_results[0].method == "structural"
        assert data_results[0].confidence >= 0.5

    def test_auto_map_matches_table(self, pptx_with_table):
        prs = Presentation(pptx_with_table)
        tables = _sample_tables()
        results = auto_map_presentation_obj(
            prs, tables, use_llm=False, write_alt=False,
        )
        data_results = [r for r in results if r.fingerprint.shape_type == "table"]
        assert len(data_results) == 1
        assert data_results[0].table_title == "Purchase Intent"

    def test_auto_map_writes_alt_text(self, pptx_with_chart, tmp_path):
        prs = Presentation(pptx_with_chart)
        out_path = str(tmp_path / "mapped.pptx")
        results = auto_map_presentation_obj(
            prs, _sample_tables(), pptx_out=out_path,
            use_llm=False, write_alt=True,
        )
        assert os.path.exists(out_path)

        prs2 = Presentation(out_path)
        for shp in prs2.slides[0].shapes:
            alt = _read_existing_alt(shp)
            if "Brand Awareness" in alt:
                assert "table_title: Brand Awareness" in alt
                assert "type: chart" in alt
                return
        pytest.fail("Alt text not written to chart shape")

    def test_auto_map_with_text_association(self, pptx_with_chart_and_text):
        prs = Presentation(pptx_with_chart_and_text)
        results = auto_map_presentation_obj(
            prs, _sample_tables(), use_llm=False, write_alt=False,
        )
        text_results = [r for r in results if r.fingerprint.shape_type == "text"]
        assert len(text_results) >= 1
        for tr in text_results:
            assert tr.table_title == "Brand Awareness"
            assert tr.method == "slide_context"

    def test_progress_callback_fires(self, pptx_with_chart):
        prs = Presentation(pptx_with_chart)
        progress_values = []
        auto_map_presentation_obj(
            prs, _sample_tables(), use_llm=False, write_alt=False,
            progress_callback=lambda p: progress_values.append(p),
        )
        assert len(progress_values) >= 1
        assert progress_values[-1] == pytest.approx(1.0)

    def test_multiple_shapes_same_table(self, tmp_path):
        """Two charts with identical categories should both match the same table."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        for _ in range(2):
            cd = CategoryChartData()
            cd.categories = ["Aided", "Unaided", "Top of Mind"]
            cd.add_series("Total", (85, 45, 22))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
            )

        results = auto_map_presentation_obj(
            prs, _sample_tables(), use_llm=False, write_alt=False,
        )
        chart_results = [r for r in results if r.fingerprint.shape_type == "chart"]
        claimed = [r.table_title for r in chart_results if r.table_title]
        assert len(claimed) == 2, "Both charts should be matched"
        assert all(t == "Brand Awareness" for t in claimed)


# ---------------------------------------------------------------------------
# Report serialization
# ---------------------------------------------------------------------------

class TestResultsToReport:
    def test_produces_serializable_dicts(self, pptx_with_chart):
        prs = Presentation(pptx_with_chart)
        results = auto_map_presentation_obj(
            prs, _sample_tables(), use_llm=False, write_alt=False,
        )
        report = results_to_report(results)
        assert isinstance(report, list)
        assert len(report) >= 1
        for entry in report:
            assert "shape_name" in entry
            assert "table_title" in entry
            assert "confidence" in entry
            assert "method" in entry
            assert isinstance(entry["candidates"], list)


# ---------------------------------------------------------------------------
# Module-level helpers
# ---------------------------------------------------------------------------

class TestHelpers:
    def test_norm(self):
        assert _norm("  Hello  World ") == "hello world"
        assert _norm("") == ""
        assert _norm(None) == ""

    def test_jaccard_identical(self):
        assert _jaccard({"a", "b"}, {"a", "b"}) == 1.0

    def test_jaccard_disjoint(self):
        assert _jaccard({"a"}, {"b"}) == 0.0

    def test_jaccard_empty(self):
        assert _jaccard(set(), set()) == 1.0

    def test_jaccard_partial(self):
        result = _jaccard({"a", "b", "c"}, {"b", "c", "d"})
        assert abs(result - 0.5) < 0.01


# ---------------------------------------------------------------------------
# Q-code extraction
# ---------------------------------------------------------------------------

class TestQCodeExtraction:
    def test_extracts_simple_qcode(self):
        assert _extract_qcode("(Q1369) 92. Would you say...") == "(Q1369)"

    def test_extracts_qcode_with_suffix(self):
        assert _extract_qcode("(Q1378_1) 101. How much...") == "(Q1378_1)"

    def test_extracts_qcode_with_guid_prefix(self):
        assert _extract_qcode("006a986d-c834-48fa (Q1383) text") == "(Q1383)"

    def test_no_qcode_returns_none(self):
        assert _extract_qcode("just some regular text") is None

    def test_empty_returns_none(self):
        assert _extract_qcode("") is None
        assert _extract_qcode(None) is None


# ---------------------------------------------------------------------------
# Uninformative series detection
# ---------------------------------------------------------------------------

class TestUninformativeSeries:
    def test_percent_is_uninformative(self):
        assert _is_uninformative_series(["%"]) is True

    def test_series1_is_uninformative(self):
        assert _is_uninformative_series(["Series 1"]) is True

    def test_empty_string_is_uninformative(self):
        assert _is_uninformative_series([""]) is True

    def test_empty_list_is_uninformative(self):
        assert _is_uninformative_series([]) is True

    def test_real_name_is_informative(self):
        assert _is_uninformative_series(["Total"]) is False

    def test_mixed_is_informative(self):
        assert _is_uninformative_series(["Total", "%"]) is False


# ---------------------------------------------------------------------------
# Trend chart detection
# ---------------------------------------------------------------------------

class TestTrendDetection:
    def test_monthly_dates(self):
        assert _is_trend_chart(["Aug 2022", "Sep 2022", "Oct 2022", "Nov 2022"]) is True

    def test_specific_dates(self):
        assert _is_trend_chart([
            "Mar 26-28, 2021", "Jul 7-9, 2021", "Aug 4-6, 2021",
        ]) is True

    def test_non_dates(self):
        assert _is_trend_chart(["Aided", "Unaided", "Top of Mind"]) is False

    def test_too_few_categories(self):
        assert _is_trend_chart(["Jan 2022"]) is False

    def test_mixed_below_threshold(self):
        assert _is_trend_chart(["Jan 2022", "Category A", "Category B", "Category C"]) is False

    def test_trend_skip_in_orchestrator(self):
        """Trend chart should get method='trend_skip' from auto_map_presentation_obj."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        cd = CategoryChartData()
        cd.categories = ["Jan 2022", "Feb 2022", "Mar 2022", "Apr 2022"]
        cd.add_series("Series 1", (10, 20, 30, 40))
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        results = auto_map_presentation_obj(
            prs, _sample_tables(), use_llm=False, write_alt=False,
        )
        chart_results = [r for r in results if r.fingerprint.shape_type == "chart"]
        assert len(chart_results) == 1
        assert chart_results[0].method == "trend_skip"
        assert chart_results[0].table_title is None


# ---------------------------------------------------------------------------
# Fuzzy label matching
# ---------------------------------------------------------------------------

class TestFuzzyLabelMatching:
    def test_exact_match(self):
        assert _fuzzy_match_label("Gas too expensive", ["Gasoline was too expensive", "Other"]) == "Gasoline was too expensive"

    def test_abbreviation_match(self):
        result = _fuzzy_match_label("Gas too expensive", ["Gasoline was too expensive", "Airfare too costly"])
        assert result == "Gasoline was too expensive"

    def test_no_match_below_threshold(self):
        assert _fuzzy_match_label("Completely different", ["Alpha", "Beta"]) is None

    def test_empty_inputs(self):
        assert _fuzzy_match_label("", ["A", "B"]) is None
        assert _fuzzy_match_label("test", []) is None

    def test_containment_exact(self):
        assert _containment({"a", "b"}, {"a", "b", "c"}) == pytest.approx(1.0)

    def test_containment_partial(self):
        assert _containment({"a", "b", "c"}, {"a", "b"}) == pytest.approx(2 / 3)

    def test_containment_empty(self):
        assert _containment(set(), {"a"}) == 0.0

    def test_fuzzy_containment_exact(self):
        result = _fuzzy_containment({"gas too expensive"}, {"gasoline was too expensive"}, threshold=0.5)
        assert result > 0.5

    def test_fuzzy_containment_no_match(self):
        result = _fuzzy_containment({"zzz"}, {"aaa"}, threshold=0.9)
        assert result == 0.0


# ---------------------------------------------------------------------------
# Axis-flip scoring
# ---------------------------------------------------------------------------

class TestAxisFlipScoring:
    def _demo_tables(self):
        """Tables with demographic columns and answer-choice rows."""
        return [
            {
                "id": "Sheet1#1",
                "title": "(Q1264) Barriers to Travel",
                "row_labels": [
                    "Travel is too expensive right now",
                    "Personal financial reasons",
                    "Gasoline was too expensive",
                    "Not enough PTO/vacation time",
                    "Base",
                ],
                "col_labels": ["Total", "Gen Z", "Millennial", "Gen X", "Boomer or older"],
                "values": [
                    [50, 45, 55, 48, 52],
                    [35, 40, 30, 38, 32],
                    [28, 22, 25, 30, 35],
                    [20, 30, 25, 15, 10],
                    [1000, 200, 250, 300, 250],
                ],
            },
        ]

    def test_normal_bar_chart_matches(self):
        """Bar chart with answer-choice categories should match in normal orientation."""
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=[
                "Travel is too expensive right now",
                "Personal financial reasons",
                "Gasoline was too expensive",
                "Not enough PTO/vacation time",
            ],
            series_names=["%"],
        )
        cands = score_fingerprint_against_tables(fp, self._demo_tables())
        best = cands[0]
        assert best.table["title"] == "(Q1264) Barriers to Travel"
        assert best.orientation == "normal"
        assert best.score >= 0.70

    def test_flipped_demo_chart_matches(self):
        """Chart with demographics as categories should match in flipped orientation."""
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=["Total", "Gen Z", "Millennial", "Gen X", "Boomer or older"],
            series_names=["Gas too expensive"],
        )
        cands = score_fingerprint_against_tables(fp, self._demo_tables())
        best = cands[0]
        assert best.table["title"] == "(Q1264) Barriers to Travel"
        assert best.orientation == "flipped"
        assert best.score >= 0.50

    def test_uninformative_series_uses_full_row_weight(self):
        """With series=['%'], score should be based entirely on row similarity."""
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="Chart1", shape_type="chart",
            categories=[
                "Travel is too expensive right now",
                "Personal financial reasons",
                "Gasoline was too expensive",
            ],
            series_names=["%"],
        )
        cands = score_fingerprint_against_tables(fp, self._demo_tables())
        best = cands[0]
        assert best.score >= 0.70
        assert best.col_score == 0.0


# ---------------------------------------------------------------------------
# Q-code matching in orchestrator
# ---------------------------------------------------------------------------

class TestQCodeMatching:
    def test_qcode_match_from_alt_text(self):
        """Shape with (Q1369) in alt text should Q-code-match the right table."""
        tables = [
            {
                "id": "s1",
                "title": "(Q1369) 92. Would you say better off or worse off?",
                "row_labels": ["Much better off", "Better off", "Worse off", "Base"],
                "col_labels": ["Total", "Gen Z"],
                "values": [[30, 35], [40, 38], [20, 15], [1000, 250]],
            },
            {
                "id": "s2",
                "title": "(Q1371) 94. LOOKING FORWARD...",
                "row_labels": ["Much better off", "Better off", "Worse off", "Base"],
                "col_labels": ["Total", "Gen Z"],
                "values": [[25, 30], [35, 32], [25, 20], [1000, 250]],
            },
        ]

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        cd = CategoryChartData()
        cd.categories = ["Much better off", "Better off", "Worse off"]
        cd.add_series("%", (30, 40, 20))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        _set_descr(chart_shape, "guid-here (Q1369) 92. Would you say better off")

        results = auto_map_presentation_obj(prs, tables, use_llm=False, write_alt=False)
        chart_results = [r for r in results if r.fingerprint.shape_type == "chart"]
        assert len(chart_results) == 1
        assert chart_results[0].method == "qcode"
        assert chart_results[0].table_title == "(Q1369) 92. Would you say better off or worse off?"
        assert chart_results[0].confidence == pytest.approx(0.95)

    def test_qcode_not_in_crosstab_falls_through(self):
        """If the Q-code in alt text doesn't match any table, use structural scoring."""
        tables = _sample_tables()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        cd = CategoryChartData()
        cd.categories = ["Aided", "Unaided", "Top of Mind"]
        cd.add_series("Total", (85, 45, 22))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(0.5), Inches(4), Inches(3), cd,
        )
        _set_descr(chart_shape, "(Q9999) Non-existent question")

        results = auto_map_presentation_obj(prs, tables, use_llm=False, write_alt=False)
        chart_results = [r for r in results if r.fingerprint.shape_type == "chart"]
        assert len(chart_results) == 1
        assert chart_results[0].method == "structural"
        assert chart_results[0].table_title == "Brand Awareness"


# ---------------------------------------------------------------------------
# Flipped alt text building
# ---------------------------------------------------------------------------

class TestFlippedAltText:
    def test_build_flipped_alt_text(self):
        alt = _build_alt_text(
            "Barriers to Travel", None, "chart",
            orientation="flipped",
            row_key="Gasoline was too expensive",
            column_keys=["Total", "Gen Z", "Millennial"],
        )
        assert "table_title: Barriers to Travel" in alt
        assert "row_key: Gasoline was too expensive" in alt
        assert "column: Total,Gen Z,Millennial" in alt
        assert "orientation: flipped" in alt
        assert "type: chart" in alt
        assert "auto_update: yes" in alt

    def test_build_normal_alt_text_no_orientation(self):
        alt = _build_alt_text("Brand Awareness", "Total", "chart")
        assert "orientation" not in alt
        assert "row_key" not in alt


# ---------------------------------------------------------------------------
# Infer row/column keys for flipped charts
# ---------------------------------------------------------------------------

class TestInferKeys:
    def _table(self):
        return {
            "title": "Test",
            "row_labels": ["Gasoline was too expensive", "Travel is too expensive", "Base"],
            "col_labels": ["Total", "Gen Z", "Millennial", "Gen X"],
            "values": [[28, 22, 25, 30], [50, 45, 55, 48], [1000, 200, 250, 300]],
        }

    def test_infer_row_key_exact(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="C1", shape_type="chart",
            series_names=["Gasoline was too expensive"],
        )
        assert _infer_row_key(fp, self._table()) == "Gasoline was too expensive"

    def test_infer_row_key_fuzzy(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="C1", shape_type="chart",
            series_names=["Gasoline too expensive"],
        )
        result = _infer_row_key(fp, self._table())
        assert result == "Gasoline was too expensive"

    def test_infer_row_key_uninformative_skipped(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="C1", shape_type="chart",
            series_names=["%"],
        )
        assert _infer_row_key(fp, self._table()) is None

    def test_infer_column_keys(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="C1", shape_type="chart",
            categories=["Total", "Gen Z", "Millennial"],
        )
        result = _infer_column_keys(fp, self._table())
        assert result == ["Total", "Gen Z", "Millennial"]

    def test_infer_column_keys_fuzzy(self):
        fp = ShapeFingerprint(
            slide_idx=0, shape_idx=0, shape_name="C1", shape_type="chart",
            categories=["Total", "Baby Boomer+"],
        )
        table = {
            "title": "T",
            "row_labels": ["A"],
            "col_labels": ["Total", "Boomer or older"],
            "values": [[1]],
        }
        result = _infer_column_keys(fp, table)
        assert result is not None
        assert "Total" in result
