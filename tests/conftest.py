"""Shared fixtures for Report Relay test suite."""

import sys
import os

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Ensure project root is importable
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _set_descr(shape, text: str):
    """Write alt-text *descr* attribute the same way deck_update reads it."""
    el = shape.element
    c_nv_pr = el.find(".//p:cNvPr", namespaces={"p": P_NS})
    if c_nv_pr is not None:
        c_nv_pr.set("descr", text)


# ---------------------------------------------------------------------------
# Sample table dict (matches crosstab_parser output schema)
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_table():
    return {
        "id": "Sheet1#1",
        "sheet": "Sheet1",
        "title": "Brand Awareness",
        "row_labels": ["Aided", "Unaided", "Top of Mind", "Base"],
        "col_labels": ["Total", "Male", "Female"],
        "values": [
            [0.85, 0.80, 0.90],
            [0.45, 0.42, 0.48],
            [0.22, 0.20, 0.24],
            [1500, 700, 800],
        ],
        "meta": {
            "block_start": 0,
            "block_end": 6,
            "col_banners": ["Total", "Male", "Female"],
            "col_groups": ["", "", ""],
        },
    }


# ---------------------------------------------------------------------------
# Minimal tagged PPTX (1 slide, chart + question + base text shapes)
# ---------------------------------------------------------------------------

@pytest.fixture
def mini_pptx(tmp_path):
    """Build a minimal PPTX with a bar chart and two tagged text shapes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # -- bar chart --
    chart_data = CategoryChartData()
    chart_data.categories = ["Cat A", "Cat B", "Cat C"]
    chart_data.add_series("Series 1", (10, 20, 30))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(0.5), Inches(4), Inches(3),
        chart_data,
    )
    _set_descr(
        chart_frame,
        "table_title: Brand Awareness\ncolumn: Total\nauto_update: yes\ntype: chart",
    )

    # -- question text shape --
    q_box = slide.shapes.add_textbox(Inches(5), Inches(0.5), Inches(4), Inches(0.5))
    q_box.text_frame.text = "Question: Old question text"
    _set_descr(
        q_box,
        "type: text_question\ntable_title: Brand Awareness",
    )

    # -- base text shape --
    b_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(0.5))
    b_box.text_frame.text = "Base: Total respondents. 1,000 complete surveys."
    _set_descr(
        b_box,
        "type: text_base\ntable_title: Brand Awareness",
    )

    path = tmp_path / "mini.pptx"
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# Sample crosstab Excel workbook (parseable by crosstab_parser)
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_crosstab_xlsx(tmp_path):
    """Create a small .xlsx that parse_workbook can consume."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    # Row 1: title
    ws.cell(row=1, column=1, value="Brand Awareness")

    # Row 2: empty separator
    # Row 3: banner header
    ws.cell(row=3, column=1, value="")
    ws.cell(row=3, column=2, value="Total")
    ws.cell(row=3, column=3, value="Male")
    ws.cell(row=3, column=4, value="Female")

    # Row 4–7: data rows
    rows = [
        ("Aided",       0.85, 0.80, 0.90),
        ("Unaided",     0.45, 0.42, 0.48),
        ("Top of Mind", 0.22, 0.20, 0.24),
        ("Base",        1500, 700,  800),
    ]
    for i, (label, *vals) in enumerate(rows, start=4):
        ws.cell(row=i, column=1, value=label)
        for j, v in enumerate(vals, start=2):
            ws.cell(row=i, column=j, value=v)

    path = tmp_path / "crosstab.xlsx"
    wb.save(str(path))
    return str(path)
